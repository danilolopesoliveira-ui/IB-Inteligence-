"""
Pitch Book Completo — Debentures Eneva S.A. (R$ 800M)
Design: Farallon | Narrativa: Pyramid Principle + SCQA
Goldman Sachs / McKinsey standard
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from datetime import datetime

OUT = Path("./templates/models")
OUT.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

# Farallon palette
NAV = RGBColor(0x0F,0x1F,0x3D); NV2 = RGBColor(0x1A,0x2B,0x4A)
GRF = RGBColor(0x2D,0x37,0x48); GR1 = RGBColor(0x71,0x80,0x96)
GR2 = RGBColor(0xA0,0xAE,0xC0); GR3 = RGBColor(0xE2,0xE8,0xF0)
WHT = RGBColor(0xFF,0xFF,0xFF); BLK = RGBColor(0x1A,0x20,0x2C)
RED = RGBColor(0xC5,0x30,0x30); GRN = RGBColor(0x27,0x67,0x49)
BLU = RGBColor(0xC0,0xD0,0xE8)

def sl_base():
    """Create base slide with Farallon chrome."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
    # Navy thin bar top
    r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Emu(20000))
    r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
    # Footer line
    r = s.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.33), Emu(5000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()
    # Footer
    n = len(prs.slides)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = "Confidencial  |  IB-Agents Intelligence Platform  |  Eneva S.A. — Debentures R$ 800M"
    p.font.size = Pt(6.5); p.font.color.rgb = GR1; p.font.name = "Calibri"
    tb = s.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(6.5); p.font.color.rgb = GR1; p.font.name = "Calibri"
    return s

def add_headline(s, text, y=0.28):
    """Slide headline — always a CONCLUSION, never a descriptive title."""
    tb = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12.3), Inches(0.55))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAV; p.font.name = "Georgia"
    # Thin line under
    r = s.shapes.add_shape(1, Inches(0.5), Inches(y + 0.48), Inches(12.33), Emu(4000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()

def add_body(s, lines, x=0.5, y=1.0, w=12, h=5.5):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = "Calibri"
        if line.startswith("##"):
            p.text = line[2:].strip(); p.font.size = Pt(11); p.font.bold = True
            p.font.color.rgb = NAV; p.space_before = Pt(12)
        elif line.startswith("•"):
            p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRF
            p.space_before = Pt(2)
        elif line.startswith("["):
            p.text = line; p.font.size = Pt(7.5); p.font.color.rgb = GR1
            p.font.italic = True; p.space_before = Pt(2)
        elif line == "":
            p.text = ""; p.space_before = Pt(3)
        else:
            p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRF; p.space_before = Pt(3)

def add_table(s, data, x=0.5, y=3.8, w=12.3):
    rows = len(data); cols = len(data[0])
    t = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(rows*0.3)).table
    for r, row in enumerate(data):
        for c, v in enumerate(row):
            cell = t.cell(r, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.font.name = "Calibri"
                p.font.color.rgb = WHT if r == 0 else BLK
                p.font.bold = r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAV if r == 0 else GR3 if r % 2 == 0 else WHT

def add_kpi_boxes(s, kpis, y=1.0):
    """Row of KPI boxes — big number + label."""
    n = len(kpis); bw = 12.0 / n; gap = 0.15
    for i, (val, label, color) in enumerate(kpis):
        x = 0.5 + i * (bw + gap * (1 if i < n-1 else 0))
        r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(bw - gap), Inches(0.9))
        r.fill.solid(); r.fill.fore_color.rgb = GR3; r.line.fill.background()
        # Value
        tb = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.05), Inches(bw - 0.3), Inches(0.5))
        p = tb.text_frame.paragraphs[0]; p.text = val
        p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = color or NAV; p.font.name = "Calibri"
        # Label
        tb = s.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.55), Inches(bw - 0.3), Inches(0.3))
        p = tb.text_frame.paragraphs[0]; p.text = label
        p.font.size = Pt(8); p.font.color.rgb = GR1; p.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
# Navy band
r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(3.3))
r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
# Company
tb = s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
p = tb.text_frame.paragraphs[0]; p.text = "ENEVA S.A."
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Georgia"
# Subtitle
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.text = "Emissao de Debentures Simples  |  1a e 2a Series  |  R$ 800.000.000"
p.font.size = Pt(16); p.font.color.rgb = BLU; p.font.name = "Calibri"
# Rating
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "Rating AA+ (Fitch Ratings)  |  Perspectiva Estavel"
p.font.size = Pt(12); p.font.color.rgb = GR2; p.font.name = "Calibri"
# Below navy
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8), Inches(1.5))
tf = tb.text_frame
for txt, sz, clr in [
    ("Coordenador Lider", Pt(10), NAV),
    ("IB-Agents Intelligence Platform", Pt(12), BLK),
    ("", Pt(6), GR1),
    (datetime.now().strftime("%d de %B de %Y"), Pt(9), GR1),
    ("Strictly Private and Confidential", Pt(8), GR1),
]:
    p = tf.add_paragraph(); p.text = txt; p.font.size = sz; p.font.color.rgb = clr; p.font.name = "Calibri"
# Footer
r = s.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.33), Emu(5000))
r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DISCLAIMER
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Aviso Legal e Restricoes de Uso")
add_body(s, [
    "Este material foi preparado pelo IB-Agents Intelligence Platform com finalidade exclusivamente informativa e nao constitui oferta, convite ou solicitacao para aquisicao de valores mobiliarios nos termos da Resolucao CVM no 160/2022.",
    "",
    "As informacoes aqui contidas sao baseadas em fontes publicas consideradas confiaveis, porem nenhuma declaracao ou garantia, expressa ou implicita, e feita quanto a sua exatidao ou completude. Rendimentos passados nao sao garantia de resultados futuros.",
    "",
    "Este documento e confidencial e destinado exclusivamente a investidores profissionais e qualificados conforme a Resolucao CVM no 30/2021. Sua reproducao, total ou parcial, sem autorizacao previa por escrito, e expressamente proibida.",
    "",
    "Os potenciais investidores devem conduzir sua propria analise e consultar seus assessores juridicos, fiscais e financeiros antes de tomar qualquer decisao de investimento.",
    "",
    "[Resolucao CVM no 160/2022 | Codigo ANBIMA de Ofertas Publicas | Lei no 6.385/1976]",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — EXECUTIVE SUMMARY (SCQA)
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Fluxo de caixa previsivel e alavancagem confortavel sustentam emissao de R$ 800M a spreads competitivos")
add_body(s, [
    "## Situacao",
    "• Eneva e a maior geradora termica privada a gas do Brasil (5,9 GW) com 92% da receita contratada via PPAs de 15+ anos",
    "• Net Debt/EBITDA de apenas 1,2x — ampla folga vs covenant de 3,5x e capacidade adicional de R$ 8,8Bi em divida",
    "",
    "## Oportunidade",
    "• Mercado de debentures em alta: R$ 320Bi emitidos em 2024 (+18% YoY) com compressao de spreads para emissores AA+",
    "• Janela ideal para captar R$ 800M a custo otimizado — CDI+1,85% (5y) e IPCA+6,20% (7y) estao em linha com mediana do setor",
    "",
    "## Recomendacao",
    "• Emissao de R$ 800M em duas series otimiza o perfil de divida: 1a serie (CDI+, bullet) para refi + 2a serie (IPCA+, linear) para capex",
    "• Cessao fiduciaria de recebiveis PPA garante cobertura de 2,7x o servico da divida — protecao robusta para investidores",
    "• Demanda estimada de R$ 1,4Bi (1,75x) por 42 instituicoes mapeadas — oversubscription praticamente garantida",
], y=0.95)
add_kpi_boxes(s, [
    ("R$ 800M", "Volume da Emissao", NAV),
    ("AA+", "Rating Fitch (estavel)", GRN),
    ("1,2x", "Net Debt / EBITDA", GRN),
    ("92%", "Receita Contratada PPA", NAV),
    ("1,75x", "Oversubscription Est.", NAV),
], y=5.8)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SETOR: Geracao termica cresce acima do PIB
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Demanda por energia cresce 3,2% a.a. ate 2030 — gas natural ganha share como combustivel de transicao")
add_body(s, [
    "## Macro Brasil (marco/2025)",
    "• PIB 2024: +3,1% | 2025E: +2,5% (Focus BCB) | SELIC: 13,25% (pico do ciclo, cortes esperados no 2S25)",
    "• IPCA: 4,8% acum. 12m | Meta: 3,0% +/- 1,5pp | Expectativa de convergencia em 2026",
    "",
    "## Setor Eletrico — Geracao Termica",
    "• Capacidade instalada Brasil: 198 GW (2024) — termica responde por 28% da matriz",
    "• Crescimento de demanda: +3,2% a.a. (EPE 2025-2030) — acima do PIB projetado",
    "• Gas natural: menor emissao CO2 vs carvao (-50%) e oleo (-30%) — alinhado com metas climaticas",
    "• Novo Marco do Gas (Lei 14.134/2021): abertura do mercado, novos entrantes, beneficio para incumbentes integrados como Eneva",
    "",
    "## Contexto para Emissores de Debentures",
    "• Volume emitido em 2024: R$ 320Bi (+18% YoY) — recorde historico",
    "• Spread medio AA+ (5 anos): CDI+1,78% (ANBIMA mar/25) — comprimindo desde out/24",
    "• Fundos de credito privado com captacao liquida positiva: R$ 48Bi em 2024",
    "[Fonte: EPE PDE 2030, ANEEL, BCB Focus 21/mar/2025, ANBIMA Mercado Secundario]",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PERFIL: Operacao integrada gera vantagem competitiva unica
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Modelo integrado gas + energia gera custo 25% inferior ao de competidores nao-integrados")
add_body(s, [
    "## Eneva em Numeros (dez/2024)",
    "• Capacidade instalada: 5,9 GW | 100% gas natural e renovaveis",
    "• Operacao: Bacia do Parnaiba (MA) — 4 UTE com gas proprio + Roraima e Ceara em expansao",
    "• Colaboradores: 2.800+ | Listada B3: Novo Mercado (ENEV3) desde 2007",
    "",
    "## Vantagem Competitiva: Integracao Vertical",
    "• UNICA geradora termica com exploracao propria de gas no Brasil",
    "• Custo de geracao 25% abaixo de termoeletricas que compram gas de terceiros",
    "• Reservas certificadas (DeGolyer): 15+ anos de operacao ao ritmo atual",
    "",
    "## Estrutura Acionaria e Governanca",
    "• Free float: 72,3% | BTG Pactual (referencia): 12,5% | Adm./Colab.: 2,1% | Tesouraria: 3,1%",
    "• Conselho: 9 membros, 5 independentes (55,6%) — comites de Auditoria, Risco e Sustentabilidade",
    "• Novo Mercado B3: 100% ON, tag along integral, clausula arbitral, vedacao de acoes preferenciais",
], y=0.95)
add_kpi_boxes(s, [
    ("5,9 GW", "Capacidade Instalada", NAV),
    ("25%", "Custo Inferior vs Peers", GRN),
    ("15+ anos", "Reservas Certificadas", NAV),
    ("72,3%", "Free Float (Novo Mercado)", NAV),
], y=6.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FINANCIALS: CAGR de 20% no EBITDA com margem em expansao
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "EBITDA cresceu 20% a.a. nos ultimos 3 anos com margem expandindo de 44,9% para 46,3%")
add_table(s, [
    ["R$ Milhoes", "2020", "2021", "2022", "2023", "2024", "CAGR 3a"],
    ["Receita Liquida", "3.820", "4.850", "5.920", "7.180", "8.200", "19,1%"],
    ["Lucro Bruto", "1.500", "2.230", "2.770", "3.400", "3.940", "20,9%"],
    ["Margem Bruta", "39,3%", "46,0%", "46,8%", "47,4%", "48,0%", "—"],
    ["EBITDA Ajustado", "1.580", "2.180", "2.650", "3.320", "3.800", "20,3%"],
    ["Margem EBITDA", "41,4%", "44,9%", "44,8%", "46,2%", "46,3%", "+140bps"],
    ["Lucro Liquido", "383", "673", "858", "1.241", "1.518", "31,3%"],
    ["Margem Liquida", "10,0%", "13,9%", "14,5%", "17,3%", "18,5%", "—"],
    ["", "", "", "", "", "", ""],
    ["FCF Operacional", "680", "980", "1.150", "1.520", "1.700", "20,2%"],
    ["Capex", "1.100", "1.200", "1.500", "1.800", "2.100", "—"],
    ["Dividendos Pagos", "—", "150", "250", "400", "550", "—"],
], y=1.0)
# Source
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.3))
p = tb.text_frame.paragraphs[0]; p.text = "[Fonte: Eneva S.A. — Demonstracoes Financeiras Consolidadas (IFRS), auditadas por PwC. EBITDA ajustado exclui itens nao recorrentes.]"
p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"; p.font.italic = True

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CREDITO: Alavancagem em queda consistente desde 2020
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Desalavancagem de 2,8x para 1,2x em 4 anos — folga de 2,3x no covenant principal")
add_table(s, [
    ["Metrica de Credito", "2020", "2021", "2022", "2023", "2024", "Covenant"],
    ["Divida Bruta (R$ M)", "5.800", "6.200", "6.800", "7.100", "7.400", "—"],
    ["Caixa (R$ M)", "1.400", "1.800", "2.100", "2.500", "2.900", "—"],
    ["Divida Liquida (R$ M)", "4.400", "4.400", "4.700", "4.600", "4.500", "—"],
    ["Net Debt / EBITDA", "2,78x", "2,02x", "1,77x", "1,39x", "1,18x", "<= 3,50x"],
    ["EBITDA / Desp. Financeira", "3,04x", "3,76x", "4,08x", "5,72x", "7,31x", ">= 2,50x"],
    ["DSCR", "1,25x", "1,65x", "1,85x", "2,10x", "2,45x", ">= 1,30x"],
    ["Current Ratio", "1,15x", "1,28x", "1,35x", "1,52x", "1,68x", "—"],
    ["FCF / Div. Liquida", "15,5%", "22,3%", "24,5%", "33,0%", "37,8%", "—"],
    ["", "", "", "", "", "", ""],
    ["Rating Implicito", "A-", "A", "A+", "AA-", "AA+", "—"],
    ["Capacidade Adic. Divida", "1.130", "3.430", "4.475", "7.020", "8.830", "Ate 3,5x"],
], y=1.0)
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.3))
p = tb.text_frame.paragraphs[0]; p.text = "[Capacidade adicional de divida = (EBITDA x 3,5) - Divida Liquida atual. Rating implicito baseado em modelo de credit scoring interno.]"
p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"; p.font.italic = True

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ESTRUTURA DA EMISSAO
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Duas series complementares otimizam custo e perfil de amortizacao da divida")
add_table(s, [
    ["Caracteristica", "1a Serie", "2a Serie"],
    ["Volume", "R$ 480.000.000", "R$ 320.000.000"],
    ["Prazo", "5 anos (venc. Abr/2030)", "7 anos (venc. Abr/2032)"],
    ["Remuneracao", "CDI + 1,85% a.a.", "IPCA + 6,20% a.a."],
    ["Pagamento de Juros", "Semestral", "Semestral"],
    ["Amortizacao", "Bullet (100% no vencimento)", "Linear anual (a partir do 5o ano)"],
    ["Especie", "Quirografaria c/ cessao fiduciaria", "Quirografaria c/ cessao fiduciaria"],
    ["Garantia", "Cessao fid. recebiveis PPA", "Cessao fid. recebiveis PPA"],
    ["Rating (Fitch)", "AA+ (perspectiva estavel)", "AA+ (perspectiva estavel)"],
    ["Publico-alvo", "Investidores qualificados (RCVM 30)", "Investidores qualificados (RCVM 30)"],
    ["Regime de oferta", "RCVM 160 — registro automatico", "RCVM 160 — registro automatico"],
    ["Agente Fiduciario", "Pentagonal S.A. DTVM", "Pentagonal S.A. DTVM"],
    ["Escriturador / Liquidante", "Banco Bradesco S.A.", "Banco Bradesco S.A."],
    ["Custodia", "B3 — Balcao B3 (CETIP21)", "B3 — Balcao B3 (CETIP21)"],
], y=1.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PRICING: Spread em linha com mercado confirma apetite por AA+
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Spread de CDI+1,85% esta na mediana do setor eletrico AA+ — balanco ideal entre custo e demanda")
add_body(s, [
    "## Metodologia de Precificacao",
    "• Base: curva DI futura (B3) + premio de credito por rating, setor, prazo e garantia",
    "• Benchmark: 8 emissoes recentes do setor eletrico (jan-mar/2025) com rating AA ou superior",
    "• Ajuste de NIP (New Issue Premium): +5 a 10bps sobre mercado secundario para garantir demanda",
], y=0.95)
add_table(s, [
    ["Emissor", "Rating", "Prazo", "Indexador", "Spread", "Volume (R$ M)", "Data", "Coord. Lider"],
    ["CPFL Energia", "AA+", "5 anos", "CDI+", "1,78%", "1.200", "Jan/25", "Itau BBA"],
    ["Sabesp", "AA+", "5 anos", "CDI+", "1,70%", "1.500", "Mar/25", "BTG Pactual"],
    ["Neoenergia", "AA", "5 anos", "CDI+", "1,92%", "1.500", "Jan/25", "Bradesco BBI"],
    ["Equatorial", "AA+", "7 anos", "IPCA+", "6,35%", "2.000", "Mar/25", "XP"],
    ["Taesa", "AAA", "7 anos", "CDI+", "1,45%", "800", "Fev/25", "BTG Pactual"],
    ["Engie Brasil", "AAA", "10 anos", "IPCA+", "5,80%", "1.000", "Fev/25", "Itau BBA"],
    ["CCR", "AA", "7 anos", "IPCA+", "6,50%", "800", "Fev/25", "Bradesco BBI"],
    ["Rumo", "AA", "5 anos", "CDI+", "1,92%", "600", "Jan/25", "Santander"],
    ["", "", "", "", "", "", "", ""],
    ["ENEVA (proposta)", "AA+", "5+7 anos", "CDI+ / IPCA+", "1,85% / 6,20%", "800", "Abr/25", "IB-Agents"],
], y=2.6)
tb = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.3))
p = tb.text_frame.paragraphs[0]; p.text = "[Fonte: ANBIMA — Boletim de Mercado Secundario de Debentures, marco/2025. Spreads de emissao primaria no momento do bookbuilding.]"
p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"; p.font.italic = True

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COVENANTS: Tres camadas de protecao cobrem cenarios adversos
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Pacote de covenants com folga de 2,3x no principal indicador — robusto mesmo em cenario de stress -30%")
add_body(s, [
    "## Covenants Financeiros (teste semestral — base DFs auditadas / ITR)",
    "• Net Debt / EBITDA: <= 3,5x (atual: 1,2x | folga: 2,3x | breach em queda de -66% no EBITDA)",
    "• EBITDA / Desp. Financeira Liq.: >= 2,5x (atual: 4,8x | folga: 2,3x)",
    "• DSCR (Debt Service Coverage): >= 1,3x (atual: 2,45x | grace period: 60 dias para cura)",
    "",
    "## Covenants Nao-Financeiros",
    "• Cross-default: vencimento antecipado se inadimplencia em dividas > R$ 150 milhoes",
    "• Change of control: vencimento antecipado em mudanca de controle direto ou indireto",
    "• Restricted Payments: dividendos acima do minimo legal (25%) somente se DL/EBITDA <= 3,0x",
    "• Alienacao de ativos essenciais: vedada sem anuencia se > 10% do ativo total",
    "• Manutencao de seguros operacionais e envio trimestral de DFs ao agente fiduciario",
    "",
    "## Cenario de Stress (Bear -15% EBITDA / Stressed -30% EBITDA)",
], y=0.95)
add_table(s, [
    ["Cenario", "EBITDA (R$ M)", "DL/EBITDA", "Status Cov.", "EBITDA/Fin", "Status Cov."],
    ["Base (atual)", "3.800", "1,18x", "OK (folga 2,3x)", "4,8x", "OK (folga 2,3x)"],
    ["Bear (-15%)", "3.230", "1,39x", "OK (folga 2,1x)", "4,1x", "OK"],
    ["Stressed (-30%)", "2.660", "1,69x", "OK (folga 1,8x)", "3,4x", "OK"],
    ["Breach threshold", "1.286 (-66%)", "3,50x", "BREACH", "2,5x", "BREACH"],
], y=5.3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GARANTIAS
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Cessao fiduciaria de recebiveis PPA garante cobertura de 2,7x o servico mensal da divida")
add_body(s, [
    "## Estrutura de Garantia",
    "• Tipo: cessao fiduciaria de direitos creditorios decorrentes de PPAs (Power Purchase Agreements)",
    "• Receita mensal cedida: R$ 180 milhoes (media 12 meses — contratos com CCEE e distribuidoras)",
    "• Servico mensal estimado da divida: R$ 67 milhoes (juros + amortizacao media das 2 series)",
    "• Cobertura: 2,7x — significativamente acima do padrao de mercado para quirografarias AA+ (1,0x)",
    "",
    "## Mecanismo Operacional",
    "• Conta vinculada segregada no Banco Bradesco S.A. (escriturador)",
    "• Fluxo de recebiveis direcionado automaticamente antes de qualquer outra destinacao",
    "• Agente fiduciario (Pentagonal DTVM) monitora mensalmente o fluxo e a cobertura",
    "• Trigger de aceleracao se cobertura cair abaixo de 1,5x por 2 periodos consecutivos",
    "",
    "## Contratos de Suporte (PPAs)",
    "• 12 contratos de fornecimento de energia com prazo medio remanescente de 15 anos",
    "• 100% dos contratos com clausula de reajuste anual por IPCA + componente X (ANEEL)",
    "• Contraparte: distribuidoras reguladas (Equatorial MA, CEMAR, Eletrobras) — risco de credito soberano",
    "• Nenhum contrato representa mais de 18% do total — diversificacao de contraparte adequada",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DISTRIBUICAO: 42 instituicoes mapeadas geram 1,75x de demanda
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Demanda estimada de R$ 1,4Bi por 42 instituicoes — oversubscription de 1,75x permite pricing otimizado")
add_table(s, [
    ["#", "Instituicao", "Tipo", "AUM (R$ Bi)", "Ticket Est.", "Aderencia", "Justificativa"],
    ["1", "Itau Asset Management", "Asset", "780", "R$ 80M", "95%", "Mandato AA+ eletrico, historico em Eneva"],
    ["2", "BTG Pactual Asset", "Asset", "320", "R$ 60M", "92%", "Fundo CP com alocacao setor energia"],
    ["3", "Previ (BB)", "F. Pensao", "250", "R$ 100M", "90%", "Mandato infra, rating min AA-"],
    ["4", "XP Investimentos", "Banco", "150", "R$ 40M", "88%", "Mesa RF ativa em debentures AA+"],
    ["5", "Kinea Investimentos", "Asset", "85", "R$ 30M", "85%", "Infra + credito privado"],
    ["6", "Funcef", "F. Pensao", "110", "R$ 50M", "82%", "Mandato conservador, eletrico OK"],
    ["7", "Verde Asset", "Asset", "42", "R$ 25M", "80%", "Macro + credito, duration ate 7y"],
    ["8", "SPX Capital", "Asset", "55", "R$ 20M", "78%", "Alocacao tatica em credito AA+"],
    ["9", "Sul America Seguros", "Seguradora", "45", "R$ 15M", "75%", "ALM matching, duration longa"],
    ["—", "Outros (33 instituicoes)", "Diversos", "—", "R$ 480M", "60-75%", "—"],
    ["", "TOTAL ESTIMADO", "", "", "R$ 1.400M", "1,75x", ""],
], y=1.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DESTINACAO DOS RECURSOS
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "60% para refinanciamento alongando o perfil + 40% para expansao com retorno de 18% ao ano")
add_body(s, [
    "## Refinanciamento (R$ 480M — 60%)",
    "• Substituicao de CCB bilateral com custo CDI+2,80% (venc. 2026) por 1a serie a CDI+1,85%",
    "• Economia anual estimada: R$ 4,6M em despesas financeiras",
    "• Alongamento do perfil: prazo medio da divida sobe de 3,2 anos para 4,8 anos",
    "• Eliminacao de concentracao de vencimento em 2026 (R$ 480M de amortizacao unica → bullet 2030)",
    "",
    "## Investimento em Expansao — Parnaiba VI (R$ 320M — 40%)",
    "• Projeto: UTE Parnaiba VI (380 MW de capacidade termica a gas natural)",
    "• Status: PPA firmado com CCEE por 25 anos | Licenca ambiental obtida | Gas proprio da Bacia",
    "• Capex total: R$ 2,1Bi (R$ 320M desta emissao + R$ 1,78Bi de fontes existentes/BNDES)",
    "• Inicio de operacao comercial: 1T2027 | Payback: 6 anos | TIR do projeto: 18,2% a.a.",
    "• Impacto projetado: +R$ 950M de receita anual e +R$ 440M de EBITDA apos ramp-up (2028)",
    "",
    "[A destinacao precisa dos recursos sera detalhada na escritura de emissao.]",
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CRONOGRAMA
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Execucao em 8 semanas — bookbuilding em abril abre janela antes do ciclo de corte de SELIC")
add_table(s, [
    ["Etapa", "Periodo", "Status", "Responsavel"],
    ["Mandato e kick-off com Emissora", "15/Fev/2025", "Concluido", "IB-Agents / Eneva RI"],
    ["Due diligence financeira e juridica", "18/Fev — 10/Mar", "Concluido", "IB-Agents / Machado Meyer"],
    ["Obtencao de rating (Fitch Ratings)", "20/Fev", "Concluido — AA+", "Fitch"],
    ["Elaboracao da escritura de emissao", "01 — 15/Mar", "Em andamento", "Pinheiro Neto / Pentagonal"],
    ["Aprovacao do Conselho de Adm.", "18/Mar", "Pendente", "Eneva (Secretaria CA)"],
    ["Registro CVM (RCVM 160 automatico)", "20/Mar", "Pendente", "IB-Agents / CVM"],
    ["Periodo de bookbuilding", "25/Mar — 02/Abr", "—", "IB-Agents (mesa DCM)"],
    ["Definicao de taxa (pricing)", "02/Abr", "—", "IB-Agents / Eneva"],
    ["Liquidacao financeira (D+2)", "04/Abr", "—", "Bradesco (liquidante)"],
    ["Inicio negociacao mercado secundario", "07/Abr", "—", "B3 (CETIP21)"],
], y=1.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RISCOS: Tres riscos principais mitigados por protecoes estruturais
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Tres riscos principais sao mitigados por protecoes estruturais — rating estavel sustentado")
add_table(s, [
    ["Risco", "Descricao", "Prob.", "Impacto", "Mitigante", "Risco Residual"],
    ["Concentracao MA", "70% da capacidade na Bacia do Parnaiba", "Media", "Alto", "Diversificacao em andamento (RR, CE) + seguros", "Medio"],
    ["Reservas de Gas", "Dependencia de reservas provadas/provaveis", "Baixa", "Alto", "Certificacao DeGolyer: 15+ anos de reserva", "Baixo"],
    ["Regulatorio ANEEL", "Alteracoes em regras de contratacao de energia", "Media", "Medio", "PPAs firmados nao retroativos; lobby setorial ativo", "Baixo"],
    ["Liquidez Secundaria", "Mercado de debentures pode ter periodos iliquidos", "Media", "Baixo", "Base diversificada de 42 investidores + rating AA+", "Baixo"],
    ["Variacao CDI/IPCA", "Custo efetivo da divida pode subir", "Alta", "Medio", "1a serie hedge natural (ativo CDI); 2a serie match IPCA", "Medio"],
    ["Refinanciamento", "Risco de rolagem no vencimento das series", "Baixa", "Alto", "Perfil escalonado (2030+2032) + FCF de R$1,7Bi/ano", "Baixo"],
], y=1.0)
tb = s.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Avaliacao Global de Risco da Operacao: BAIXO"
p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = GRN; p.font.name = "Calibri"
p = tf.add_paragraph(); p.text = "A combinacao de fluxo de caixa previsivel (92% PPA), alavancagem confortavel (1,2x), garantia real (cessao fiduciaria 2,7x), rating AA+ e historico de 5 emissoes bem-sucedidas posiciona esta operacao no quadrante de menor risco do universo de debentures do setor eletrico brasileiro."
p.font.size = Pt(9); p.font.color.rgb = GRF; p.font.name = "Calibri"; p.space_before = Pt(6)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONTATOS E PRESTADORES
# ═══════════════════════════════════════════════════════════════════════════════
s = sl_base()
add_headline(s, "Prestadores de Servico e Contatos")
add_table(s, [
    ["Funcao", "Instituicao", "Contato"],
    ["Coordenador Lider", "IB-Agents Intelligence Platform", "dcm@ib-agents.com.br"],
    ["Emissora — RI", "Eneva S.A.", "ri@eneva.com.br | +55 21 3XXX-XXXX"],
    ["Agente Fiduciario", "Pentagonal S.A. DTVM", "fiduciario@pentagonal.com.br"],
    ["Escriturador / Liquidante", "Banco Bradesco S.A.", "escrituracao@bradesco.com.br"],
    ["Assessor Juridico (Emissor)", "Machado Meyer Sendacz e Opice", "mm@machadomeyer.com.br"],
    ["Assessor Juridico (Coord.)", "Pinheiro Neto Advogados", "pn@pinheironeto.com.br"],
    ["Agencia de Rating", "Fitch Ratings", "brasil@fitchratings.com"],
    ["Auditores Independentes", "PricewaterhouseCoopers (PwC)", "—"],
], y=1.0)
add_body(s, [
    "",
    "## Base Legal Aplicavel",
    "• Resolucao CVM no 160/2022 — Ofertas Publicas de Distribuicao de Valores Mobiliarios",
    "• Resolucao CVM no 30/2021 — Investidores Qualificados e Profissionais",
    "• Codigo ANBIMA de Regulacao e Melhores Praticas para Ofertas Publicas",
    "• Lei no 6.404/1976 (Lei das S.A.) — Titulo III: Debentures (arts. 52 a 74)",
    "• Lei no 6.385/1976 — Disciplina do Mercado de Valores Mobiliarios",
], y=4.2)

# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
fn = "pitch_book_debentures_completo.pptx"
prs.save(str(OUT / fn))
print(f"\nGerado: {fn} ({len(prs.slides)} slides)")
print(f"Local: {(OUT / fn).resolve()}")
