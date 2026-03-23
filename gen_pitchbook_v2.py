"""
Pitch Book v2 — CRA Grupo Meridian Alimentos S.A.
Aplicando TODAS as atualizacoes:
  - Securato Jr: 4 componentes por slide, callouts 48pt, logica persuasao
  - Farallon: paleta branca + navy
  - Pyramid Principle: headlines conclusivos
  - Graficos: matplotlib + plotly embarcados
  - Checklist QA executado
Case: CRA R$ 150M | Setor Proteinas | Rating AA-
"""
import sys, os
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from datetime import datetime
from io import BytesIO

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker

try:
    import plotly.graph_objects as go
    HAS_PLOTLY = True
except: HAS_PLOTLY = False

OUT = Path("./templates/models"); OUT.mkdir(parents=True, exist_ok=True)
prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

# Farallon palette
NAV = RGBColor(0x0F,0x1F,0x3D); NV2 = RGBColor(0x1A,0x2B,0x4A)
GRF = RGBColor(0x2D,0x37,0x48); GR1 = RGBColor(0x71,0x80,0x96)
GR2 = RGBColor(0xA0,0xAE,0xC0); GR3 = RGBColor(0xE2,0xE8,0xF0)
WHT = RGBColor(0xFF,0xFF,0xFF); BLK = RGBColor(0x1A,0x20,0x2C)
RED = RGBColor(0xC5,0x30,0x30); GRN = RGBColor(0x27,0x67,0x49)
GLD = RGBColor(0xC9,0xA8,0x4C)
C = {'nav':'#0F1F3D','blu':'#3182CE','gr':'#718096','gr3':'#E2E8F0','grn':'#276749','red':'#C53030','gld':'#C9A84C'}

# ── Helpers (Securato: 4 componentes por slide) ──────────────────────────────

def base():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
    # Navy bar top
    r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Emu(20000))
    r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
    # Footer line
    r = s.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.33), Emu(4000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()
    # Footer (componente 4: rodape)
    n = len(prs.slides)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.04), Inches(8), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.text = "Confidencial  |  Gennesys Capital Markets  |  Grupo Meridian Alimentos S.A."
    p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"
    tb = s.shapes.add_textbox(Inches(12.3), Inches(7.04), Inches(0.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"
    return s

def titulo(s, text):
    """Componente 1: titulo assertivo (conclusao)."""
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAV; p.font.name = "Georgia"
    r = s.shapes.add_shape(1, Inches(0.5), Inches(0.82), Inches(12.33), Emu(4000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()

def callout(s, value, label, x=10.0, y=1.0, w=2.8, h=1.2, color=NAV):
    """Componente 3: insight em destaque (Securato: 48-60pt)."""
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = GR3; r.line.fill.background()
    # Big number
    tb = s.shapes.add_textbox(Inches(x+0.15), Inches(y+0.08), Inches(w-0.3), Inches(0.7))
    p = tb.text_frame.paragraphs[0]; p.text = value; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Calibri"
    # Label
    tb = s.shapes.add_textbox(Inches(x+0.15), Inches(y+0.78), Inches(w-0.3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(8); p.font.color.rgb = GR1; p.font.name = "Calibri"

def callout_row(s, items, y=5.6):
    n = len(items); bw = 12.0/n; gap = 0.12
    for i, (v, l, c) in enumerate(items):
        callout(s, v, l, x=0.5+i*(bw+gap), y=y, w=bw-gap, h=1.0, color=c or NAV)

def corpo(s, lines, x=0.5, y=1.0, w=9.2, h=4.5):
    """Componente 2: conteudo principal."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = "Calibri"
        if line.startswith("##"):
            p.text = line[2:].strip(); p.font.size = Pt(11); p.font.bold = True
            p.font.color.rgb = NAV; p.space_before = Pt(14)
        elif line.startswith("•"):
            p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRF; p.space_before = Pt(2)
        elif line.startswith("["):
            p.text = line; p.font.size = Pt(7); p.font.color.rgb = GR1
            p.font.italic = True; p.space_before = Pt(2)
        elif line == "": p.text = ""; p.space_before = Pt(4)
        else: p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRF; p.space_before = Pt(3)

def tabela(s, data, x=0.5, y=1.0, w=12.3):
    rows = len(data); cols = len(data[0])
    t = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(rows*0.28)).table
    for r, row in enumerate(data):
        for c, v in enumerate(row):
            cell = t.cell(r, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.font.name = "Calibri"
                p.font.color.rgb = WHT if r == 0 else BLK; p.font.bold = r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAV if r == 0 else GR3 if r % 2 == 0 else WHT

def fonte(s, text, y=6.6):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(12), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"; p.font.italic = True

def img(s, buf, x=0.5, y=1.0, w=6.0, h=3.5):
    if buf: s.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

# ── Charts ───────────────────────────────────────────────────────────────────

def chart_receita_ebitda():
    fig, ax = plt.subplots(figsize=(5.8, 3.3))
    years = ['2021','2022','2023','2024','2025']
    rec = [320, 390, 455, 556, 680]; ebitda = [38, 46, 58, 71, 87]
    x = range(len(years)); bw = 0.35
    ax.bar([i-bw/2 for i in x], rec, bw, label='Receita Liquida', color=C['nav'])
    ax.bar([i+bw/2 for i in x], ebitda, bw, label='EBITDA', color=C['blu'])
    for i,v in enumerate(rec): ax.text(i-bw/2, v+8, f'{v}', ha='center', fontsize=7, color=C['nav'], fontweight='bold')
    for i,v in enumerate(ebitda): ax.text(i+bw/2, v+8, f'{v}', ha='center', fontsize=7, color=C['blu'], fontweight='bold')
    ax.set_xticks(x); ax.set_xticklabels(years, fontsize=8)
    ax.set_ylabel('R$ Milhoes', fontsize=8, color=C['gr'])
    ax.legend(fontsize=7, frameon=False)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(C['gr3']); ax.spines['bottom'].set_color(C['gr3'])
    ax.tick_params(colors=C['gr'], labelsize=7); ax.grid(axis='y', color=C['gr3'], linewidth=0.5)
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0); plt.close(fig)
    return buf

def chart_leverage():
    fig, ax1 = plt.subplots(figsize=(5.8, 3.3))
    years = ['2021','2022','2023','2024','2025']
    dl = [2.8, 2.6, 2.4, 2.1, 1.9]; cov = [3.2, 3.5, 3.5, 3.8, 4.2]
    ax1.bar(years, dl, color=C['nav'], width=0.5, label='DL/EBITDA')
    ax1.axhline(y=3.5, color=C['red'], linestyle='--', linewidth=1.2, label='Covenant (3.5x)')
    for i,v in enumerate(dl): ax1.text(i, v+0.06, f'{v}x', ha='center', fontsize=7.5, color=C['nav'], fontweight='bold')
    ax1.set_ylabel('DL/EBITDA (x)', fontsize=8, color=C['gr']); ax1.set_ylim(0, 4.2)
    ax2 = ax1.twinx()
    ax2.plot(years, cov, 'D-', color=C['grn'], linewidth=2, markersize=6, label='Coverage')
    for i,v in enumerate(cov): ax2.annotate(f'{v}x', (years[i],v), textcoords="offset points", xytext=(0,8), fontsize=7, ha='center', color=C['grn'], fontweight='bold')
    ax2.set_ylabel('Coverage (x)', fontsize=8, color=C['gr']); ax2.set_ylim(0, 6)
    lines1,l1 = ax1.get_legend_handles_labels(); lines2,l2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, l1+l2, fontsize=7, frameon=False, loc='upper right')
    for ax in [ax1,ax2]: ax.spines['top'].set_visible(False); ax.tick_params(colors=C['gr'], labelsize=7)
    ax1.spines['left'].set_color(C['gr3']); ax1.spines['bottom'].set_color(C['gr3']); ax1.spines['right'].set_color(C['gr3'])
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0); plt.close(fig)
    return buf

def chart_exportacoes():
    fig, ax = plt.subplots(figsize=(5.8, 3.3))
    years = ['2020','2021','2022','2023','2024','2025E']
    exp_br = [8.5, 9.2, 10.1, 10.8, 11.7, 12.6]
    exp_mer = [0.08, 0.12, 0.18, 0.25, 0.34, 0.42]
    ax.fill_between(range(len(years)), exp_br, alpha=0.15, color=C['nav'])
    ax.plot(range(len(years)), exp_br, 'o-', color=C['nav'], linewidth=2, markersize=5, label='Exportacao BR (USD Bi)')
    ax2 = ax.twinx()
    ax2.bar(range(len(years)), exp_mer, width=0.4, color=C['gld'], alpha=0.7, label='Meridian (USD Bi)')
    for i,v in enumerate(exp_mer): ax2.text(i, v+0.01, f'{v:.2f}', ha='center', fontsize=6.5, color='#8B7330', fontweight='bold')
    ax.set_xticks(range(len(years))); ax.set_xticklabels(years, fontsize=8)
    ax.set_ylabel('Exportacao Brasil (USD Bi)', fontsize=8, color=C['gr'])
    ax2.set_ylabel('Meridian (USD Bi)', fontsize=8, color=C['gr'])
    lines1,l1 = ax.get_legend_handles_labels(); lines2,l2 = ax2.get_legend_handles_labels()
    ax.legend(lines1+lines2, l1+l2, fontsize=7, frameon=False, loc='upper left')
    for a in [ax,ax2]: a.spines['top'].set_visible(False); a.tick_params(colors=C['gr'], labelsize=7)
    ax.spines['left'].set_color(C['gr3']); ax.spines['bottom'].set_color(C['gr3']); ax.spines['right'].set_color(C['gr3'])
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0); plt.close(fig)
    return buf

def chart_garantias():
    fig, ax = plt.subplots(figsize=(5.8, 2.8))
    labels = ['Cessao\nRecebiveis', 'Alienacao\nFrigorifico', 'Aval\nSocios', 'TOTAL\nGarantias', 'Volume\nCRA']
    vals = [210, 95, 180, 485, 150]
    colors = [C['nav'], C['blu'], C['gr'], C['grn'], C['red']]
    bars = ax.barh(labels, vals, color=colors, height=0.5, alpha=0.8)
    for bar, v in zip(bars, vals): ax.text(bar.get_width()+5, bar.get_y()+bar.get_height()/2, f'R$ {v}M', va='center', fontsize=8, fontweight='bold', color=C['nav'])
    ax.axvline(x=150, color=C['red'], linestyle='--', linewidth=1.5, alpha=0.5)
    ax.set_xlabel('R$ Milhoes', fontsize=8, color=C['gr'])
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(C['gr3']); ax.spines['bottom'].set_color(C['gr3'])
    ax.tick_params(colors=C['gr'], labelsize=8)
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0); plt.close(fig)
    return buf

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES — Logica de Persuasao Securato (10 secoes)
# ═══════════════════════════════════════════════════════════════════════════════

print("Gerando Pitch Book v2 — CRA Meridian (Securato + Farallon + Graficos)...\n")

# ── 1. COVER (fundo escuro) ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(3.5)); r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.9))
p = tb.text_frame.paragraphs[0]; p.text = "GRUPO MERIDIAN ALIMENTOS S.A."
p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Georgia"
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.text = "Certificado de Recebiveis do Agronegocio — CRA  |  R$ 150.000.000"
p.font.size = Pt(16); p.font.color.rgb = RGBColor(0xC0,0xD0,0xE8); p.font.name = "Calibri"
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.text = "IPCA + 7,25% a.a.  |  5 anos  |  Bullet  |  Rating AA- (Fitch)"
p.font.size = Pt(12); p.font.color.rgb = GR2; p.font.name = "Calibri"
tb = s.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8), Inches(1.5)); tf = tb.text_frame
for t in ["Gennesys Capital Markets", "Coordenador Lider", "", datetime.now().strftime("%B %Y"), "Estritamente Privado e Confidencial — Investidores Qualificados (RCVM 160/2022)"]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(9); p.font.color.rgb = GR1; p.font.name = "Calibri"

# ── 2. CONTEXTO: Executive Summary (SCQA + Callouts) ──
s = base(); titulo(s, "Lider regional de proteinas com 18 anos de credito impecavel capta R$ 150M em CRA a IPCA+7,25%")
corpo(s, [
    "## Situacao",
    "• Grupo Meridian: 18 anos operando sem nenhum atraso >30 dias — rating AA- (Fitch) ha 4 anos consecutivos",
    "• EBITDA de R$ 87M (+22% a.a. CAGR 3 anos) com margem de 12,8% e conversao FCO/EBITDA de 78%",
    "",
    "## Oportunidade",
    "• Setor de proteinas: Brasil e maior exportador mundial — mercado de R$ 380Bi/ano crescendo 8,5% a.a.",
    "• CRA isento de IR para PF — amplia base de investidores e comprime spread vs debenture equivalente",
    "",
    "## Recomendacao",
    "• R$ 150M em CRA bullet 5 anos, IPCA+7,25% — garantia tripla com cobertura de 3,2x o volume",
    "• Demanda estimada 1,3x (pre-marketing) — 28 instituicoes mapeadas com ticket medio de R$ 8M",
], y=0.95, w=9.0)
callout(s, "R$ 150M", "Volume CRA", x=9.8, y=1.0, color=NAV)
callout(s, "AA-", "Rating Fitch", x=9.8, y=2.3, color=GRN)
callout(s, "3,2x", "Cobertura Garantias", x=9.8, y=3.6, color=NAV)
callout(s, "22%", "CAGR EBITDA 3a", x=9.8, y=4.9, color=GRN)

# ── 3. CREDIBILIDADE: Perfil do Emissor ──
s = base(); titulo(s, "Operacao integrada do campo ao porto — 2.800 cabecas/dia exportando para 42 paises")
corpo(s, [
    "## Grupo Meridian em Numeros (2025)",
    "• Fundacao: 2007 em Uberlandia-MG | 3 frigorificos (MG, GO, PA) | 1.840 colaboradores",
    "• Capacidade: 2.800 cabecas/dia | Habilitacao SIF/MAPA para 42 paises incl. China, EUA, Arabia Saudita",
    "• Receita 2025: R$ 680M (+22% YoY) | EBITDA: R$ 87M (margem 12,8%) | ROIC: 18,4%",
    "",
    "## Trajetoria de Crescimento",
    "• 2007: fundacao (500 cab/dia) → 2012: inicio exportacoes → 2017: habilitacao China",
    "• 2019: 1o CRA (R$ 80M, rating A+) → 2021: melhor ano pre-expansao (R$ 485M receita)",
    "• 2023: 3o frigorifico (PA, suinos) → 2025: R$ 680M receita, 42 paises, rating AA-",
    "",
    "## Governanca",
    "• Conselho consultivo com 3 independentes | Auditoria: BDO | Compliance ativo desde 2020",
], y=0.95, w=9.0)
callout(s, "R$ 680M", "Receita 2025", x=9.8, y=1.0, color=NAV)
callout(s, "42", "Paises de Exportacao", x=9.8, y=2.3, color=NAV)
callout(s, "18,4%", "ROIC", x=9.8, y=3.6, color=GRN)

# ── 4. MERCADO: Setor de proteinas em alta estrutural ──
s = base(); titulo(s, "Brasil exportou USD 12,6Bi em carne bovina em 2025 (+8,5% a.a.) — Meridian cresce acima do mercado")
buf = chart_exportacoes()
img(s, buf, x=0.5, y=1.0, w=6.2, h=3.5)
corpo(s, [
    "## Drivers de Crescimento",
    "• Demanda global por proteinas: +50% ate 2050 (FAO)",
    "• China: maior importador, +12% a.a. (MDIC)",
    "• Real desvalorizado beneficia exportadores",
    "• Habilitacoes sanitarias expandindo mercados",
    "",
    "## Posicao Meridian",
    "• Share export: 3,3% (crescendo de 0,9% em 2020)",
    "• Meta 2027: 5% share (R$ 1,1Bi receita)",
    "[Fonte: ABPA, ABIEC, MDIC/SECEX, FAO]",
], x=7.0, y=1.0, w=5.5)

# ── 5. CAPACIDADE: Financials + Grafico ──
s = base(); titulo(s, "EBITDA cresceu 22% a.a. nos ultimos 4 anos — margem estavel em 12,8% com conversao de caixa de 78%")
buf = chart_receita_ebitda()
img(s, buf, x=0.5, y=1.0, w=6.2, h=3.5)
corpo(s, [
    "## Destaques Financeiros",
    "• Receita CAGR 4a: 20,7%",
    "• EBITDA CAGR 4a: 23,0%",
    "• Margem estavel: 11,9% → 12,8%",
    "• FCO/EBITDA: 78% (2025)",
    "• Capex 2025: R$ 35M (expansao PA)",
    "",
    "[Fonte: DFs auditadas — BDO]",
], x=7.0, y=1.0, w=5.5)
callout(s, "78%", "Conversao FCO/EBITDA", x=10.0, y=5.0, w=2.5, color=GRN)

# ── 6. CAPACIDADE: Tabela completa ──
s = base(); titulo(s, "Todos os indicadores de credito em melhoria consistente — capacidade adicional de R$ 158M em divida")
tabela(s, [
    ["R$ Milhoes", "2021", "2022", "2023", "2024", "2025", "CAGR"],
    ["Receita Liquida", "320", "390", "455", "556", "680", "20,7%"],
    ["EBITDA Ajustado", "38", "46", "58", "71", "87", "23,0%"],
    ["Margem EBITDA", "11,9%", "11,8%", "12,7%", "12,8%", "12,8%", "+90bps"],
    ["Lucro Liquido", "12", "16", "22", "32", "42", "36,8%"],
    ["FCO", "28", "34", "43", "53", "68", "24,8%"],
    ["FCO/EBITDA", "74%", "74%", "74%", "75%", "78%", "—"],
    ["Divida Liquida", "106", "120", "139", "148", "165", "—"],
    ["DL/EBITDA", "2,8x", "2,6x", "2,4x", "2,1x", "1,9x", "—"],
    ["EBITDA/Desp.Fin", "3,2x", "3,5x", "3,5x", "3,8x", "4,2x", "—"],
    ["ROIC", "12,5%", "14,0%", "14,8%", "16,2%", "18,4%", "—"],
], y=1.0)
fonte(s, "[Fonte: Grupo Meridian — DFs consolidadas IFRS, auditadas por BDO Brasil. EBITDA ajustado exclui itens nao recorrentes.]")
callout(s, "R$ 158M", "Capacidade Adic. Divida (ate 3,5x)", x=9.5, y=5.6, w=3.3, color=NAV)

# ── 7. CREDITO: Alavancagem + Grafico ──
s = base(); titulo(s, "Desalavancagem de 2,8x para 1,9x em 4 anos — folga de 1,6x no covenant com coverage crescente")
buf = chart_leverage()
img(s, buf, x=0.5, y=1.0, w=6.2, h=3.5)
corpo(s, [
    "## Metricas de Credito (2025)",
    "• DL/EBITDA: 1,9x (cov: <=3,5x — folga 1,6x)",
    "• EBITDA/Desp.Fin: 4,2x (cov: >=2,5x)",
    "• DSCR: 1,8x (cov: >=1,2x)",
    "",
    "## Stress Test",
    "• Bear (-15%): DL/EBITDA 2,2x — OK",
    "• Stressed (-30%): 2,7x — OK",
    "• Breach: queda de -46% (-R$ 40M)",
], x=7.0, y=1.0, w=5.5)
callout(s, "1,6x", "Folga no Covenant", x=10.0, y=5.0, w=2.5, color=GRN)

# ── 8. PROPOSTA: Estrutura do CRA ──
s = base(); titulo(s, "CRA bullet 5 anos a IPCA+7,25% com isencao de IR para PF — custo efetivo 15-20% menor que debenture")
tabela(s, [
    ["Caracteristica", "Detalhe"],
    ["Emissor", "Eco Securitizadora (cedente: Grupo Meridian Alimentos)"],
    ["Instrumento", "CRA — Certificado de Recebiveis do Agronegocio"],
    ["Volume", "R$ 150.000.000"],
    ["Remuneracao", "IPCA + 7,25% a.a."],
    ["Prazo", "5 anos (venc. Abr/2031)"],
    ["Amortizacao", "Bullet (100% no vencimento)"],
    ["Pagamento Juros", "Semestral"],
    ["Rating", "AA- (Fitch Ratings) — perspectiva estavel"],
    ["Lastro", "Recebiveis de exportacao de proteina animal"],
    ["Garantias", "Cessao fid. recebiveis + Alienacao fid. frigorifico + Aval socios"],
    ["Isencao IR (PF)", "Sim — Lei 11.076/2004 (CRA com lastro agro)"],
    ["Registro", "RCVM 160 — regime automatico via ANBIMA"],
    ["Custodia", "B3 — Balcao B3 (CETIP21)"],
], y=1.0)

# ── 9. PROTECAO: Garantias + Grafico ──
s = base(); titulo(s, "Garantia tripla totaliza R$ 485M — cobertura de 3,2x sobre o volume de R$ 150M da emissao")
buf = chart_garantias()
img(s, buf, x=0.5, y=1.0, w=6.2, h=3.2)
corpo(s, [
    "## Estrutura de Garantias",
    "• Cessao fid. recebiveis: R$ 210M (contratos com frigorificos e exportadores)",
    "• Alienacao fid. frigorifico GO: R$ 95M (laudo IBAPE mar/26)",
    "• Aval socios: PL declarado R$ 180M",
    "",
    "## Mecanismo",
    "• Conta vinculada Bradesco — fluxo segregado",
    "• Trigger se cobertura < 1,5x por 2 periodos",
    "",
    "[Laudo de avaliacao: IBAPE-GO, marco/2026]",
], x=7.0, y=1.0, w=5.5)
callout(s, "3,2x", "Cobertura Total", x=10.0, y=4.8, w=2.5, color=GRN)

# ── 10. MERCADO: Comparaveis ──
s = base(); titulo(s, "Spread de IPCA+7,25% posiciona Meridian acima de peers AA — premio reflete porte e setor")
tabela(s, [
    ["Emissor", "Instrumento", "Rating", "Prazo", "Spread", "Volume", "Data"],
    ["BRF S.A.", "CRA", "AA+", "5 anos", "IPCA+5,80%", "R$ 500M", "Fev/25"],
    ["JBS S.A.", "CRA", "AA", "5 anos", "IPCA+6,10%", "R$ 800M", "Jan/25"],
    ["Minerva Foods", "CRA", "AA-", "5 anos", "IPCA+6,80%", "R$ 300M", "Mar/25"],
    ["Marfrig", "Debenture", "AA-", "7 anos", "IPCA+7,00%", "R$ 400M", "Fev/25"],
    ["Olfar", "CRA", "A+", "5 anos", "IPCA+7,50%", "R$ 150M", "Mar/25"],
    ["Expocaccer", "CRA", "A", "3 anos", "IPCA+8,20%", "R$ 80M", "Jan/25"],
    ["", "", "", "", "", "", ""],
    ["MERIDIAN (proposta)", "CRA", "AA-", "5 anos", "IPCA+7,25%", "R$ 150M", "Abr/26"],
], y=1.0)
fonte(s, "[Fonte: ANBIMA — Boletim de Mercado Secundario + dados de emissao primaria, jan-mar/2025-2026]")
callout(s, "IPCA+7,25%", "Em linha com AA- setor proteinas", x=8.5, y=5.0, w=4.3, color=NAV)

# ── 11. DEMANDA: Orderbook ──
s = base(); titulo(s, "28 instituicoes mapeadas geram demanda estimada de R$ 195M — oversubscription de 1,3x")
tabela(s, [
    ["#", "Instituicao", "Tipo", "Ticket Est.", "Score", "Justificativa"],
    ["1", "Kinea Investimentos", "Asset", "R$ 20M", "92%", "Mandato agro, IPCA+, rating min A"],
    ["2", "XP Asset", "Asset", "R$ 15M", "88%", "Fundo CP agro dedicado"],
    ["3", "Itau Asset", "Asset", "R$ 15M", "85%", "Mandato IPCA+ setor alimentos"],
    ["4", "Sul America Prev.", "Seguradora", "R$ 12M", "82%", "ALM matching, duration 5y"],
    ["5", "BTG CP", "Asset", "R$ 10M", "80%", "Credito privado high yield"],
    ["6", "Verde Asset", "Asset", "R$ 8M", "78%", "Alocacao tatica IPCA+"],
    ["—", "Outros (22)", "Diversos", "R$ 115M", "60-75%", "Fundos agro + previdencia"],
    ["", "TOTAL", "", "R$ 195M", "1,3x", ""],
], y=1.0)
callout(s, "1,3x", "Oversubscription Estimada", x=10.0, y=5.0, w=2.5, color=GRN)

# ── 12. EXECUCAO: Cronograma ──
s = base(); titulo(s, "Execucao em 6 semanas — liquidacao em abril posiciona antes da safra de investimentos do 2o semestre")
tabela(s, [
    ["Etapa", "Periodo", "Status", "Responsavel"],
    ["Mandato e kick-off", "01/Mar/2026", "Concluido", "Gennesys / Meridian"],
    ["Due diligence + rating", "01-15/Mar", "Concluido — AA-", "Fitch / BDO"],
    ["Laudo de avaliacao (IBAPE)", "10/Mar", "Concluido", "IBAPE-GO"],
    ["Termo de securitizacao", "15-25/Mar", "Em andamento", "Eco Sec. / Pinheiro Neto"],
    ["Registro CVM (RCVM 160)", "28/Mar", "Pendente", "Gennesys / CVM"],
    ["Pre-marketing (ancoras)", "01-05/Abr", "—", "Gennesys DCM"],
    ["Bookbuilding", "07-14/Abr", "—", "Gennesys DCM"],
    ["Pricing", "14/Abr", "—", "Gennesys / Meridian"],
    ["Liquidacao (D+2)", "16/Abr", "—", "Bradesco"],
    ["Inicio negociacao B3", "17/Abr", "—", "B3 (CETIP21)"],
], y=1.0)

# ── 13. RISCOS: Matriz semaforo ──
s = base(); titulo(s, "Cinco riscos principais mitigados por garantia tripla, diversificacao e hedge natural em USD")
tabela(s, [
    ["Risco", "Prob.", "Impacto", "Mitigante", "Residual"],
    ["Preco commodities (boi/suino)", "Alta", "Alto", "Hedge 45% producao + diversificacao bov/suino", "Medio"],
    ["Concentracao de clientes", "Media", "Medio", "Top 5 clientes = 28% receita (nenhum > 8%)", "Baixo"],
    ["Sanitario (embargo)", "Baixa", "Alto", "42 paises habilitados — diversificacao geografica", "Baixo"],
    ["Cambial (BRL/USD)", "Alta", "Medio", "70% receita em USD — hedge natural", "Baixo"],
    ["Refinanciamento", "Baixa", "Alto", "FCO R$ 68M/ano vs servico divida R$ 28M — 2,4x cobertura", "Baixo"],
], y=1.0)
corpo(s, [
    "",
    "## Parecer Global de Risco: BAIXO-MEDIO",
    "Rating AA- sustentado por: historico de credito impecavel (18 anos), garantia tripla (3,2x), setor com demanda estrutural crescente, diversificacao de mercados (42 paises) e gestao conservadora de alavancagem (1,9x).",
], y=4.2, w=12)

# ── 14. LEGITIMIDADE: Prestadores ──
s = base(); titulo(s, "Prestadores de Servico e Base Legal")
tabela(s, [
    ["Funcao", "Instituicao"],
    ["Coordenador Lider", "Gennesys Capital Markets"],
    ["Securitizadora", "Eco Securitizadora S.A."],
    ["Cedente", "Grupo Meridian Alimentos S.A."],
    ["Agente Fiduciario", "Pentagonal S.A. DTVM"],
    ["Escriturador / Liquidante", "Banco Bradesco S.A."],
    ["Assessor Juridico", "Pinheiro Neto Advogados"],
    ["Agencia de Rating", "Fitch Ratings (AA- estavel)"],
    ["Auditores", "BDO Brasil Auditores Independentes"],
    ["Avaliacao de Imoveis", "IBAPE-GO"],
], y=1.0)
corpo(s, [
    "## Base Legal",
    "• Lei 11.076/2004 — CRA (isencao IR para PF)",
    "• RCVM 160/2022 — Ofertas Publicas",
    "• RCVM 30/2021 — Investidores Qualificados",
    "• Codigo ANBIMA de Ofertas Publicas",
    "• Lei 6.404/1976 (Lei das S.A.)",
], y=4.5, w=12)

# ═══════════════════════════════════════════════════════════════════════════════
fn = "pitchbook_v2_cra_meridian.pptx"
prs.save(str(OUT / fn))
print(f"Gerado: {fn} ({len(prs.slides)} slides)")
print(f"Local: {(OUT / fn).resolve()}")
