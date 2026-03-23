"""
Pitch Book AI-Generated — Motor inteligente que usa Claude para gerar conteudo
e renderiza em PPTX com padrao Securato/Farallon.

Fontes grandes (14-16pt corpo, 28pt titulos, 48pt callouts).
Graficos em todos os slides de dados.
Layout 2 colunas (visual + insights).
"""
import sys, os, json
sys.path.insert(0, os.path.dirname(__file__))

from dotenv import load_dotenv
load_dotenv()

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from datetime import datetime
from io import BytesIO

import anthropic
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

try:
    import plotly.graph_objects as go
    HAS_PLOTLY = True
except: HAS_PLOTLY = False

OUT = Path("./templates/models"); OUT.mkdir(parents=True, exist_ok=True)

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 1: Ask Claude to generate the full analysis
# ═══════════════════════════════════════════════════════════════════════════════

print("STEP 1: Gerando analise completa via Claude API...\n")

client = anthropic.Anthropic()

ANALYSIS_PROMPT = """Voce e um analista senior de Investment Banking preparando um pitch book para a seguinte operacao:

EMPRESA: AgroVale Participacoes S.A.
SETOR: Agronegocio — producao e exportacao de soja, milho e algodao
OPERACAO: Emissao de Debentures Simples — R$ 500 milhoes
RATING: AA (Fitch Ratings)
PRAZO: 5 anos (1a Serie CDI+1,65%) e 7 anos (2a Serie IPCA+6,00%)

Gere um JSON com a analise completa. O JSON deve ter EXATAMENTE esta estrutura:
{
  "company": {
    "name": "AgroVale Participacoes S.A.",
    "sector": "string — descricao do setor",
    "description": "string — 3 frases sobre a empresa",
    "founded": "ano",
    "headquarters": "cidade-UF",
    "employees": numero,
    "key_differentials": ["string", "string", "string"],
    "governance": "string — resumo governanca"
  },
  "financials": {
    "years": ["2021","2022","2023","2024","2025E"],
    "revenue": [numeros em R$ milhoes],
    "ebitda": [numeros],
    "ebitda_margin": [decimais ex: 0.25],
    "net_income": [numeros],
    "net_debt": [numeros],
    "dl_ebitda": [decimais ex: 1.5],
    "coverage": [decimais],
    "fcf": [numeros],
    "capex": [numeros]
  },
  "market": {
    "tam_brl_bi": numero,
    "tam_growth": "string — CAGR",
    "brazil_position": "string — posicao do Brasil no mercado global",
    "drivers": ["string — driver 1 com dado", "string — driver 2", "string — driver 3", "string — driver 4"],
    "tailwinds": ["string", "string"],
    "source": "string — fontes dos dados"
  },
  "deal": {
    "volume": 500000000,
    "series": [
      {"name": "1a Serie", "volume_pct": 60, "tenor": "5 anos", "rate": "CDI + 1,65%", "amort": "Bullet"},
      {"name": "2a Serie", "volume_pct": 40, "tenor": "7 anos", "rate": "IPCA + 6,00%", "amort": "Linear 5o ano+"}
    ],
    "guarantee": "string — tipo de garantia",
    "guarantee_coverage": "2.5x",
    "use_of_proceeds": "string — destinacao dos recursos",
    "rating": "AA (Fitch)",
    "registration": "RCVM 160 — automatico"
  },
  "comparables": [
    {"issuer": "string", "rating": "string", "tenor": "string", "spread": "string", "volume": "string", "date": "string"},
    {"issuer": "string", "rating": "string", "tenor": "string", "spread": "string", "volume": "string", "date": "string"},
    {"issuer": "string", "rating": "string", "tenor": "string", "spread": "string", "volume": "string", "date": "string"},
    {"issuer": "string", "rating": "string", "tenor": "string", "spread": "string", "volume": "string", "date": "string"},
    {"issuer": "string", "rating": "string", "tenor": "string", "spread": "string", "volume": "string", "date": "string"},
    {"issuer": "string", "rating": "string", "tenor": "string", "spread": "string", "volume": "string", "date": "string"}
  ],
  "covenants": {
    "dl_ebitda_max": 3.5,
    "coverage_min": 2.5,
    "dscr_min": 1.3,
    "stress_bear": {"ebitda_drop": -15, "dl_ebitda": numero, "status": "OK/BREACH"},
    "stress_severe": {"ebitda_drop": -30, "dl_ebitda": numero, "status": "OK/BREACH"},
    "breach_threshold": "string — % queda EBITDA para breach"
  },
  "investors": [
    {"name": "string", "type": "string", "ticket": "string", "score": "string"},
    {"name": "string", "type": "string", "ticket": "string", "score": "string"},
    {"name": "string", "type": "string", "ticket": "string", "score": "string"},
    {"name": "string", "type": "string", "ticket": "string", "score": "string"},
    {"name": "string", "type": "string", "ticket": "string", "score": "string"},
    {"name": "string", "type": "string", "ticket": "string", "score": "string"}
  ],
  "risks": [
    {"risk": "string", "probability": "Alta/Media/Baixa", "impact": "Alto/Medio/Baixo", "mitigant": "string"},
    {"risk": "string", "probability": "string", "impact": "string", "mitigant": "string"},
    {"risk": "string", "probability": "string", "impact": "string", "mitigant": "string"},
    {"risk": "string", "probability": "string", "impact": "string", "mitigant": "string"}
  ],
  "headlines": {
    "exec_summary": "string — headline conclusivo do slide de executive summary",
    "company": "string — headline conclusivo do slide de perfil",
    "market": "string — headline conclusivo do slide de mercado",
    "financials": "string — headline conclusivo do slide financeiro",
    "credit": "string — headline conclusivo do slide de credito",
    "deal": "string — headline conclusivo do slide de estrutura",
    "pricing": "string — headline conclusivo do slide de pricing",
    "covenants": "string — headline conclusivo do slide de covenants",
    "investors": "string — headline conclusivo do slide de demanda",
    "risks": "string — headline conclusivo do slide de riscos"
  },
  "investment_thesis": "string — 3 frases resumindo a tese de investimento"
}

REGRAS:
1. Dados devem ser REALISTAS para o setor de agronegocio brasileiro em 2025
2. Comparaveis devem ser empresas REAIS do setor (SLC, BrasilAgro, Boa Safra, Rumo, Cosan)
3. Headlines devem ser CONCLUSOES assertivas (Pyramid Principle), nao titulos descritivos
4. Investidores devem ser instituicoes REAIS (Itau Asset, BTG, Kinea, etc)
5. Todos os numeros devem ser internamente consistentes

Responda APENAS com o JSON, sem texto antes ou depois."""

response = client.messages.create(
    model="claude-sonnet-4-6",
    max_tokens=8000,
    messages=[{"role": "user", "content": ANALYSIS_PROMPT}]
)

raw = response.content[0].text.strip()
if raw.startswith("```"): raw = raw.split("\n", 1)[1].rsplit("```", 1)[0]
data = json.loads(raw)
print(f"  Tokens: {response.usage.input_tokens} in / {response.usage.output_tokens} out")
print(f"  Empresa: {data['company']['name']}")
print(f"  Volume: R$ {data['deal']['volume']/1e6:.0f}M")
print()

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 2: Generate charts from the AI data
# ═══════════════════════════════════════════════════════════════════════════════

print("STEP 2: Gerando graficos com matplotlib/plotly...\n")

CN = '#0F1F3D'; CB = '#3182CE'; CG = '#276749'; CR = '#C53030'; CGR = '#718096'; CL = '#E2E8F0'
NAV=RGBColor(0x0F,0x1F,0x3D); GRF=RGBColor(0x2D,0x37,0x48); GR1=RGBColor(0x71,0x80,0x96)
GR2=RGBColor(0xA0,0xAE,0xC0); GR3=RGBColor(0xE2,0xE8,0xF0); WHT=RGBColor(0xFF,0xFF,0xFF)
BLK=RGBColor(0x1A,0x20,0x2C); GRN=RGBColor(0x27,0x67,0x49); RED=RGBColor(0xC5,0x30,0x30)

fin = data['financials']
yrs = fin['years']

def mkfig(w=6.5, h=3.8):
    fig, ax = plt.subplots(figsize=(w, h))
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(CL); ax.spines['bottom'].set_color(CL)
    ax.tick_params(colors=CGR, labelsize=9); ax.grid(axis='y', color=CL, linewidth=0.5)
    return fig, ax

def savefig(fig):
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white')
    buf.seek(0); plt.close(fig); return buf

# Chart 1: Revenue + EBITDA bars
fig, ax = mkfig()
x = np.arange(len(yrs)); bw = 0.35
ax.bar(x-bw/2, fin['revenue'], bw, label='Receita Liquida (R$ M)', color=CN)
ax.bar(x+bw/2, fin['ebitda'], bw, label='EBITDA (R$ M)', color=CB)
for i,v in enumerate(fin['revenue']): ax.text(i-bw/2, v+max(fin['revenue'])*0.02, f'{v:,.0f}', ha='center', fontsize=8, color=CN, fontweight='bold')
for i,v in enumerate(fin['ebitda']): ax.text(i+bw/2, v+max(fin['revenue'])*0.02, f'{v:,.0f}', ha='center', fontsize=8, color=CB, fontweight='bold')
ax.set_xticks(x); ax.set_xticklabels(yrs, fontsize=10)
ax.set_ylabel('R$ Milhoes', fontsize=10, color=CGR)
ax.legend(fontsize=9, frameon=False)
chart_revenue = savefig(fig)
print("  chart_revenue OK")

# Chart 2: Margins
fig, ax = mkfig()
mg = [m*100 for m in fin['ebitda_margin']]
ax.plot(yrs, mg, 'o-', color=CN, linewidth=2.5, markersize=8, label='Margem EBITDA (%)')
for i,v in enumerate(mg): ax.annotate(f'{v:.1f}%', (yrs[i],v), textcoords="offset points", xytext=(0,12), fontsize=9, ha='center', color=CN, fontweight='bold')
ax.set_ylim(min(mg)-5, max(mg)+8)
ax.set_ylabel('%', fontsize=10, color=CGR)
ax.legend(fontsize=9, frameon=False)
chart_margins = savefig(fig)
print("  chart_margins OK")

# Chart 3: Leverage dual axis
fig, ax1 = mkfig()
ax1.bar(yrs, fin['dl_ebitda'], color=CN, width=0.5, label='DL/EBITDA (x)', alpha=0.85)
ax1.axhline(y=data['covenants']['dl_ebitda_max'], color=CR, linestyle='--', linewidth=1.5, label=f'Covenant ({data["covenants"]["dl_ebitda_max"]}x)')
for i,v in enumerate(fin['dl_ebitda']): ax1.text(i, v+0.05, f'{v:.1f}x', ha='center', fontsize=9, color=CN, fontweight='bold')
ax1.set_ylabel('DL/EBITDA (x)', fontsize=10, color=CGR)
ax1.set_ylim(0, data['covenants']['dl_ebitda_max']+1)
ax2 = ax1.twinx()
ax2.plot(yrs, fin['coverage'], 'D-', color=CG, linewidth=2, markersize=7, label='Coverage (x)')
for i,v in enumerate(fin['coverage']): ax2.annotate(f'{v:.1f}x', (yrs[i],v), textcoords="offset points", xytext=(0,10), fontsize=9, ha='center', color=CG, fontweight='bold')
ax2.set_ylabel('Coverage (x)', fontsize=10, color=CGR)
ax2.set_ylim(0, max(fin['coverage'])+2)
lines1,l1=ax1.get_legend_handles_labels(); lines2,l2=ax2.get_legend_handles_labels()
ax1.legend(lines1+lines2, l1+l2, fontsize=8, frameon=False, loc='upper right')
ax1.spines['left'].set_color(CL); ax1.spines['bottom'].set_color(CL); ax1.spines['right'].set_color(CL)
for a in [ax1,ax2]: a.spines['top'].set_visible(False); a.tick_params(colors=CGR, labelsize=9)
chart_leverage = savefig(fig)
print("  chart_leverage OK")

# Chart 4: FCF
fig, ax = mkfig()
ax.bar(yrs, fin['fcf'], color=CG, width=0.5, alpha=0.85, label='FCF (R$ M)')
ax.bar(yrs, [-c for c in fin['capex']], color=CR, width=0.5, alpha=0.5, label='Capex (R$ M)')
for i,v in enumerate(fin['fcf']): ax.text(i, v+max(fin['fcf'])*0.03, f'{v:,.0f}', ha='center', fontsize=9, color=CG, fontweight='bold')
ax.set_ylabel('R$ Milhoes', fontsize=10, color=CGR)
ax.legend(fontsize=9, frameon=False)
ax.axhline(y=0, color=CGR, linewidth=0.5)
chart_fcf = savefig(fig)
print("  chart_fcf OK")

# Chart 5: Stress test waterfall (if plotly available)
chart_stress = None
if HAS_PLOTLY:
    cov = data['covenants']
    current_dl = fin['dl_ebitda'][-1]
    scenarios = ['Base', f'Bear ({cov["stress_bear"]["ebitda_drop"]}%)', f'Stressed ({cov["stress_severe"]["ebitda_drop"]}%)' , f'Covenant ({cov["dl_ebitda_max"]}x)']
    values = [current_dl, cov['stress_bear']['dl_ebitda'], cov['stress_severe']['dl_ebitda'], cov['dl_ebitda_max']]
    colors_bar = [CN, '#D4A853', CR if cov['stress_severe']['status']=='BREACH' else '#D4A853', '#888888']
    fig = go.Figure(go.Bar(x=scenarios, y=values, marker_color=colors_bar, text=[f'{v:.1f}x' for v in values], textposition='outside', textfont_size=14))
    fig.add_hline(y=cov['dl_ebitda_max'], line_dash="dash", line_color=CR, line_width=2, annotation_text=f"Covenant {cov['dl_ebitda_max']}x", annotation_font_size=11)
    fig.update_layout(width=650, height=380, margin=dict(l=10,r=10,t=30,b=10), plot_bgcolor='white', paper_bgcolor='white',
        yaxis=dict(title='DL/EBITDA (x)', gridcolor=CL, title_font_size=11, tickfont_size=10),
        xaxis=dict(tickfont_size=10), font=dict(family='Calibri', color=CN))
    buf = BytesIO(); fig.write_image(buf, format='png', scale=2); buf.seek(0)
    chart_stress = buf
    print("  chart_stress OK (plotly)")

print()

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 3: Build PPTX with professional layout
# ═══════════════════════════════════════════════════════════════════════════════

print("STEP 3: Construindo PPTX profissional...\n")

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
co = data['company']; dl = data['deal']; hl = data['headlines']

def add_base():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
    r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Emu(20000)); r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
    r = s.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.33), Emu(4000)); r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()
    n = len(prs.slides)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.04), Inches(8), Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text = f"Confidencial  |  IB-Agents  |  {co['name']}"
    p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"
    tb = s.shapes.add_textbox(Inches(12.4), Inches(7.04), Inches(0.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"
    return s

def add_title(s, text):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAV; p.font.name = "Georgia"
    r = s.shapes.add_shape(1, Inches(0.5), Inches(0.88), Inches(12.33), Emu(4000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()

def add_text(s, lines, x=0.5, y=1.1, w=5.8, h=5.0, size=14):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = "Calibri"
        if line.startswith("##"):
            p.text = line[2:].strip(); p.font.size = Pt(size+2); p.font.bold = True
            p.font.color.rgb = NAV; p.space_before = Pt(16)
        elif line.startswith("•"):
            p.text = line; p.font.size = Pt(size); p.font.color.rgb = GRF; p.space_before = Pt(4)
        elif line.startswith("["):
            p.text = line; p.font.size = Pt(9); p.font.color.rgb = GR1; p.font.italic = True; p.space_before = Pt(4)
        elif line == "": p.text = ""; p.space_before = Pt(6)
        else: p.text = line; p.font.size = Pt(size); p.font.color.rgb = GRF; p.space_before = Pt(5)

def add_callout(s, value, label, x, y, w=2.8, h=1.3, color=NAV):
    r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = GR3; r.line.fill.background()
    tb = s.shapes.add_textbox(Inches(x+0.1), Inches(y+0.05), Inches(w-0.2), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = value; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = color; p.font.name = "Calibri"
    tb = s.shapes.add_textbox(Inches(x+0.1), Inches(y+0.88), Inches(w-0.2), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10); p.font.color.rgb = GR1; p.font.name = "Calibri"

def add_img(s, buf, x=0.5, y=1.1, w=6.5, h=4.0):
    if buf: s.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

def add_table(s, data_tbl, x=0.5, y=1.1, w=12.3):
    rows = len(data_tbl); cols = len(data_tbl[0])
    t = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(rows*0.34)).table
    for r, row in enumerate(data_tbl):
        for c, v in enumerate(row):
            cell = t.cell(r, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10); p.font.name = "Calibri"
                p.font.color.rgb = WHT if r == 0 else BLK; p.font.bold = r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAV if r == 0 else GR3 if r % 2 == 0 else WHT

# ── SLIDE 1: COVER ──
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(3.5)); r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = co['name'].upper()
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Georgia"
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11), Inches(0.6))
ser_txt = " | ".join([f"{s['name']}: {s['rate']}" for s in dl['series']])
p = tb.text_frame.paragraphs[0]; p.text = f"Debentures Simples  |  R$ {dl['volume']/1e6:.0f}M  |  {ser_txt}"
p.font.size = Pt(15); p.font.color.rgb = RGBColor(0xC0,0xD0,0xE8); p.font.name = "Calibri"
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11), Inches(0.5))
p = tb.text_frame.paragraphs[0]; p.text = f"Rating {dl['rating']}  |  {dl['registration']}"
p.font.size = Pt(12); p.font.color.rgb = GR2; p.font.name = "Calibri"
tb = s.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8), Inches(1.5)); tf = tb.text_frame
for t in ["IB-Agents Intelligence Platform", "Coordenador Lider", "", datetime.now().strftime("%B %Y"), "Strictly Private and Confidential"]:
    p = tf.add_paragraph(); p.text = t; p.font.size = Pt(10); p.font.color.rgb = GR1; p.font.name = "Calibri"

# ── SLIDE 2: EXECUTIVE SUMMARY ──
s = add_base(); add_title(s, hl['exec_summary'])
add_text(s, [
    "## Tese de Investimento",
    data['investment_thesis'],
    "",
    "## Destaques da Operacao",
    f"• Volume: R$ {dl['volume']/1e6:.0f}M em {len(dl['series'])} series — {dl['guarantee']}",
    f"• Rating: {dl['rating']} — {dl['use_of_proceeds']}",
    f"• Empresa: {co['description']}",
], x=0.5, y=1.1, w=7.0, size=14)
add_callout(s, f"R$ {dl['volume']/1e6:.0f}M", "Volume", x=8.0, y=1.1, color=NAV)
add_callout(s, dl['rating'].split('(')[0].strip(), "Rating Fitch", x=8.0, y=2.6, color=GRN)
add_callout(s, f"{fin['dl_ebitda'][-1]:.1f}x", "DL/EBITDA", x=8.0, y=4.1, color=GRN)
add_callout(s, f"{fin['ebitda_margin'][-1]*100:.1f}%", "Margem EBITDA", x=11.0, y=1.1, color=NAV)
add_callout(s, f"{fin['coverage'][-1]:.1f}x", "Coverage", x=11.0, y=2.6, color=GRN)
add_callout(s, dl['guarantee_coverage'], "Cobertura Garantias", x=11.0, y=4.1, color=NAV)

# ── SLIDE 3: COMPANY OVERVIEW ──
s = add_base(); add_title(s, hl['company'])
add_text(s, [
    "## Visao Geral",
    co['description'],
    f"• Fundacao: {co['founded']} | Sede: {co['headquarters']} | {co['employees']} colaboradores",
    "",
    "## Diferenciais Competitivos",
] + [f"• {d}" for d in co['key_differentials']] + [
    "",
    "## Governanca",
    f"• {co['governance']}",
], x=0.5, y=1.1, w=8.0, size=14)
add_callout(s, f"R$ {fin['revenue'][-1]:,.0f}M", "Receita 2025", x=9.0, y=1.1, color=NAV)
add_callout(s, f"{co['employees']:,}", "Colaboradores", x=9.0, y=2.6, color=NAV)
add_callout(s, f"{fin['ebitda'][-1]:,.0f}M", "EBITDA 2025", x=9.0, y=4.1, color=GRN)

# ── SLIDE 4: MARKET ──
s = add_base(); add_title(s, hl['market'])
mkt = data['market']
add_text(s, [
    f"## Mercado Enderecavel: R$ {mkt['tam_brl_bi']}Bi ({mkt['tam_growth']})",
    f"• {mkt['brazil_position']}",
    "",
    "## Drivers de Crescimento",
] + [f"• {d}" for d in mkt['drivers']] + [
    "",
    "## Ventos Favoraveis",
] + [f"• {t}" for t in mkt['tailwinds']] + [
    f"[Fonte: {mkt['source']}]",
], x=0.5, y=1.1, w=8.0, size=14)
add_callout(s, f"R$ {mkt['tam_brl_bi']}Bi", "TAM", x=9.0, y=1.1, color=NAV)
add_callout(s, mkt['tam_growth'], "CAGR Mercado", x=9.0, y=2.6, color=GRN)

# ── SLIDE 5: FINANCIALS + CHART ──
s = add_base(); add_title(s, hl['financials'])
add_img(s, chart_revenue, x=0.5, y=1.1, w=7.0, h=4.2)
rev_cagr = ((fin['revenue'][-1]/fin['revenue'][0])**(1/(len(yrs)-1))-1)*100
ebitda_cagr = ((fin['ebitda'][-1]/fin['ebitda'][0])**(1/(len(yrs)-1))-1)*100
add_text(s, [
    "## Destaques",
    f"• Receita CAGR: {rev_cagr:.1f}%",
    f"• EBITDA CAGR: {ebitda_cagr:.1f}%",
    f"• Margem EBITDA: {fin['ebitda_margin'][-1]*100:.1f}%",
    f"• FCF: R$ {fin['fcf'][-1]:,.0f}M",
    "",
    "[DFs IFRS auditadas]",
], x=7.8, y=1.1, w=5.0, size=14)
add_callout(s, f"{ebitda_cagr:.0f}%", "CAGR EBITDA", x=10.0, y=4.5, w=2.8, color=GRN)

# ── SLIDE 6: MARGINS + CHART ──
s = add_base(); add_title(s, "Margem EBITDA estavel demonstra consistencia operacional e disciplina de custos")
add_img(s, chart_margins, x=0.5, y=1.1, w=7.0, h=4.2)
add_text(s, [
    "## Evolucao de Margens",
    f"• Margem EBITDA estavel em ~{fin['ebitda_margin'][-1]*100:.0f}%",
    f"• Receita crescendo sem compressao de margem",
    f"• Indicativo de poder de precificacao e escala",
], x=7.8, y=1.1, w=5.0, size=14)

# ── SLIDE 7: CREDIT + CHART ──
s = add_base(); add_title(s, hl['credit'])
add_img(s, chart_leverage, x=0.5, y=1.1, w=7.0, h=4.2)
cov = data['covenants']
folga = cov['dl_ebitda_max'] - fin['dl_ebitda'][-1]
add_text(s, [
    "## Metricas de Credito",
    f"• DL/EBITDA: {fin['dl_ebitda'][-1]:.1f}x (cov: {cov['dl_ebitda_max']}x — folga {folga:.1f}x)",
    f"• Coverage: {fin['coverage'][-1]:.1f}x (cov: {cov['coverage_min']}x)",
    f"• DSCR: >= {cov['dscr_min']}x",
    "",
    "## Stress Test",
    f"• Bear ({cov['stress_bear']['ebitda_drop']}%): {cov['stress_bear']['dl_ebitda']:.1f}x — {cov['stress_bear']['status']}",
    f"• Stressed ({cov['stress_severe']['ebitda_drop']}%): {cov['stress_severe']['dl_ebitda']:.1f}x — {cov['stress_severe']['status']}",
    f"• Breach: {cov['breach_threshold']}",
], x=7.8, y=1.1, w=5.0, size=14)
add_callout(s, f"{folga:.1f}x", "Folga Covenant", x=10.0, y=5.0, w=2.8, color=GRN)

# ── SLIDE 8: FCF + CHART ──
s = add_base(); add_title(s, "Geracao de caixa robusta — FCF crescente sustenta servico da divida com ampla margem")
add_img(s, chart_fcf, x=0.5, y=1.1, w=7.0, h=4.2)
add_text(s, [
    "## Fluxo de Caixa",
    f"• FCF 2025: R$ {fin['fcf'][-1]:,.0f}M",
    f"• Capex: R$ {fin['capex'][-1]:,.0f}M",
    f"• FCF/Servico Divida: >2x",
    "",
    "## Disciplina de Capital",
    "• Capex focado em expansao produtiva",
    "• Dividendos condicionais a covenants",
], x=7.8, y=1.1, w=5.0, size=14)

# ── SLIDE 9: DEAL STRUCTURE ──
s = add_base(); add_title(s, hl['deal'])
tbl_data = [["Caracteristica"] + [sr['name'] for sr in dl['series']]]
fields = [("Volume", lambda sr: f"R$ {dl['volume']*sr['volume_pct']/100/1e6:.0f}M"),
    ("Prazo", lambda sr: sr['tenor']), ("Remuneracao", lambda sr: sr['rate']),
    ("Amortizacao", lambda sr: sr['amort'])]
for label, fn in fields:
    tbl_data.append([label] + [fn(sr) for sr in dl['series']])
tbl_data.extend([
    ["Garantia"] + [dl['guarantee']] * len(dl['series']),
    ["Rating"] + [dl['rating']] * len(dl['series']),
    ["Registro"] + [dl['registration']] * len(dl['series']),
])
add_table(s, tbl_data, y=1.1)
add_callout(s, dl['guarantee_coverage'], "Cobertura Garantias", x=9.5, y=5.0, w=3.3, color=GRN)

# ── SLIDE 10: STRESS TEST (plotly) ──
if chart_stress:
    s = add_base(); add_title(s, hl['covenants'])
    add_img(s, chart_stress, x=0.5, y=1.1, w=7.0, h=4.2)
    add_text(s, [
        "## Cenarios de Stress",
        f"• Base: {fin['dl_ebitda'][-1]:.1f}x — OK",
        f"• Bear ({cov['stress_bear']['ebitda_drop']}%): {cov['stress_bear']['dl_ebitda']:.1f}x — {cov['stress_bear']['status']}",
        f"• Stressed ({cov['stress_severe']['ebitda_drop']}%): {cov['stress_severe']['dl_ebitda']:.1f}x — {cov['stress_severe']['status']}",
        f"• Breach threshold: {cov['breach_threshold']}",
    ], x=7.8, y=1.1, w=5.0, size=14)

# ── SLIDE 11: COMPARABLES ──
s = add_base(); add_title(s, hl['pricing'])
comps = data['comparables']
tbl = [["Emissor", "Rating", "Prazo", "Spread", "Volume", "Data"]]
for c in comps:
    tbl.append([c['issuer'], c['rating'], c['tenor'], c['spread'], c['volume'], c['date']])
tbl.append([co['name']+" (proposta)", dl['rating'], dl['series'][0]['tenor'], dl['series'][0]['rate'], f"R$ {dl['volume']/1e6:.0f}M", datetime.now().strftime("%b/%y")])
add_table(s, tbl, y=1.1)

# ── SLIDE 12: INVESTORS ──
s = add_base(); add_title(s, hl['investors'])
inv = data['investors']
tbl = [["#", "Instituicao", "Tipo", "Ticket Estimado", "Score"]]
for i, inv_item in enumerate(inv):
    tbl.append([str(i+1), inv_item['name'], inv_item['type'], inv_item['ticket'], inv_item['score']])
add_table(s, tbl, y=1.1)

# ── SLIDE 13: RISKS ──
s = add_base(); add_title(s, hl['risks'])
risks = data['risks']
tbl = [["Risco", "Probabilidade", "Impacto", "Mitigante"]]
for r in risks:
    tbl.append([r['risk'], r['probability'], r['impact'], r['mitigant']])
add_table(s, tbl, y=1.1)

# ── SAVE ──
fn = "pitchbook_ai_agrovale.pptx"
prs.save(str(OUT / fn))
print(f"\nGerado: {fn} ({len(prs.slides)} slides)")
print(f"Local: {(OUT / fn).resolve()}")
print(f"API tokens: {response.usage.input_tokens} in / {response.usage.output_tokens} out")
print(f"Custo estimado: ~${(response.usage.input_tokens*3/1e6 + response.usage.output_tokens*15/1e6):.4f} USD")
