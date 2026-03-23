"""
Gerador de PPTX — Book de Credito CRA Grupo Meridian Alimentos S.A.
Padrao Goldman Sachs / McKinsey | Design Farallon-inspired
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── PALETA FARALLON ──────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0F, 0x1F, 0x3D)
NAVY_MED  = RGBColor(0x1A, 0x2B, 0x4A)
GRAPHITE  = RGBColor(0x2D, 0x37, 0x48)
MUTED     = RGBColor(0x71, 0x80, 0x96)
LIGHT_BG  = RGBColor(0xE2, 0xE8, 0xF0)
BORDER    = RGBColor(0xA0, 0xAE, 0xC0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
NEGATIVE  = RGBColor(0xC5, 0x30, 0x30)
POSITIVE  = RGBColor(0x27, 0x67, 0x49)
GOLD      = RGBColor(0xD4, 0xA8, 0x53)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── HELPERS ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=11, bold=False, color=None, align=PP_ALIGN.LEFT,
             font="Calibri", wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or GRAPHITE
    return txb

def slide_base(title_text, subtitle_text=None, show_footer=True):
    """Cria slide com estrutura padrao: header navy + area branca + footer"""
    slide = prs.slides.add_slide(BLANK)

    # Fundo branco
    add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)

    # Header navy
    add_rect(slide, 0, 0, 13.33, 1.1, fill=NAVY)

    # Headline (conclusao, nao titulo)
    add_text(slide, title_text,
             l=0.4, t=0.15, w=12.5, h=0.75,
             size=18, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.LEFT)

    if subtitle_text:
        add_text(slide, subtitle_text,
                 l=0.4, t=0.78, w=12.5, h=0.3,
                 size=9, bold=False, color=LIGHT_BG, font="Calibri")

    # Linha separadora
    add_rect(slide, 0, 1.1, 13.33, 0.02, fill=BORDER)

    if show_footer:
        add_rect(slide, 0, 7.15, 13.33, 0.35, fill=LIGHT_BG)
        add_text(slide, "Grupo Meridian Alimentos S.A. | Book de Credito — CRA R$150M | Confidencial — RCVM 160/2022 | Exclusivo para Investidores Qualificados",
                 l=0.3, t=7.17, w=10, h=0.2,
                 size=7, color=MUTED, font="Calibri")
        add_text(slide, "2026",
                 l=12.8, t=7.17, w=0.5, h=0.2,
                 size=7, color=MUTED, font="Calibri", align=PP_ALIGN.RIGHT)
    return slide


def kpi_box(slide, l, t, w, h, value, label, color=NAVY):
    add_rect(slide, l, t, w, h, fill=LIGHT_BG)
    add_text(slide, value,
             l=l+0.08, t=t+0.08, w=w-0.16, h=h*0.55,
             size=22, bold=True, color=color, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, label,
             l=l+0.08, t=t+h*0.58, w=w-0.16, h=h*0.38,
             size=8, color=MUTED, font="Calibri", align=PP_ALIGN.CENTER)


def bullet_block(slide, items, l, t, w, h, title=None, title_color=NAVY):
    if title:
        add_text(slide, title, l=l, t=t, w=w, h=0.22,
                 size=9, bold=True, color=title_color, font="Calibri")
        t += 0.24
        h -= 0.24
    row_h = min(h / len(items), 0.32)
    for i, item in enumerate(items):
        add_text(slide, f"  {item}", l=l, t=t + i*row_h, w=w, h=row_h,
                 size=9, color=GRAPHITE, font="Calibri")


def table_header(slide, cols, l, t, col_w, row_h=0.28):
    for j, col in enumerate(cols):
        add_rect(slide, l + j*col_w, t, col_w - 0.02, row_h, fill=NAVY)
        add_text(slide, col,
                 l=l + j*col_w + 0.05, t=t+0.04, w=col_w-0.1, h=row_h-0.08,
                 size=8, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)


def table_row(slide, cols, l, t, col_w, row_h=0.24, bg=WHITE, text_color=GRAPHITE, bold=False):
    for j, val in enumerate(cols):
        add_rect(slide, l + j*col_w, t, col_w - 0.02, row_h, fill=bg)
        add_text(slide, val,
                 l=l + j*col_w + 0.05, t=t+0.03, w=col_w-0.1, h=row_h-0.06,
                 size=8, bold=bold, color=text_color, font="Calibri", align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 1 — CAPA
# ============================================================================
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(slide, 0, 5.8, 13.33, 1.7, fill=NAVY_MED)
add_rect(slide, 0.4, 1.2, 0.08, 4.0, fill=GOLD)  # linha dourada vertical

add_text(slide, "GRUPO MERIDIAN ALIMENTOS S.A.",
         l=0.65, t=1.2, w=11, h=0.6,
         size=28, bold=True, color=WHITE, font="Georgia")
add_text(slide, "Book de Credito — Oferta de CRA",
         l=0.65, t=1.85, w=9, h=0.45,
         size=18, bold=False, color=GOLD, font="Georgia", italic=True)
add_rect(slide, 0.65, 2.45, 5, 0.03, fill=GOLD)

add_text(slide, "R$ 150.000.000",
         l=0.65, t=2.65, w=6, h=0.55,
         size=22, bold=True, color=WHITE, font="Calibri")
add_text(slide, "Certificados de Recebiveis do Agronegocio — IPCA + 7,25% a.a. — Serie Unica — Prazo: 5 anos",
         l=0.65, t=3.22, w=10, h=0.35,
         size=11, color=LIGHT_BG, font="Calibri")

add_text(slide, "Rating: AA- (Fitch) | Regime: 476/CVM 160 | Lastro: Recebiveis Agroindustriais",
         l=0.65, t=3.7, w=10, h=0.3,
         size=10, color=MUTED, font="Calibri")

add_text(slide, "Coordenador Lider: IB Agents Capital | Securitizadora: Eco Securitizacao S.A.",
         l=0.65, t=4.15, w=10, h=0.28,
         size=9, color=MUTED, font="Calibri")

add_text(slide, "Maro / 2026",
         l=0.65, t=5.9, w=4, h=0.28,
         size=9, color=MUTED, font="Calibri")
add_text(slide, "ESTRITAMENTE CONFIDENCIAL — Destinado exclusivamente a Investidores Qualificados conforme RCVM 160/2022\nEste material nao constitui oferta publica de valores mobiliarios. As informacoes aqui contidas sao baseadas em fontes consideradas confiaveis.",
         l=0.4, t=6.25, w=12.5, h=0.5,
         size=7, color=MUTED, font="Calibri")


# ============================================================================
# SLIDE 2 — INVESTMENT HIGHLIGHTS
# ============================================================================
slide = slide_base(
    "Cinco razoes estruturais para investir no CRA Meridian — retorno IPCA+7,25% com protecao AA-",
    "Investment Highlights | Executive Summary"
)

highlights = [
    ("1", "Lider no segmento de aves premium no Centro-Oeste com 18% de share de mercado",
     "IPCA+7,25%", "Taxa de emissao"),
    ("2", "EBITDA de R$187M em 2025 com CAGR de 22,3% a.a. nos ultimos 5 anos",
     "AA-", "Rating Fitch"),
    ("3", "Cobertura de divida (DSCR) de 2,8x — covenant minimo de 1,5x com folga de 87%",
     "2,8x", "DSCR 2025"),
    ("4", "Garantias de R$225M (1,5x de cobertura): recebiveis agroindustriais + imoveis rurais",
     "1,5x", "Cobertura Garantias"),
    ("5", "Destinacao exclusiva: expansao da planta de Jatai (GO) — ROI projetado de 28% a.a.",
     "R$150M", "Volume CRA"),
]

for i, (num, text, kpi, kpi_label) in enumerate(highlights):
    y = 1.3 + i * 1.05
    add_rect(slide, 0.3, y, 0.35, 0.35, fill=NAVY)
    add_text(slide, num, l=0.3, t=y, w=0.35, h=0.35,
             size=14, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)
    add_text(slide, text, l=0.75, t=y+0.02, w=9.0, h=0.32,
             size=10, color=GRAPHITE, font="Calibri")
    add_rect(slide, 10.0, y-0.02, 2.9, 0.42, fill=LIGHT_BG)
    add_text(slide, kpi, l=10.05, t=y-0.02, w=2.8, h=0.28,
             size=16, bold=True, color=NAVY, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, kpi_label, l=10.05, t=y+0.24, w=2.8, h=0.18,
             size=7, color=MUTED, font="Calibri", align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 3 — VISAO GERAL DA EMPRESA
# ============================================================================
slide = slide_base(
    "Meridian e a maior processadora independente de proteinas premium no Centro-Oeste — R$1,2Bi de receita",
    "Modulo 3 — Visao Geral da Empresa"
)

add_text(slide, "PERFIL DA EMPRESA", l=0.4, t=1.25, w=5.5, h=0.22,
         size=9, bold=True, color=NAVY, font="Calibri")
perfil = [
    "Fundacao: 1998, Jatai (GO)",
    "Atividade: Abate, processamento e exportacao de frangos e suinos",
    "Capacidade: 280.000 aves/dia (3 plantas: Jatai, Rio Verde, Mineiros)",
    "Clientes: 340+ redes de varejo + 18 paises exportacao",
    "Colaboradores: 6.800 diretos",
    "Modelo: Verticalizacao total — genetica, incubatorio, granja, abate, distribuicao",
]
for i, item in enumerate(perfil):
    add_text(slide, f"  • {item}", l=0.4, t=1.52+i*0.35, w=5.8, h=0.32,
             size=9.5, color=GRAPHITE, font="Calibri")

add_text(slide, "HISTORICO & TRAJETORIA", l=7.2, t=1.25, w=5.8, h=0.22,
         size=9, bold=True, color=NAVY, font="Calibri")
timeline = [
    ("1998", "Fundacao — 1 planta, 40.000 aves/dia"),
    ("2008", "Certificacao habilitacao sanitaria MAPA — inicio exportacoes"),
    ("2014", "Abertura 2a planta (Rio Verde) — capacidade 120.000 aves/dia"),
    ("2018", "Aquisicao Mineiros — integracao suinos (20.000 cab/mes)"),
    ("2022", "Receita cruza R$800M — 1a emissao de debentures (R$80M, CDI+2,1%)"),
    ("2025", "Receita R$1,2Bi | EBITDA R$187M | Rating AA- Fitch"),
]
for i, (year, desc) in enumerate(timeline):
    add_rect(slide, 7.2, 1.52+i*0.35, 0.55, 0.28, fill=NAVY)
    add_text(slide, year, l=7.2, t=1.52+i*0.35, w=0.55, h=0.28,
             size=8, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, desc, l=7.82, t=1.52+i*0.35+0.02, w=5.1, h=0.26,
             size=9, color=GRAPHITE, font="Calibri")

# KPI boxes
kpis = [("R$1,2Bi","Receita 2025"), ("R$187M","EBITDA 2025"), ("15,8%","Margem EBITDA"), ("2,1x","Divida/EBITDA")]
for i, (v, l) in enumerate(kpis):
    kpi_box(slide, 0.3 + i*3.1, 5.85, 2.85, 0.88, v, l, NAVY)

add_rect(slide, 6.3, 1.1, 0.02, 5.8, fill=BORDER)


# ============================================================================
# SLIDE 4 — MERCADO: TAM/SAM/SOM
# ============================================================================
slide = slide_base(
    "TAM de R$379Bi em crescimento de 8,3% a.a. — Meridian capturou apenas 0,3% do mercado enderecavel",
    "Modulo 4 — Analise de Mercado | TAM / SAM / SOM"
)

# Circulos concentricos (representados por retangulos aninhados)
add_rect(slide, 0.4, 1.25, 6.0, 5.5, fill=RGBColor(0xEB, 0xF4, 0xFF))
add_rect(slide, 1.2, 1.8, 4.4, 4.4, fill=RGBColor(0xBF, 0xD9, 0xFF))
add_rect(slide, 2.1, 2.4, 2.6, 3.2, fill=NAVY)

add_text(slide, "TAM", l=0.5, t=1.35, w=1.5, h=0.3, size=9, bold=True, color=NAVY, font="Calibri")
add_text(slide, "R$379Bi", l=0.5, t=1.65, w=1.5, h=0.3, size=11, bold=True, color=NAVY, font="Calibri")
add_text(slide, "Proteinas animais Brasil\nCAGR 8,3% a.a. (ABPA/ABIEC)", l=0.5, t=1.95, w=1.8, h=0.4, size=7, color=MUTED, font="Calibri")

add_text(slide, "SAM", l=1.3, t=1.9, w=1.5, h=0.3, size=9, bold=True, color=NAVY_MED, font="Calibri")
add_text(slide, "R$42Bi", l=1.3, t=2.2, w=1.5, h=0.3, size=11, bold=True, color=NAVY_MED, font="Calibri")
add_text(slide, "Aves premium CO\nCAGR 11,2% a.a.", l=1.3, t=2.5, w=1.8, h=0.35, size=7, color=MUTED, font="Calibri")

add_text(slide, "SOM", l=2.2, t=2.5, w=2.2, h=0.28, size=9, bold=True, color=WHITE, font="Calibri")
add_text(slide, "R$1,2Bi", l=2.2, t=2.78, w=2.2, h=0.35, size=16, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
add_text(slide, "Share atual: 0,3% TAM\n18% SAM premium CO", l=2.2, t=3.15, w=2.2, h=0.4, size=7.5, color=LIGHT_BG, font="Calibri", align=PP_ALIGN.CENTER)

# Metodologia
add_text(slide, "METODOLOGIA BOTTOM-UP", l=6.8, t=1.25, w=6.2, h=0.22,
         size=9, bold=True, color=NAVY, font="Calibri")
metodologia = [
    "TAM: 14,2M ton abatidas/ano (ABPA 2025) x ticket medio R$26,7/kg = R$379Bi",
    "SAM: segmento premium (cortes especiais, congelados, exportacao habilitada)",
    "  • Centro-Oeste: 18% do abate nacional | premium: 11% do segmento",
    "  • SAM = R$379Bi x 18% x 65% = R$42Bi",
    "SOM: share atual da Meridian no SAM = R$1,2Bi / R$42Bi = 2,9%",
    "  • Meta 2030: R$2,1Bi = 5,0% do SAM (CAGR 15%)",
]
for i, line in enumerate(metodologia):
    add_text(slide, line, l=6.8, t=1.52+i*0.38, w=6.1, h=0.35,
             size=9, color=GRAPHITE if not line.startswith("  ") else MUTED, font="Calibri")

add_text(slide, "Fontes: ABPA Relatorio Anual 2025 | ABIEC | MAPA/DIPOA | IBGE PAM 2024 | Calculo proprio Meridian",
         l=0.4, t=6.95, w=12.5, h=0.18, size=7, color=MUTED, font="Calibri", italic=True)


# ============================================================================
# SLIDE 5 — 4 DRIVERS DE CRESCIMENTO
# ============================================================================
slide = slide_base(
    "Quatro drivers estruturais sustentam CAGR de 11% no mercado de proteinas premium ate 2030",
    "Modulo 4 — Analise de Mercado | Drivers de Crescimento"
)

drivers = [
    ("DEMANDA GLOBAL",
     "Populacao mundial atinge 8,5Bi em 2030 — proteina animal e o segmento de maior crescimento (FAO: +14% ate 2032)",
     ["Consumo per capita de frango no Brasil: 42,2kg/ano (3o maior do mundo — UBABEF 2025)",
      "Exportacoes brasileiras de frango: US$10,2Bi em 2025 (+18% vs 2024) — SECEX/MDIC",
      "Habilitacoes sanitarias MAPA: Brasil habilita 3 novos paises/ano em media"]),
    ("EXPORTACOES & CAMBIO",
     "BRL/USD de R$5,80 (BACEN Focus dez/25) beneficia exportadores — Meridian exporta 28% da receita",
     ["Mercados-alvo: Asia (38%), Oriente Medio (31%), Europa (22%), outros (9%)",
      "Margens de exportacao 4-6pp acima do mercado interno (hedge natural cambial)",
      "Novos mercados habilitados em 2025: Indonesia, Marrocos, Jordania"]),
    ("PREMIUMIZACAO",
     "Cortes especiais e semi-elaborados crescem 2,3x mais rapido que frango inteiro (Nielsen 2025)",
     ["Participacao de produtos de VA na receita Meridian: 41% em 2025 vs 28% em 2021",
      "Ticket medio de cortes premium: R$34/kg vs R$19/kg no frango inteiro",
      "Margem bruta premium: 38% vs 24% no produto in natura"]),
    ("CONSOLIDACAO DO SETOR",
     "M&A acelera: 14 transacoes no setor de proteinas em 2023-25 — independentes buscam escala",
     ["Top 5 players concentram 62% do abate nacional — fragmentacao no mid-market",
      "Meridian e alvo potencial de M&A / plataforma de consolidacao regional",
      "CAPEX de expansao Jatai gera +35.000 aves/dia = +12,5% de capacidade"]),
]

for i, (title, subtitle, bullets) in enumerate(drivers):
    col = i % 2
    row = i // 2
    l = 0.3 + col * 6.5
    t = 1.25 + row * 2.85

    add_rect(slide, l, t, 6.2, 2.7, fill=LIGHT_BG)
    add_rect(slide, l, t, 6.2, 0.28, fill=NAVY)
    add_text(slide, title, l=l+0.1, t=t+0.04, w=6.0, h=0.22,
             size=9, bold=True, color=WHITE, font="Calibri")
    add_text(slide, subtitle, l=l+0.1, t=t+0.32, w=6.0, h=0.38,
             size=8.5, color=NAVY_MED, font="Calibri", bold=True)
    for j, b in enumerate(bullets):
        add_text(slide, f"• {b}", l=l+0.1, t=t+0.74+j*0.55, w=6.0, h=0.52,
                 size=8.5, color=GRAPHITE, font="Calibri")

add_text(slide, "Fontes: FAO Food Outlook 2025 | ABPA | SECEX/MDIC | Nielsen Retail 2025 | BACEN Focus dez/2025",
         l=0.3, t=6.95, w=12.5, h=0.18, size=7, color=MUTED, font="Calibri", italic=True)


# ============================================================================
# SLIDE 6 — LANDSCAPE COMPETITIVO
# ============================================================================
slide = slide_base(
    "Meridian lidera o segmento premium no Centro-Oeste — vantagem competitiva em 4 dos 5 atributos criticos",
    "Modulo 4 — Analise de Mercado | Landscape Competitivo"
)

cols_h = ["Empresa", "Market Cap", "Receita", "Abate/dia", "Premium", "Exportacao", "Certificacoes", "Rating"]
col_w = 12.6 / len(cols_h)
table_header(slide, cols_h, l=0.35, t=1.3, col_w=col_w)

rows = [
    (["JBS (JBSS3)", "R$68Bi", "R$385Bi", "5.200.000", "Sim", "52%", "EUA/China/UE/Halal", "BBB- S&P"], WHITE, GRAPHITE, False),
    (["Marfrig (MRFG3)", "R$9,2Bi", "R$89Bi", "1.800.000", "Sim", "41%", "EUA/UE/Halal", "B+ Fitch"], WHITE, GRAPHITE, False),
    (["Minerva (BEEF3)", "R$4,1Bi", "R$38Bi", "890.000", "Parcial", "68%", "UE/Oriente Medio", "BB- S&P"], WHITE, GRAPHITE, False),
    (["Cooperativa XYZ", "N/A", "R$3,2Bi", "310.000", "Parcial", "12%", "MAPA basico", "N/R"], WHITE, GRAPHITE, False),
    (["MERIDIAN (alvo)", "Nao listada", "R$1,2Bi", "280.000", "Sim", "28%", "MAPA/USDA/HALAL/CE", "AA- Fitch"], RGBColor(0xE8, 0xF0, 0xFF), NAVY, True),
]

for i, (data, bg, tc, bd) in enumerate(rows):
    table_row(slide, data, l=0.35, t=1.58+i*0.3, col_w=col_w, bg=bg, text_color=tc, bold=bd)

add_text(slide, "POSICIONAMENTO ESTRATEGICO", l=0.35, t=3.2, w=12.5, h=0.22,
         size=9, bold=True, color=NAVY, font="Calibri")
pos_items = [
    "Meridian e o maior independente no segmento premium de aves no Centro-Oeste — sem concorrente direto na faixa de R$1-5Bi com o mesmo perfil de verticalizacao e certificacoes",
    "Vantagem em certificacoes: USDA + Halal + Uniao Europeia = acesso a mercados premium com spreads de R$4-8/kg acima da media",
    "Escala suficiente para capturar eficiencias mas abaixo do radar de M&A defensivo dos grandes players — 'sweet spot' para crescimento organico + inorganico",
]
for i, item in enumerate(pos_items):
    add_text(slide, f"  • {item}", l=0.35, t=3.46+i*0.45, w=12.5, h=0.42,
             size=9, color=GRAPHITE, font="Calibri")

add_text(slide, "Fontes: B3 (Nov/2025) | Relatorios anuais JBS/Marfrig/Minerva | MAPA/DIPOA | Estimativas Meridian",
         l=0.35, t=6.95, w=12.5, h=0.18, size=7, color=MUTED, font="Calibri", italic=True)


# ============================================================================
# SLIDE 7 — VANTAGENS COMPETITIVAS
# ============================================================================
slide = slide_base(
    "Tres moats estruturais protegem a posicao de Meridian — replicacao estimada em R$800M e 7 anos",
    "Modulo 5 — Vantagens Competitivas"
)

moats = [
    ("INTEGRACAO VERTICAL TOTAL",
     "Do ovo ao prato: genetica propria → incubatorio → granja integrada (380 parceiros) → abate → distribuicao refrigerada propria",
     [("Custo de producao", "R$3,82/kg", "vs R$4,65/kg media do setor (-18%)"),
      ("Mortalidade", "2,1%", "vs 3,8% media do setor — genetica propria"),
      ("Conversao alimentar", "1,71kg/kg", "vs 1,85kg/kg media — ganho de escala")]),
    ("CERTIFICACOES PREMIUM",
     "Portfolio de certificacoes habilita os 3 mercados premium globais simultaneamente — barreira quasi-intransponivel para novos entrantes",
     [("USDA-FSIS", "Habilitado desde 2011", "=> mercado EUA aberto"),
      ("Halal (BVQI)", "12 paises habilitados", "=> Oriente Medio + Asia"),
      ("EU Health Cert.", "Habilitado CE/UK", "=> UE + UK pos-Brexit")]),
    ("ESCALA + LOCALIZACAO",
     "Centro-Oeste: maior produtor de graos do Brasil — custo de racao 12% abaixo da media nacional (proximidade CONAB/Jatai)",
     [("Racao propria", "62% da necessidade", "=> hedge de custo de graos"),
      ("Distancia porto", "1.050km Paranagua", "vs 1.800km media GO"),
      ("Energia", "R$0,31/kWh", "vs R$0,42/kWh media industrial SP")]),
]

for i, (title, desc, kpis_m) in enumerate(moats):
    l = 0.3 + i * 4.3
    add_rect(slide, l, 1.25, 4.1, 5.5, fill=LIGHT_BG)
    add_rect(slide, l, 1.25, 4.1, 0.32, fill=NAVY)
    add_text(slide, f"MOAT {i+1}", l=l+0.1, t=1.27, w=0.65, h=0.28,
             size=7, bold=True, color=GOLD, font="Calibri")
    add_text(slide, title, l=l+0.78, t=1.27, w=3.2, h=0.28,
             size=8, bold=True, color=WHITE, font="Calibri")
    add_text(slide, desc, l=l+0.1, t=1.62, w=3.9, h=0.75,
             size=8.5, color=GRAPHITE, font="Calibri")
    for j, (k, v, note) in enumerate(kpis_m):
        y = 2.45 + j * 1.0
        add_rect(slide, l+0.1, y, 3.9, 0.9, fill=WHITE)
        add_text(slide, k, l=l+0.2, t=y+0.06, w=3.7, h=0.22,
                 size=8, color=MUTED, font="Calibri")
        add_text(slide, v, l=l+0.2, t=y+0.28, w=3.7, h=0.3,
                 size=14, bold=True, color=NAVY, font="Calibri")
        add_text(slide, note, l=l+0.2, t=y+0.58, w=3.7, h=0.25,
                 size=7.5, color=POSITIVE, font="Calibri")


# ============================================================================
# SLIDE 8 — DRE HISTORICA
# ============================================================================
slide = slide_base(
    "EBITDA cresceu de R$52M para R$187M em 5 anos (CAGR 22,3%) — expansao de 620bps na margem",
    "Modulo 6 — Desempenho Financeiro | DRE Historica 2021-2025"
)

anos = ["2021", "2022", "2023", "2024", "2025E"]
col_w = 2.0
cols_h = ["Indicador (R$M)"] + anos + ["CAGR"]
table_header(slide, cols_h, l=0.35, t=1.3, col_w=12.0/len(cols_h))
cw = 12.0 / len(cols_h)

dre_data = [
    (["Receita Liquida", "614", "748", "921", "1.051", "1.218", "18,7%"], WHITE, GRAPHITE),
    (["Custo dos Produtos Vendidos", "(487)", "(586)", "(708)", "(798)", "(915)", "—"], LIGHT_BG, GRAPHITE),
    (["Lucro Bruto", "127", "162", "213", "253", "303", "24,3%"], WHITE, GRAPHITE),
    (["Margem Bruta", "20,7%", "21,7%", "23,1%", "24,1%", "24,9%", "+620bps"], WHITE, POSITIVE),
    (["EBITDA Ajustado", "52", "72", "105", "151", "187", "22,3%"], RGBColor(0xE8,0xF0,0xFF), NAVY),
    (["Margem EBITDA", "8,5%", "9,6%", "11,4%", "14,4%", "15,4%", "+690bps"], RGBColor(0xE8,0xF0,0xFF), POSITIVE),
    (["D&A", "(28)", "(31)", "(35)", "(40)", "(44)", "—"], WHITE, GRAPHITE),
    (["EBIT", "24", "41", "70", "111", "143", "56,4%"], WHITE, GRAPHITE),
    (["Resultado Financeiro Liq.", "(18)", "(22)", "(28)", "(36)", "(41)", "—"], WHITE, NEGATIVE),
    (["Lucro Liquido", "4", "13", "29", "52", "71", "104,5%"], WHITE, NAVY),
    (["Margem Liquida", "0,7%", "1,7%", "3,2%", "4,9%", "5,8%", "+510bps"], WHITE, POSITIVE),
]

for i, (data, bg, tc) in enumerate(dre_data):
    bd = data[0] in ["EBITDA Ajustado", "Margem EBITDA", "Lucro Liquido"]
    row_data = [data[0]] + data[1:]
    table_row(slide, row_data, l=0.35, t=1.58+i*0.42, col_w=cw, bg=bg, text_color=tc, bold=bd)

add_text(slide, "Fontes: DFs auditadas Meridian (EY, 2021-2024) | 2025E: estimativa da administracao (nao auditada) | EBITDA ajustado exclui LTIP e itens nao recorrentes",
         l=0.35, t=6.95, w=12.5, h=0.18, size=7, color=MUTED, font="Calibri", italic=True)


# ============================================================================
# SLIDE 9 — METRICAS DE CREDITO
# ============================================================================
slide = slide_base(
    "Alavancagem de 2,1x e DSCR de 2,8x demonstram capacidade robusta de servico da divida",
    "Modulo 6 — Desempenho Financeiro | Metricas de Credito 2021-2025"
)

cols_h2 = ["Metrica", "2021", "2022", "2023", "2024", "2025E", "Covenant"]
cw2 = 12.6 / len(cols_h2)
table_header(slide, cols_h2, l=0.35, t=1.3, col_w=cw2)

credit_rows = [
    (["Divida Bruta Total (R$M)", "158", "189", "234", "278", "310", "—"], WHITE, GRAPHITE),
    (["(-) Caixa e Equiv. (R$M)", "(42)", "(51)", "(68)", "(83)", "(98)", "—"], WHITE, GRAPHITE),
    (["Divida Liquida (R$M)", "116", "138", "166", "195", "212", "—"], LIGHT_BG, GRAPHITE),
    (["EBITDA Ajustado (R$M)", "52", "72", "105", "151", "187", "—"], LIGHT_BG, GRAPHITE),
    (["Divida Liq./EBITDA", "2,2x", "1,9x", "1,6x", "1,3x", "1,1x", "max 3,5x"], RGBColor(0xE8,0xF0,0xFF), NAVY),
    (["Divida Bruta/EBITDA", "3,0x", "2,6x", "2,2x", "1,8x", "1,7x", "max 4,0x"], WHITE, GRAPHITE),
    (["EBITDA/Juros (ICR)", "2,1x", "2,4x", "2,8x", "3,2x", "3,8x", "min 2,0x"], WHITE, POSITIVE),
    (["DSCR (EBITDA/Servico Div.)", "1,4x", "1,6x", "1,9x", "2,4x", "2,8x", "min 1,5x"], RGBColor(0xE8,0xF0,0xFF), NAVY),
    (["FCO/EBITDA (conversao)", "68%", "71%", "74%", "77%", "79%", "min 60%"], WHITE, POSITIVE),
    (["Capex/Receita", "6,8%", "6,2%", "5,9%", "5,4%", "8,1%*", "—"], WHITE, GRAPHITE),
]

for i, (data, bg, tc) in enumerate(credit_rows):
    bd = data[0] in ["Divida Liq./EBITDA", "DSCR (EBITDA/Servico Div.)"]
    table_row(slide, data, l=0.35, t=1.58+i*0.42, col_w=cw2, bg=bg, text_color=tc, bold=bd)

add_text(slide, "* Capex 2025 inclui R$95M da expansao da planta Jatai (CRA) | Fontes: DFs auditadas Meridian | Calculos: IB Agents Capital",
         l=0.35, t=6.95, w=12.5, h=0.18, size=7, color=MUTED, font="Calibri", italic=True)


# ============================================================================
# SLIDE 10 — STRESS TEST
# ============================================================================
slide = slide_base(
    "Stress test: covenants mantidos mesmo em queda de 30% do EBITDA — margem de seguranca estrutural",
    "Modulo 6 — Analise de Cenarios | Stress Test 2026-2030"
)

add_text(slide, "PREMISSAS DOS CENARIOS", l=0.35, t=1.25, w=12.5, h=0.22,
         size=9, bold=True, color=NAVY, font="Calibri")
premissas = [
    ("BULL (+15%)", "Expansao de volumes acelerada, BRL/USD acima de R$5,90, premium mix 48%", POSITIVE),
    ("BASE (0%)", "Crescimento organico conforme historico — projecoes base da administracao", NAVY),
    ("BEAR (-15%)", "Retorno de grippe aviaria em mercados-chave, compressao de margens, cambio R$5,40", GOLD),
    ("STRESSED (-30%)", "Combinacao de choques: surto sanitario + credito restrito + cambio valorizado", NEGATIVE),
]
for i, (label, desc, color) in enumerate(premissas):
    add_rect(slide, 0.35, 1.52+i*0.3, 1.5, 0.26, fill=color)
    add_text(slide, label, l=0.35, t=1.52+i*0.3, w=1.5, h=0.26,
             size=8, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, desc, l=1.9, t=1.54+i*0.3, w=11.0, h=0.24,
             size=8.5, color=GRAPHITE, font="Calibri")

# Tabela de resultados
cols_st = ["Cenario", "EBITDA 2026", "Div.Liq/EBITDA", "ICR", "DSCR", "Covenant OK?"]
cw_st = 12.6 / len(cols_st)
table_header(slide, cols_st, l=0.35, t=2.82, col_w=cw_st)

st_rows = [
    (["Bull (+15%)", "R$215M", "0,9x", "4,4x", "3,2x", "SIM"], WHITE, POSITIVE),
    (["Base (0%)", "R$187M", "1,1x", "3,8x", "2,8x", "SIM"], RGBColor(0xE8,0xF0,0xFF), NAVY),
    (["Bear (-15%)", "R$159M", "1,3x", "3,2x", "2,4x", "SIM"], WHITE, GOLD),
    (["Stressed (-30%)", "R$131M", "1,6x", "2,7x", "2,0x", "SIM *"], LIGHT_BG, NEGATIVE),
]
for i, (data, bg, tc) in enumerate(st_rows):
    bd = data[0] == "Base (0%)"
    table_row(slide, data, l=0.35, t=3.1+i*0.32, col_w=cw_st, bg=bg, text_color=tc, bold=bd)

add_text(slide, "* No cenario stressed, DSCR de 2,0x esta exatamente no limite do covenant minimo (1,5x). Breach ocorre apenas com queda de EBITDA > 46% — nao projetado em nenhum cenario historico.",
         l=0.35, t=4.4, w=12.5, h=0.3, size=8.5, color=NEGATIVE, font="Calibri", bold=True)

add_text(slide, "THRESHOLD DE BREACH (EBITDA minimo para manutencao dos covenants)", l=0.35, t=4.8, w=12.5, h=0.22,
         size=9, bold=True, color=NAVY, font="Calibri")
breach_data = [
    ("Divida Liq./EBITDA max 3,5x", "EBITDA min: R$61M", "Margem: -67% vs EBITDA 2025 — altamente improvavel"),
    ("ICR min 2,0x", "EBITDA min: R$82M", "Margem: -56% vs EBITDA 2025 — nao projetado"),
    ("DSCR min 1,5x", "EBITDA min: R$100M", "Margem: -46% vs EBITDA 2025 — cenario extremo"),
]
for i, (cov, threshold, comment) in enumerate(breach_data):
    add_rect(slide, 0.35, 5.08+i*0.38, 3.5, 0.33, fill=LIGHT_BG)
    add_text(slide, cov, l=0.45, t=5.1+i*0.38, w=3.3, h=0.28, size=8.5, bold=True, color=NAVY, font="Calibri")
    add_rect(slide, 3.9, 5.08+i*0.38, 2.2, 0.33, fill=RGBColor(0xFF,0xED,0xED))
    add_text(slide, threshold, l=4.0, t=5.1+i*0.38, w=2.1, h=0.28, size=8.5, bold=True, color=NEGATIVE, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, comment, l=6.2, t=5.12+i*0.38, w=6.9, h=0.26, size=8.5, color=POSITIVE, font="Calibri")

add_text(slide, "Metodologia: choque de EBITDA aplicado sobre base 2025 com estrutura de divida e amortizacao fixas conforme term sheet proposto",
         l=0.35, t=6.95, w=12.5, h=0.18, size=7, color=MUTED, font="Calibri", italic=True)


# ============================================================================
# SLIDE 11 — ESTRUTURA DE GARANTIAS
# ============================================================================
slide = slide_base(
    "Garantias de R$225M (1,5x de cobertura) — recebiveis agroindustriais + imoveis rurais + aval",
    "Modulo 8 — Estrutura da Operacao | Garantias"
)

garantias = [
    ("RECEBIVEIS AGROINDUSTRIAIS", "R$150M", "PRIMARIA",
     ["Cessao fiduciaria de recebiveis de 340 clientes", "Concentracao maxima: 8% por cliente (JBS, Carrefour, GPA)", "Taxa de inadimplencia historica: 0,3% a.a.", "Advance rate: 78% sobre carteira bruta de R$192M", "Monitoramento mensal — agente fiduciario: Oliveira Trust"],
     POSITIVE),
    ("IMOVEIS RURAIS (MATRICULAS)", "R$62M", "COMPLEMENTAR",
     ["4 propriedades rurais em Jatai/Mineiros/Rio Verde (GO)", "Area total: 3.840 hectares + benfeitorias industriais", "Laudo ENGEFOTO (01/2026): valor de mercado R$103M", "LTV: 60% = R$62M | Haircut padrao imovel rural", "Alienacao fiduciaria registrada em cartorio"],
     NAVY),
    ("AVAL DOS SOCIOS CONTROLADORES", "R$30M", "TERCIARIA",
     ["Aval solidario: Familia Meireles (71,2% do capital)", "Patrimonio declarado comprovado: R$420M", "Ratio patrimonio/aval: 14x", "Restricao de alienacao de ativos > R$5M sem ciencia do agente fiduciario"],
     GOLD),
]

for i, (title, value, rank, items, color) in enumerate(garantias):
    l = 0.3 + i * 4.3
    add_rect(slide, l, 1.25, 4.1, 5.55, fill=LIGHT_BG)
    add_rect(slide, l, 1.25, 4.1, 0.55, fill=color)
    add_text(slide, rank, l=l+0.1, t=1.27, w=1.2, h=0.22,
             size=7, bold=True, color=WHITE, font="Calibri")
    add_text(slide, title, l=l+0.1, t=1.48, w=3.9, h=0.22,
             size=8, bold=True, color=WHITE, font="Calibri")
    add_text(slide, value, l=l+2.4, t=1.27, w=1.6, h=0.22,
             size=14, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.RIGHT)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", l=l+0.1, t=1.88+j*0.53, w=3.9, h=0.5,
                 size=8.5, color=GRAPHITE, font="Calibri")

# Resumo de cobertura
add_rect(slide, 0.3, 6.85, 12.7, 0.25, fill=NAVY)
add_text(slide, "COBERTURA TOTAL: R$242M para emissao de R$150M = 1,61x  |  Minimo exigido: 1,2x  |  Excesso: R$92M",
         l=0.4, t=6.86, w=12.5, h=0.22,
         size=9, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 12 — EQUIPE DE GESTAO
# ============================================================================
slide = slide_base(
    "Equipe de gestao com 180+ anos de experiencia combinada — track record de entrega em todos os ciclos",
    "Modulo 7 — Equipe de Gestao"
)

executivos = [
    ("RM", "Roberto Meireles", "CEO & Fundador", NAVY,
     ["Fundou a Meridian em 1998 aos 32 anos", "Engenheiro Agronomo — ESALQ/USP", "MBA — Wharton School (2004)", "Presidente da ABPA Regional CO (2019-2023)", "Track record: 5x crescimento da receita em 10 anos"]),
    ("CM", "Claudia Meireles", "CFO", NAVY_MED,
     ["15 anos na Meridian (CFO desde 2012)", "Economista FGV-SP + CPA-10/CPA-20 ANBIMA", "Liderou 3 emissoes de debentures (2016/2019/2022)", "Implantou ERP SAP S/4HANA em 2021", "Supervisiona equipe de 18 profissionais de financas"]),
    ("LA", "Luiz Andrade", "COO", GRAPHITE,
     ["Ex-BRF (8 anos) — Diretor Industrial Frangosul", "Engenheiro de Alimentos — UFSC", "Responsavel pela expansao de Rio Verde (2014)", "Implantou WCM (World Class Manufacturing) em 2020", "Reducao de 23% no custo de producao em 4 anos"]),
    ("PA", "Patricia Alves", "Diretora Comercial", GRAPHITE,
     ["Ex-JBS Exports (6 anos) — VP Asia Pacific", "Fluente em Ingles, Mandarin e Arabe", "Abertura de 12 novos mercados desde 2019", "Exportacoes cresceram de 8% para 28% da receita", "Contatos com 45 importadores habilitados MAPA"]),
]

for i, (initials, name, role, color, items) in enumerate(executivos):
    col = i % 2
    row = i // 2
    l = 0.35 + col * 6.45
    t = 1.3 + row * 2.7

    add_rect(slide, l, t, 0.6, 0.6, fill=color)
    add_text(slide, initials, l=l, t=t, w=0.6, h=0.6,
             size=16, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, name, l=l+0.7, t=t+0.04, w=5.4, h=0.28,
             size=12, bold=True, color=NAVY, font="Calibri")
    add_text(slide, role, l=l+0.7, t=t+0.34, w=5.4, h=0.22,
             size=9, color=MUTED, font="Calibri")
    for j, item in enumerate(items):
        add_text(slide, f"  • {item}", l=l+0.1, t=t+0.7+j*0.38, w=6.1, h=0.35,
                 size=8.5, color=GRAPHITE, font="Calibri")


# ============================================================================
# SLIDE 13 — TESE DE INVESTIMENTO
# ============================================================================
slide = slide_base(
    "Tese: retorno IPCA+7,25% com risco AA- em empresa com DSCR de 2,8x e garantias de 1,5x de cobertura",
    "Modulo 8 — Tese de Investimento e Proximos Passos"
)

add_text(slide, "POR QUE INVESTIR AGORA", l=0.4, t=1.25, w=6.5, h=0.22,
         size=10, bold=True, color=NAVY, font="Calibri")

razoes = [
    ("RETORNO SUPERIOR AO MERCADO",
     "IPCA+7,25% vs IPCA+6,10% media ANBIMA para AA- em jan/26 — premio de 115bps por liquidez"),
    ("EMPRESA EM MOMENTO DE ACELERACAO",
     "CAPEX de R$150M financiado pelo CRA gera +35k aves/dia e ROI de 28% — EBITDA deve superar R$250M em 2027"),
    ("ESTRUTURA DE GARANTIAS ROBUSTA",
     "R$242M de garantias para R$150M emitidos — recebiveis com inadimplencia de 0,3% historica"),
    ("SETOR BENEFICIADO POR MACRO",
     "BRL/USD > R$5,80 + exportacoes recorde + premiumizacao = ventos favoraveis para 3-5 anos"),
    ("RATING E COVENANTS CONSERVADORES",
     "AA- Fitch com outlook estavel — 4 covenants com folgas de 46-67% no cenario stressed"),
]

for i, (title, desc) in enumerate(razoes):
    add_rect(slide, 0.4, 1.52+i*0.9, 0.5, 0.5, fill=GOLD)
    add_text(slide, str(i+1), l=0.4, t=1.52+i*0.9, w=0.5, h=0.5,
             size=18, bold=True, color=NAVY, font="Georgia", align=PP_ALIGN.CENTER)
    add_text(slide, title, l=1.0, t=1.55+i*0.9, w=5.5, h=0.25,
             size=9, bold=True, color=NAVY, font="Calibri")
    add_text(slide, desc, l=1.0, t=1.82+i*0.9, w=5.5, h=0.32,
             size=8.5, color=GRAPHITE, font="Calibri")

add_rect(slide, 7.0, 1.1, 0.02, 5.8, fill=BORDER)

add_text(slide, "CATALISADORES", l=7.3, t=1.25, w=5.8, h=0.22,
         size=10, bold=True, color=NAVY, font="Calibri")
cats = [
    ("Mar/26", "Fechamento do CRA — emissao e liquidacao financeira"),
    ("Abr-Jun/26", "Inicio das obras de expansao da planta Jatai (+35k aves/dia)"),
    ("Set/26", "Revisao de rating Fitch — expectativa de confirmacao AA-"),
    ("Dez/26", "Operacionalizacao da nova linha de congelados premium"),
    ("2027", "Meta: Receita R$1,55Bi | EBITDA R$250M | Exportacoes 32%"),
    ("2028-30", "Pipeline de M&A: 2 alvos identificados em GO e MT"),
]
for i, (period, event) in enumerate(cats):
    add_rect(slide, 7.3, 1.52+i*0.72, 1.1, 0.62, fill=NAVY)
    add_text(slide, period, l=7.3, t=1.52+i*0.72, w=1.1, h=0.62,
             size=8, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, event, l=8.5, t=1.55+i*0.72, w=4.6, h=0.55,
             size=9, color=GRAPHITE, font="Calibri")


# ============================================================================
# SLIDE 14 — TERM SHEET RESUMIDO
# ============================================================================
slide = slide_base(
    "Estrutura da Emissao — CRA Grupo Meridian Alimentos S.A. | R$150M | IPCA+7,25% | 5 anos",
    "Modulo 8 — Term Sheet Resumido"
)

ts_items = [
    ("Emissora / Cedente", "Grupo Meridian Alimentos S.A. | CNPJ 12.345.678/0001-90"),
    ("Securitizadora", "Eco Securitizacao S.A. | CVM Reg. 21.234"),
    ("Instrumento", "Certificados de Recebiveis do Agronegocio (CRA) — Lei 11.076/04 e Lei 13.986/20"),
    ("Volume Total", "R$ 150.000.000 (Cento e cinquenta milhoes de reais)"),
    ("Serie", "Serie Unica"),
    ("Remuneracao", "IPCA (IBGE) + 7,25% a.a. — base 252 DU"),
    ("Prazo", "5 (cinco) anos | Vencimento: marco/2031"),
    ("Amortizacao", "Bullet (principal no vencimento) | Juros semestrais"),
    ("Rating", "AA- (Fitch Ratings) | Perspectiva: Estavel"),
    ("Regime de Oferta", "RCVM 160/2022 art. 3o, I — Esforcos Restritos (476)"),
    ("Publico Alvo", "Investidores Profissionais (RCVM 30/2021)"),
    ("Lastro", "Direitos creditarios agroindustriais — cessao fiduciaria"),
    ("Garantias", "Cessao fiduc. de recebiveis (R$150M) + AFF imoveis rurais (R$62M) + Aval (R$30M)"),
    ("Agente Fiduciario", "Oliveira Trust Distribuidora de Titulos e Valores Mobiliarios S.A."),
    ("Coordenador Lider", "IB Agents Capital S.A. CTVM"),
    ("Destinacao", "95% expansao industrial planta Jatai (GO) + 5% capital de giro"),
    ("Comissionamento", "2,0% flat sobre volume total emitido"),
    ("Data Base", "Marco/2026"),
]

col_w_ts = [3.2, 9.2]
for i, (label, value) in enumerate(ts_items):
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, 0.35, 1.3+i*0.32, col_w_ts[0], 0.3, fill=NAVY if i == 0 else bg)
    add_text(slide, label,
             l=0.45, t=1.32+i*0.32, w=col_w_ts[0]-0.1, h=0.26,
             size=8, bold=True, color=WHITE if i == 0 else NAVY, font="Calibri")
    add_rect(slide, 0.35+col_w_ts[0], 1.3+i*0.32, col_w_ts[1], 0.3, fill=NAVY if i == 0 else bg)
    add_text(slide, value,
             l=0.45+col_w_ts[0], t=1.32+i*0.32, w=col_w_ts[1]-0.1, h=0.26,
             size=8, bold=(i == 0), color=WHITE if i == 0 else GRAPHITE, font="Calibri")


# ============================================================================
# SLIDE 15 — FATORES DE RISCO
# ============================================================================
slide = slide_base(
    "Quatro riscos principais — todos mitigados por protecoes estruturais ja em vigor (RCVM 160)",
    "Modulo 8 — Fatores de Risco | Ordenados por Materialidade"
)

add_text(slide, "Escala qualitativa obrigatoria RCVM 160: Maior / Medio / Menor | Ordenado do maior para o menor impacto potencial",
         l=0.4, t=1.25, w=12.5, h=0.2, size=8, color=MUTED, font="Calibri", italic=True)

riscos = [
    ("MAIOR", "Risco Sanitario (Grippe Aviaria)",
     "Surto de HPAI pode fechar mercados de exportacao e gerar queda de receita",
     ["Diversificacao: 18 paises (nenhum > 22% das exportacoes)",
      "Seguro agropecuario cobrindo 80% do plantel",
      "Protocolo de biosseguridade Nivel 3 — auditado MAPA anualmente",
      "Covenant: trigger de cash sweep se exportacoes caem > 30%"],
     NEGATIVE),
    ("MEDIO", "Risco de Mercado (Commodities)",
     "Alta do milho/soja eleva custo de racao e comprime margens",
     ["62% da racao produzida internamente — hedge parcial de custo",
      "Historico: margem EBITDA nunca caiu abaixo de 7,5% em 5 anos",
      "Politica de hedge: compra antecipada de 40% da necessidade anual",
      "Covenant de margem: alert se EBITDA/Receita < 9% por 2 trimestres"],
     GOLD),
    ("MEDIO", "Risco de Taxa de Juros (IPCA)",
     "IPCA elevado aumenta o custo de servico da divida em termos nominais",
     ["CRA IPCA+7,25%: receita Meridian tambem indexada ao IPCA parcialmente",
      "Exportacoes em USD funcionam como hedge natural contra IPCA",
      "Cenario stressed: IPCA a 8%aa ainda mantem DSCR acima de 2,0x"],
     GOLD),
    ("MENOR", "Risco de Liquidez (Mercado Secundario)",
     "CRA de emissao privada (476) tem liquidez secundaria limitada",
     ["Estrutura bullet reduz necessidade de liquidez ao longo do prazo",
      "Oferecido a 35+ investidores institucionais com horizonte LP",
      "Agente de liquidez: IB Agents Capital disponibiliza cotacao semanal"],
     MUTED),
]

for i, (nivel, titulo, descricao, mitigantes, color) in enumerate(riscos):
    t = 1.55 + i * 1.28
    add_rect(slide, 0.35, t, 0.8, 1.2, fill=color)
    add_text(slide, nivel, l=0.35, t=t+0.45, w=0.8, h=0.3,
             size=7, bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER)
    add_text(slide, titulo, l=1.25, t=t+0.03, w=5.0, h=0.28,
             size=10, bold=True, color=NAVY, font="Calibri")
    add_text(slide, descricao, l=1.25, t=t+0.33, w=5.0, h=0.35,
             size=8.5, color=GRAPHITE, font="Calibri")
    add_text(slide, "Mitigantes:", l=6.5, t=t+0.03, w=1.2, h=0.22,
             size=8, bold=True, color=NAVY, font="Calibri")
    for j, mit in enumerate(mitigantes):
        add_text(slide, f"• {mit}", l=6.5, t=t+0.28+j*0.23, w=6.5, h=0.22,
                 size=8, color=GRAPHITE, font="Calibri")

add_rect(slide, 0, 6.88, 13.33, 0.27, fill=RGBColor(0xFF,0xED,0xED))
add_text(slide, "AVISO LEGAL: Este material e baseado em informacoes publicas e fornecidas pela Meridian. Nao constitui oferta de valores mobiliarios. Investidores devem conduzir sua propria due diligence. Rentabilidade passada nao garante resultados futuros.",
         l=0.3, t=6.9, w=12.7, h=0.22, size=7, color=NEGATIVE, font="Calibri", italic=True)


# ============================================================================
# SALVAR
# ============================================================================
import os
out_dir = r"C:\Users\dloge\OneDrive\Área de Trabalho\Teste 1\ib-agents\frontend\public\knowledge"
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "book_credito_meridian_alimentos_cra.pptx")
prs.save(out_path)
print(f"PPTX salvo em: {out_path}")
print(f"Slides gerados: {len(prs.slides)}")
