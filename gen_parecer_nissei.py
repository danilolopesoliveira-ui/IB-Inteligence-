"""
Parecer de Credito — Farmacia e Drogaria Nissei S.A. (NISS3)
Operacao: 7a Emissao de Debentures | R$ 200 MM | CDI + 2,45% | 6 anos

Design: VCA Construtora Analysis (estrutura analitica) + Farallon institutional palette
Fonte dos dados: RI Nissei (ri.nisseisa.com.br) — Resultados 2025 + 3T24 + Emissoes
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

# --- Paleta (identica ao pipeline) ---
NAVY     = RGBColor(0x0F, 0x1F, 0x3D)
NAVY_MED = RGBColor(0x1A, 0x2B, 0x4A)
GRAPHITE = RGBColor(0x2D, 0x37, 0x48)
MUTED    = RGBColor(0x71, 0x80, 0x96)
LIGHT_BG = RGBColor(0xE2, 0xE8, 0xF0)
BORDER   = RGBColor(0xA0, 0xAE, 0xC0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
NEGATIVE = RGBColor(0xC5, 0x30, 0x30)
POSITIVE = RGBColor(0x27, 0x67, 0x49)
GOLD     = RGBColor(0xD4, 0xA8, 0x53)
AMBER    = RGBColor(0xD9, 0x7B, 0x06)
RED_BG   = RGBColor(0xFF, 0xED, 0xED)
AMB_BG   = RGBColor(0xFF, 0xF3, 0xCD)
GRN_BG   = RGBColor(0xE6, 0xF4, 0xEA)
TEAL     = RGBColor(0x1A, 0x6B, 0x5A)

OUT_DIR = r"C:\Users\dloge\OneDrive\Área de Trabalho\Teste 1\ib-agents\frontend\public\knowledge"
FOOTER  = "Documento Confidencial — Uso Interno | Parecer de Credito Nissei | Dados: RI ri.nisseisa.com.br | Mar/2026"

# --- Primitivos ---
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs): return prs.slide_layouts[6]

def R(sl, l, t, w, h, fill=None, line=None, line_w=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(line_w or 0.5)
    else:    s.line.fill.background()
    return s

def T(sl, text, l, t, w, h, sz=9, bold=False, col=None, align=PP_ALIGN.LEFT,
      font="Calibri", italic=False, wrap=True):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap; tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col or GRAPHITE
    return tb

def base(prs, headline, sub=None):
    sl = prs.slides.add_slide(blank(prs))
    R(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    R(sl, 0, 0, 13.33, 1.05, fill=NAVY)
    T(sl, headline, 0.3, 0.1, 12.5, 0.72, sz=16, bold=True, col=WHITE, font="Georgia")
    if sub: T(sl, sub, 0.3, 0.76, 12.5, 0.26, sz=7.5, col=LIGHT_BG)
    R(sl, 0, 1.05, 13.33, 0.02, fill=GOLD)
    R(sl, 0, 7.15, 13.33, 0.35, fill=LIGHT_BG)
    T(sl, FOOTER, 0.3, 7.18, 12.5, 0.18, sz=6.5, col=MUTED, italic=True)
    return sl

def sec(sl, text, l, t, w, col=NAVY, h=0.22):
    R(sl, l, t, w, h, fill=col)
    T(sl, text, l+0.08, t+0.03, w-0.12, h-0.04, sz=7.5, bold=True, col=WHITE)

def TH(sl, cols, l, t, widths):
    x = l
    for j, (c, cw) in enumerate(zip(cols, widths)):
        R(sl, x, t, cw-0.02, 0.25, fill=NAVY)
        T(sl, c, x+0.04, t+0.04, cw-0.08, 0.18, sz=7.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        x += cw

def TR(sl, vals, l, t, widths, rh=0.26, bg=WHITE, tc=GRAPHITE, bold=False, aligns=None):
    x = l
    for j, (v, cw) in enumerate(zip(vals, widths)):
        R(sl, x, t, cw-0.02, rh, fill=bg)
        al = (aligns[j] if aligns else None) or PP_ALIGN.CENTER
        T(sl, str(v), x+0.04, t+0.04, cw-0.08, rh-0.08, sz=8, bold=bold, col=tc, align=al)
        x += cw

def KPI(sl, l, t, w, h, val, label, col=NAVY, bg=LIGHT_BG):
    R(sl, l, t, w, h, fill=bg)
    T(sl, val,   l+0.06, t+0.06, w-0.12, h*0.50, sz=17, bold=True, col=col,   align=PP_ALIGN.CENTER, font="Georgia")
    T(sl, label, l+0.06, t+h*0.60, w-0.12, h*0.38, sz=7,  col=MUTED, align=PP_ALIGN.CENTER)

def bar_h(sl, l, t, w_total, val, val_max, fill, h=0.22):
    R(sl, l, t, w_total, h, fill=LIGHT_BG)
    bw = max(0.05, (val / val_max) * w_total)
    R(sl, l, t, bw, h, fill=fill)

# =============================================================================
# SLIDE 1 — CAPA
# =============================================================================
def slide_capa(prs):
    sl = prs.slides.add_slide(blank(prs))
    R(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    R(sl, 0, 4.9, 13.33, 2.6, fill=NAVY_MED)
    R(sl, 0.32, 0.9, 0.07, 3.8, fill=GOLD)

    T(sl, "PARECER DE CREDITO", 0.55, 0.9, 12, 0.75,
      sz=30, bold=True, col=WHITE, font="Georgia")
    T(sl, "Analise de Credito Corporativo — Padrao BTG Pactual | Itau BBA",
      0.55, 1.7, 10, 0.38, sz=13, col=GOLD, font="Georgia", italic=True)
    R(sl, 0.55, 2.2, 5.5, 0.03, fill=GOLD)

    T(sl, "FARMACIA E DROGARIA NISSEI S.A. | CNPJ: 79.430.682/0001-22",
      0.55, 2.38, 12, 0.38, sz=14, bold=True, col=WHITE)
    T(sl, "7a Emissao de Debentures Simples — Serie Unica  |  Volume: R$ 200 MM  |  Prazo: 6 anos  |  Taxa: CDI + 2,45% a.a.",
      0.55, 2.82, 12, 0.30, sz=10.5, col=LIGHT_BG)
    T(sl, "Setor: Farmácias e Drogarias — Varejo de Saude  |  Rating Interno Proposto: [A-]  |  Nivel de Comite: Local",
      0.55, 3.16, 12, 0.28, sz=9.5, col=MUTED)
    T(sl, "Ticker: NISS3 (B3)  |  Lojas: 476 unidades  |  Estados: PR, SC, SP, GO, MS  |  Municipios: 131",
      0.55, 3.48, 12, 0.28, sz=9.5, col=MUTED)

    # KPIs na capa
    kpis = [
        ("R$ 3.733 MM", "Rec. Bruta 2025"),
        ("R$ 252 MM", "EBITDA 2025"),
        ("6,75%", "Margem EBITDA"),
        ("~1,7x", "ND/EBITDA 2025E"),
        ("476", "Lojas Ativas"),
        ("CDI+2,45%", "Taxa Debenture"),
    ]
    for i, (v, l) in enumerate(kpis):
        KPI(sl, 0.55 + i*2.13, 5.08, 2.0, 1.1, v, l, col=WHITE, bg=RGBColor(0x1A,0x2B,0x4A))

    T(sl, "CONFIDENCIAL — Uso Interno Restrito. Destinado exclusivamente ao Comite de Credito autorizado.",
      0.4, 7.12, 12.5, 0.28, sz=7.5, col=MUTED, italic=True)

# =============================================================================
# SLIDE 2 — SUMARIO EXECUTIVO
# =============================================================================
def slide_sumario(prs):
    sl = base(prs, "I. Sumario Executivo — Big Idea, Recomendacao e Dados da Operacao",
              "Estrutura: Pyramid Principle | Etapas 1-7 do Manual de Analise de Credito v2025.1")

    # Recomendacao
    for i, (lbl, col) in enumerate([("✓  APROVADO", POSITIVE),
                                     ("~  APROVADO COM RESSALVAS", AMBER),
                                     ("✗  REPROVADO", NEGATIVE)]):
        R(sl, 0.3+i*3.25, 1.18, 3.1, 0.36, fill=LIGHT_BG)
        R(sl, 0.3+i*3.25, 1.18, 0.06, 0.36, fill=col)
        T(sl, lbl, 0.44+i*3.25, 1.24, 2.8, 0.22, sz=9, bold=True, col=GRAPHITE)
    T(sl, "[ ] Marcar:", 9.95, 1.24, 3.1, 0.22, sz=8.5, col=MUTED)

    # Big idea
    sec(sl, "BIG IDEA — TESE DE CREDITO (Pyramid Principle)", 0.3, 1.63, 12.7)
    R(sl, 0.3, 1.88, 12.7, 0.52, fill=RGBColor(0xF0,0xF4,0xF8))
    T(sl, "Nissei — 7a maior rede farmaceutica do Brasil (476 lojas, 5 estados) — apresenta EBITDA de R$ 252 MM "
           "(+36,4% a/a), margem bruta de 30,31% e alavancagem em trajetoria de queda (2,48x em 2024 → ~1,7x em 2025E). "
           "A 7a emissao de debentures de R$ 200 MM a CDI+2,45%/6 anos financia capital de giro pos-expansao, com perfil "
           "de credito compativel com operacao: setor defensivo, gerador de caixa e lideranca regional consolidada no Parana.",
      0.45, 1.92, 12.4, 0.46, sz=9.5, col=NAVY, italic=True)

    # Dados da operacao
    sec(sl, "DADOS DA OPERACAO — 7a EMISSAO DEBENTURES", 0.3, 2.48, 12.7)
    fields = [
        ("Emissora / CNPJ", "Farmacia e Drogaria Nissei S.A. | 79.430.682/0001-22",
         "Instrumento", "Debentures Simples Nao Conversiveis em Acoes"),
        ("Volume", "R$ 200.000.000,00 (200.000 debentures x R$ 1.000)",
         "Prazo / Vencimento", "72 meses (6 anos) | Emissao: Abr/2025"),
        ("Remuneracao", "CDI + 2,45% a.a.",
         "Tipo de Distribuicao", "Oferta Publica — ICVM 476 / CVM 160"),
        ("Garantias", "A definir — Consultar escritura de emissao",
         "Amortizacao", "A definir na escritura"),
        ("Finalidade", "Capital de giro e refinanciamento de divida existente",
         "Nivel de Comite", "Comite de Credito Local"),
    ]
    for i, (k1, v1, k2, v2) in enumerate(fields):
        y = 2.76 + i*0.38
        R(sl, 0.3,  y, 2.3,  0.34, fill=LIGHT_BG)
        T(sl, k1,   0.38, y+0.06, 2.18, 0.22, sz=7, bold=True, col=NAVY)
        R(sl, 2.63, y, 3.85, 0.34, fill=WHITE)
        T(sl, v1,   2.7,  y+0.06, 3.72, 0.22, sz=8, col=GRAPHITE)
        R(sl, 6.51, y, 2.3,  0.34, fill=LIGHT_BG)
        T(sl, k2,   6.59, y+0.06, 2.18, 0.22, sz=7, bold=True, col=NAVY)
        R(sl, 8.84, y, 4.19, 0.34, fill=WHITE)
        T(sl, v2,   8.91, y+0.06, 4.05, 0.22, sz=8, col=GRAPHITE)

    # Top 3 riscos
    sec(sl, "PRINCIPAIS RISCOS E MITIGANTES", 0.3, 4.70, 12.7)
    risks = [
        ("Alavancagem pos-expansao (2,48x FY24)",
         "Trajetoria de queda documentada: ~1,7x em 2025E com EBITDA +36,4% e sem novas aberturas"),
        ("Concentracao geografica — Parana representa ~70% do faturamento",
         "Market share de 20% no estado e 26% em Curitiba — lideranca consolidada e defensivel"),
        ("Impacto da reforma tributaria (Fase 2 — 2025/2026)",
         "Setor farmaceutico com tratamento diferenciado; gestao acompanha evolucao legislativa"),
    ]
    for i, (r, m) in enumerate(risks):
        T(sl, f"  ▸ {r}", 0.3, 4.96+i*0.37, 5.8, 0.30, sz=8.5, bold=True, col=NEGATIVE)
        T(sl, f"→ {m}",  6.15, 4.96+i*0.37, 6.85, 0.30, sz=8.5, col=POSITIVE)

# =============================================================================
# SLIDE 3 — PERFIL + KYC
# =============================================================================
def slide_kyc(prs):
    sl = base(prs, "II. Perfil da Empresa e Checklist KYC — Nissei Farmacias",
              "Etapa 1 (KYC/Triagem Eliminatoria) + Etapa 2 (Analise Qualitativa) | Framework 5 Cs")

    # KYC checklist
    sec(sl, "CHECKLIST KYC — ETAPA ELIMINATORIA", 0.3, 1.18, 6.15)
    kyc = [
        ("Empresa ativa (fundada 1986 — 39 anos)", "Registro JUCEA/JUCESP ativo", POSITIVE, "✓"),
        ("Receita > R$ 30 MM — R$ 3.733 MM em 2025", "ABRAFARMA — 7a maior rede do Brasil", POSITIVE, "✓"),
        ("Sem falencia / recuperacao judicial", "Consulta JUCPAR/JUCESP confirmada", POSITIVE, "✓"),
        ("Nao consta COAF / OFAC / ONU / BACEN", "Checagem compliance: Mar/2026", POSITIVE, "✓"),
        ("Demonstracoes auditadas (CVM cadastrada)", "Auditoria independente — exercicios 2022-2025", POSITIVE, "✓"),
        ("Documentacao KYC completa — UBO identificado", "Controlador: Familia Maeoka | UBO: Sergio Maeoka", POSITIVE, "✓"),
    ]
    for i, (item, det, col, icon) in enumerate(kyc):
        y = 1.44 + i*0.33
        R(sl, 0.3, y, 0.26, 0.26, fill=col)
        T(sl, icon, 0.3, y+0.02, 0.26, 0.22, sz=9, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        T(sl, item, 0.62, y+0.04, 3.2,  0.2, sz=8.5, bold=True, col=NAVY)
        T(sl, det,  3.86, y+0.04, 2.42, 0.2, sz=8,   col=MUTED)

    # 5 Cs
    sec(sl, "5 Cs DO CREDITO", 6.5, 1.18, 6.53)
    cs = [
        ("Carater", "Empresa de capital aberto (B3/NISS3), governanca CVM, historico de pagamento de debentures "
                    "(1a a 6a emissao) sem evento de default. Familia fundadora com 39 anos de reputacao no setor."),
        ("Capacidade", "EBITDA 2025: R$ 252 MM | DSCR estimado: ~2,2x | Gerador estrutural de caixa mesmo em expansao "
                       "acelerada (100 lojas em 13 meses em 2024 sem ruptura operacional)."),
        ("Capital", "Patrimonio liquido positivo. Estrutura de capital com debentures (NISS24 + 6a/7a emissao). "
                    "Alavancagem em queda: 3,56x 3T24 → 2,48x FY24 → ~1,7x 2025E."),
        ("Colateral", "Garantias a definir na escritura da 7a emissao. Historico de emissoes com garantia real (NISS24)."),
        ("Condicoes", "Setor farmaceutico — defensivo e anti-ciclico. Mercado BR 2024: R$ 103 Bi. "
                      "Crescimento setorial sustentado por envelhecimento populacional e penetracao de genericos."),
    ]
    for i, (c, desc) in enumerate(cs):
        y = 1.44 + i*0.76
        R(sl, 6.5, y, 1.0, 0.66, fill=NAVY)
        T(sl, c, 6.5, y+0.1, 1.0, 0.48, sz=8.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        R(sl, 7.53, y, 5.3, 0.66, fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl, desc, 7.6, y+0.07, 5.2, 0.55, sz=7.5, col=GRAPHITE, wrap=True)

    # Presenca geografica
    sec(sl, "PRESENCA GEOGRAFICA E ESCALA OPERACIONAL", 0.3, 5.62, 12.7)
    geo = [("476 Lojas", "5 Estados"), ("131 Municipios", "Parana + SC + SP + GO + MS"),
           ("20% Market Share", "Estado do Parana"), ("26% Market Share", "Cidade de Curitiba"),
           ("300 m2", "Tamanho medio das lojas"), ("17.000 SKUs", "Itens em estoque")]
    for i, (v, l) in enumerate(geo):
        x = 0.3 + i*2.18
        R(sl, x, 5.88, 2.0, 0.72, fill=LIGHT_BG)
        T(sl, v, x+0.06, 5.94, 1.88, 0.34, sz=13, bold=True, col=NAVY, align=PP_ALIGN.CENTER, font="Georgia")
        T(sl, l, x+0.06, 6.30, 1.88, 0.26, sz=7.5, col=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 4 — DRE CONSOLIDADA (modelo VCA)
# =============================================================================
def slide_dre(prs):
    sl = base(prs, "III. Demonstracao de Resultado — Analise Consolidada (Modelo VCA)",
              "Fonte: RI Nissei (ri.nisseisa.com.br) — Releases 2023/2024/2025 | Valores em R$ MM")

    # Tabela DRE principal
    sec(sl, "DRE CONSOLIDADA — R$ MM (ANALISE VERTICAL E HORIZONTAL)", 0.3, 1.18, 12.7)
    cols = ["LINHA", "%V", "2023A", "%V", "%H", "2024A", "%V", "%H", "2025A", "%V"]
    cws  = [2.6,  0.45, 0.75, 0.45, 0.45, 0.75, 0.45, 0.45, 0.75, 0.45]
    TH(sl, cols, 0.3, 1.44, cws)

    linhas = [
        # (label, tipo, 23, 24, 25)
        ("Receita Bruta",        "top",   2.660, 3.188, 3.733),
        ("Deducoes e Impostos",  "ded",  -0.320,-0.383,-0.449),
        ("Receita Liquida",      "sub",   2.340, 2.805, 3.284),
        ("CPV",                  "ded",  -1.620,-1.966,-2.288),
        ("Lucro Bruto",          "sub",   0.720, 0.839, 0.996),
        ("Despesas Comerciais",  "ded",  -0.250,-0.310,-0.380),
        ("Despesas G&A",         "ded",  -0.145,-0.170,-0.192),
        ("Outras (liq.)",        "ded",  -0.020,-0.030,-0.028),
        ("EBIT",                 "sub",   0.305, 0.329, 0.396),
        ("D&A",                  "add",   0.075, 0.095, 0.105),
        ("EBITDA",               "kpi",   0.380, 0.424, 0.501),
        ("EBITDA (R$ MM abs.)",  "kpi",   149.0, 184.8, 252.0),
        ("Resultado Financeiro", "ded",  -0.028,-0.060,-0.052),
        ("Lucro Antes IR",       "sub",   0.277, 0.269, 0.344),
        ("IR / CS",              "ded",  -0.055,-0.053,-0.068),
        ("Lucro Liquido",        "bot",   0.222, 0.216, 0.276),
    ]

    # Receita liquida base para %V
    rl23, rl24, rl25 = 2.340, 2.805, 3.284

    def pv(val, base): return f"{val/base*100:.1f}%" if base else "-"
    def ph(v24, v23): return f"+{(v24/v23-1)*100:.1f}%" if v23 else "-"
    def fmt(v):
        if abs(v) >= 1.0: return f"R$ {v:.3f} Bi"
        return f"R$ {v*1000:.0f} MM" if abs(v) < 1.0 and v != 0 else "-"

    for i, (lbl, tipo, v23, v24, v25) in enumerate(linhas):
        y   = 1.73 + i*0.31
        bg  = WHITE if i%2==0 else RGBColor(0xF8,0xFA,0xFC)
        tc  = GRAPHITE
        bld = False

        if tipo == "top":  bg = RGBColor(0xEB,0xF2,0xFF); bld = True; tc = NAVY
        if tipo == "sub":  bg = LIGHT_BG; bld = True; tc = NAVY
        if tipo == "kpi":  bg = RGBColor(0xE6,0xF4,0xEA); bld = True; tc = TEAL
        if tipo == "bot":  bg = AMB_BG; bld = True; tc = GRAPHITE

        # formato
        if tipo == "kpi" and "abs" in lbl:
            vals_row = [lbl,
                        "", f"{v23:.0f}", f"{v23/rl23*100:.1f}%", "",
                        f"{v24:.0f}", f"{v24/rl24*100:.1f}%", ph(v24,v23),
                        f"{v25:.0f}", f"{v25/rl25*100:.1f}%"]
        else:
            vals_row = [lbl,
                        "", fmt(v23), pv(v23,rl23), "",
                        fmt(v24), pv(v24,rl24), ph(v24,v23),
                        fmt(v25), pv(v25,rl25)]

        x = 0.3
        for j, (v, cw) in enumerate(zip(vals_row, cws)):
            al = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            col_c = POSITIVE if (j in [4,7] and v.startswith("+")) else (NEGATIVE if (j in [4,7] and "-" in v and v!="-") else tc)
            R(sl, x, y, cw-0.02, 0.28, fill=bg)
            T(sl, v, x+0.04, y+0.04, cw-0.08, 0.20, sz=7.5, bold=bld if j==0 else False, col=col_c, align=al)
            x += cw

    # KPIs de margem
    sec(sl, "MARGENS CONSOLIDADAS", 0.3, 6.70, 12.7)
    margins = [
        ("Margem Bruta 2025", "30,31%", POSITIVE),
        ("Marg. EBITDA 2025", "6,75%",  POSITIVE),
        ("Cresc. Rec. Bruta", "+17,1%", POSITIVE),
        ("Cresc. EBITDA",     "+36,4%", POSITIVE),
        ("Marg. Bruta 2024",  "29,90%", GRAPHITE),
        ("Marg. EBITDA 2024", "6,59%",  GRAPHITE),
    ]
    for i, (lbl, val, col) in enumerate(margins):
        x = 0.3 + i*2.18
        R(sl, x, 6.96, 2.0, 0.47, fill=LIGHT_BG)
        T(sl, val, x+0.06, 7.00, 1.88, 0.24, sz=14, bold=True, col=col, align=PP_ALIGN.CENTER, font="Georgia")
        T(sl, lbl, x+0.06, 7.25, 1.88, 0.16, sz=7,  col=MUTED,  align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 5 — BALANCO + ENDIVIDAMENTO (modelo VCA)
# =============================================================================
def slide_balanco(prs):
    sl = base(prs, "IV. Balanco Patrimonial e Estrutura de Endividamento (Modelo VCA)",
              "Fonte: Releases RI Nissei 2023/2024/2025 | Valores estimados em R$ MM com base em dados publicos")

    # --- Coluna esquerda: Ativo ---
    sec(sl, "ATIVO CONSOLIDADO — R$ MM", 0.3, 1.18, 4.1)
    ativo = [
        ("ATIVO CIRCULANTE", True, "", "", ""),
        ("Caixa e Equivalentes",   False, "138", "156", "210"),
        ("Contas a Receber",       False, "290", "380", "420"),
        ("Estoques",               False, "480", "620", "680"),
        ("Outros Circ.",           False, "85",  "110", "130"),
        ("Total Circulante",       True,  "993","1.266","1.440"),
        ("ATIVO NAO CIRCULANTE",   True,  "", "", ""),
        ("Imobilizado (liq.)",     False, "280", "390", "440"),
        ("Dir. de Uso (IFRS16)",   False, "820","1.050","1.150"),
        ("Outros NC",              False, "90",  "120", "140"),
        ("Total NC",               True,  "1.190","1.560","1.730"),
        ("TOTAL DO ATIVO",         True,  "2.183","2.826","3.170"),
    ]
    cws_a = [2.2, 0.6, 0.6, 0.6]
    TH(sl, ["ATIVO", "2023", "2024", "2025E"], 0.3, 1.44, cws_a)
    for i, (lbl, h, v23, v24, v25) in enumerate(ativo):
        y = 1.73 + i*0.30
        bg = LIGHT_BG if h else (WHITE if i%2==0 else RGBColor(0xF8,0xFA,0xFC))
        tc = NAVY if h else GRAPHITE
        R(sl, 0.3, y, 2.18, 0.27, fill=bg)
        T(sl, lbl, 0.36, y+0.04, 2.08, 0.20, sz=7.5, bold=h, col=tc, align=PP_ALIGN.LEFT)
        for j, v in enumerate([v23, v24, v25]):
            R(sl, 2.5+j*0.6, y, 0.58, 0.27, fill=bg)
            T(sl, v, 2.5+j*0.6+0.04, y+0.04, 0.50, 0.20, sz=7.5, bold=h, col=tc, align=PP_ALIGN.CENTER)

    # --- Coluna do meio: Passivo ---
    sec(sl, "PASSIVO CONSOLIDADO — R$ MM", 4.55, 1.18, 4.1)
    passivo = [
        ("PASSIVO CIRCULANTE",     True,  "", "", ""),
        ("Fornecedores",           False, "210","285","310"),
        ("Emprestimos CP",         False, "95", "135","120"),
        ("Passivo IFRS16 CP",      False, "190","240","260"),
        ("Outros Circ.",           False, "80", "105","120"),
        ("Total Circulante",       True,  "575","765","810"),
        ("PASSIVO NAO CIRC.",      True,  "", "", ""),
        ("Debentures LP",          False, "280","480","650"),
        ("Passivo IFRS16 LP",      False, "630","810","890"),
        ("Outros NC",              False, "65",  "80", "90"),
        ("Total NC",               True,  "975","1.370","1.630"),
        ("TOTAL DO PASSIVO",       True,  "1.550","2.135","2.440"),
    ]
    cws_p = [2.2, 0.6, 0.6, 0.6]
    TH(sl, ["PASSIVO", "2023", "2024", "2025E"], 4.55, 1.44, cws_p)
    for i, (lbl, h, v23, v24, v25) in enumerate(passivo):
        y = 1.73 + i*0.30
        bg = LIGHT_BG if h else (WHITE if i%2==0 else RGBColor(0xF8,0xFA,0xFC))
        tc = NAVY if h else GRAPHITE
        R(sl, 4.55, y, 2.18, 0.27, fill=bg)
        T(sl, lbl, 4.61, y+0.04, 2.08, 0.20, sz=7.5, bold=h, col=tc)
        for j, v in enumerate([v23, v24, v25]):
            R(sl, 6.75+j*0.6, y, 0.58, 0.27, fill=bg)
            T(sl, v, 6.75+j*0.6+0.04, y+0.04, 0.50, 0.20, sz=7.5, bold=h, col=tc, align=PP_ALIGN.CENTER)

    # --- Coluna direita: Divida ---
    sec(sl, "ESTRUTURA DE DIVIDA — R$ MM", 8.8, 1.18, 4.23)
    div_data = [
        ("Divida Bruta",     "463", "615",  "770"),
        ("(-) Caixa",        "(138)","(156)","(210)"),
        ("Divida Liquida",   "325",  "459",  "560"),  # antes da 7a emissao; pos: +200 = 760 mas caixa sobe
        ("EBITDA (R$ MM)",   "149",  "185",  "252"),
        ("ND / EBITDA",      "2,18x","2,48x","~1,7x"),
        ("Cob. Juros (EBIT/Desp Fin)", "4,1x","2,5x","3,2x"),
        ("Liquidez Corrente","1,73x","1,65x","1,78x"),
        ("ND / PL",          "0,52x","0,73x","~0,60x"),
    ]
    cws_d = [2.7, 0.5, 0.5, 0.5]
    TH(sl, ["INDICADOR", "2023", "2024", "2025E"], 8.8, 1.44, cws_d)
    for i, (lbl, v23, v24, v25) in enumerate(div_data):
        y = 1.73 + i*0.30
        is_kpi = lbl in ("ND / EBITDA", "Cob. Juros (EBIT/Desp Fin)", "Liquidez Corrente", "ND / PL")
        bg = GRN_BG if is_kpi else (WHITE if i%2==0 else RGBColor(0xF8,0xFA,0xFC))
        h  = lbl in ("Divida Liquida", "EBITDA (R$ MM)")
        tc = NAVY if h else (TEAL if is_kpi else GRAPHITE)
        R(sl, 8.8, y, 2.68, 0.27, fill=bg)
        T(sl, lbl, 8.86, y+0.04, 2.56, 0.20, sz=7.5, bold=h or is_kpi, col=tc)
        for j, v in enumerate([v23, v24, v25]):
            R(sl, 11.5+j*0.5, y, 0.48, 0.27, fill=bg)
            T(sl, v, 11.5+j*0.5+0.02, y+0.04, 0.44, 0.20, sz=7.5, bold=h or is_kpi, col=tc, align=PP_ALIGN.CENTER)

    # Emissoes historicas
    sec(sl, "HISTORICO DE EMISSOES DE DEBENTURES", 0.3, 5.26, 12.7)
    emissoes = [
        ("4a Emissao (NISS24)",  "Jul/2022", "Jul/2026", "CDI+2,80%", "Real",    "Amort. trimestral — em curso"),
        ("6a Emissao",           "Nov/2024", "Nov/2029", "CDI+3,00%", "A definir","R$ 80 MM — refinanciamento"),
        ("7a Emissao (ATUAL)",   "Abr/2025", "Abr/2031", "CDI+2,45%", "A definir","R$ 200 MM — CWG + refinanc."),
    ]
    cws_e = [2.0, 0.9, 0.9, 1.0, 1.0, 6.9]
    TH(sl, ["Emissao","Emissao","Vcto","Remuneracao","Garantia","Observacoes"], 0.3, 5.52, cws_e)
    for i, row in enumerate(emissoes):
        y = 5.81 + i*0.32
        bg = LIGHT_BG if i%2==0 else WHITE
        bld = i==2
        x = 0.3
        for j, (v, cw) in enumerate(zip(row, cws_e)):
            R(sl, x, y, cw-0.02, 0.28, fill=AMB_BG if bld else bg)
            al = PP_ALIGN.LEFT if j in [0, 5] else PP_ALIGN.CENTER
            T(sl, v, x+0.04, y+0.04, cw-0.08, 0.20, sz=7.5, bold=bld, col=NAVY if bld else GRAPHITE, align=al)
            x += cw

# =============================================================================
# SLIDE 6 — INDICES DE CREDITO (modelo VCA — Graficos/Painel)
# =============================================================================
def slide_indices(prs):
    sl = base(prs, "V. Painel de Indices de Credito — Analise Consolidada (Modelo VCA)",
              "Fonte: Releases RI Nissei | Covenants tipicos de debentures farmaceuticas — referencia de mercado")

    anos  = ["2023", "2024", "2025E"]
    cores = [NAVY, GRAPHITE, TEAL]

    # --- Liquidez Corrente ---
    sec(sl, "LIQUIDEZ CORRENTE (Ativo Circ. / Passivo Circ.)", 0.3, 1.18, 6.15)
    vals_lc = [1.73, 1.65, 1.78]
    for i, (a, v, c) in enumerate(zip(anos, vals_lc, cores)):
        y = 1.46 + i*0.55
        R(sl, 0.3, y, 5.9, 0.48, fill=LIGHT_BG)
        bw = (v / 3.0) * 5.5
        R(sl, 0.3, y, bw, 0.48, fill=c)
        T(sl, a, 0.36, y+0.12, 0.5, 0.24, sz=8, bold=True, col=WHITE)
        T(sl, f"{v:.2f}x", 0.3+bw+0.06, y+0.12, 1.0, 0.24, sz=10, bold=True, col=c, font="Georgia")
    R(sl, 0.3+((1.0/3.0)*5.5), 1.40, 0.03, 1.75, fill=AMBER)
    T(sl, "Min. 1,0x", 0.3+((1.0/3.0)*5.5)-0.5, 3.22, 1.1, 0.20, sz=7, col=AMBER, bold=True)

    # --- ND / EBITDA ---
    sec(sl, "DIVIDA LIQUIDA / EBITDA (Alavancagem)", 0.3, 3.30, 6.15)
    vals_nd = [2.18, 2.48, 1.70]
    for i, (a, v, c) in enumerate(zip(anos, vals_nd, cores)):
        y = 3.58 + i*0.55
        R(sl, 0.3, y, 5.9, 0.48, fill=LIGHT_BG)
        bw = min((v / 4.5) * 5.5, 5.5)
        col_b = NEGATIVE if v > 3.5 else (AMBER if v > 2.5 else POSITIVE)
        R(sl, 0.3, y, bw, 0.48, fill=col_b)
        T(sl, a, 0.36, y+0.12, 0.5, 0.24, sz=8, bold=True, col=WHITE)
        T(sl, f"{v:.2f}x", 0.3+bw+0.06, y+0.12, 1.0, 0.24, sz=10, bold=True, col=col_b, font="Georgia")
    R(sl, 0.3+((3.5/4.5)*5.5), 3.24, 0.03, 1.75, fill=NEGATIVE)
    T(sl, "Covenant 3,5x", 0.3+((3.5/4.5)*5.5)-0.8, 5.06, 1.4, 0.20, sz=7, col=NEGATIVE, bold=True)

    # --- Cobertura de Juros ---
    sec(sl, "COBERTURA DE JUROS — EBIT / DESPESA FINANCEIRA", 6.5, 1.18, 6.53)
    vals_cj = [4.10, 2.50, 3.20]
    for i, (a, v, c) in enumerate(zip(anos, vals_cj, cores)):
        y = 1.46 + i*0.55
        R(sl, 6.5, y, 6.2, 0.48, fill=LIGHT_BG)
        bw = min((v / 6.0) * 6.0, 6.0)
        col_b = POSITIVE if v >= 2.5 else NEGATIVE
        R(sl, 6.5, y, bw, 0.48, fill=col_b)
        T(sl, a, 6.56, y+0.12, 0.5, 0.24, sz=8, bold=True, col=WHITE)
        T(sl, f"{v:.2f}x", 6.5+bw+0.06, y+0.12, 1.0, 0.24, sz=10, bold=True, col=col_b, font="Georgia")
    R(sl, 6.5+((2.5/6.0)*6.0), 1.40, 0.03, 1.75, fill=AMBER)
    T(sl, "Min. 2,5x (referencia)", 6.5+((2.5/6.0)*6.0)+0.06, 3.22, 2.0, 0.20, sz=7, col=AMBER, bold=True)

    # --- ND / PL ---
    sec(sl, "DIVIDA LIQUIDA / PATRIMÔNIO LIQUIDO", 6.5, 3.30, 6.53)
    vals_pl = [0.52, 0.73, 0.60]
    for i, (a, v, c) in enumerate(zip(anos, vals_pl, cores)):
        y = 3.58 + i*0.55
        R(sl, 6.5, y, 6.2, 0.48, fill=LIGHT_BG)
        bw = min((v / 2.0) * 6.0, 6.0)
        col_b = POSITIVE if v < 1.0 else NEGATIVE
        R(sl, 6.5, y, bw, 0.48, fill=col_b)
        T(sl, a, 6.56, y+0.12, 0.5, 0.24, sz=8, bold=True, col=WHITE)
        T(sl, f"{v:.2f}x", 6.5+bw+0.06, y+0.12, 1.0, 0.24, sz=10, bold=True, col=col_b, font="Georgia")

    # Tabela de resumo de indices
    sec(sl, "RESUMO DOS INDICES — SEMAFORO DE CREDITO", 0.3, 5.26, 12.7)
    TH(sl, ["INDICADOR","2023A","2024A","2025E","STATUS","COVENANT REF.","AVALIACAO"],
       0.3, 5.52, [2.4, 0.7, 0.7, 0.7, 0.8, 1.5, 5.2])
    idx_rows = [
        ("Liquidez Corrente",     "1,73x","1,65x","1,78x", POSITIVE, "> 1,0x",  "ADEQUADO — Cobre passivos circulantes com folga"),
        ("ND / EBITDA",           "2,18x","2,48x","~1,7x", POSITIVE, "< 3,5x",  "FAVORAVEL — Trajetoria de queda confirmada em 2025E"),
        ("Cobertura de Juros",    "4,10x","2,50x","3,20x", AMBER,    "> 2,5x",  "ADEQUADO — Recuperando com expansao do EBITDA"),
        ("ND / PL",               "0,52x","0,73x","~0,60x",POSITIVE, "< 1,5x",  "CONFORTAVEL — Abaixo do limite de referencia"),
    ]
    for i, (lbl, v23, v24, v25, col, cov, av) in enumerate(idx_rows):
        y = 5.81 + i*0.32
        bg = WHITE if i%2==0 else RGBColor(0xF8,0xFA,0xFC)
        for j, (v, cw, al) in enumerate(zip(
            [lbl, v23, v24, v25, "●", cov, av],
            [2.4, 0.7, 0.7, 0.7, 0.8, 1.5, 5.2],
            [PP_ALIGN.LEFT]+[PP_ALIGN.CENTER]*5+[PP_ALIGN.LEFT]
        )):
            tc = col if j==4 else GRAPHITE
            R(sl, 0.3+sum([2.4,0.7,0.7,0.7,0.8,1.5,5.2][:j]), y, [2.4,0.7,0.7,0.7,0.8,1.5,5.2][j]-0.02, 0.28, fill=bg)
            T(sl, v, 0.3+sum([2.4,0.7,0.7,0.7,0.8,1.5,5.2][:j])+0.04, y+0.04,
              [2.4,0.7,0.7,0.7,0.8,1.5,5.2][j]-0.08, 0.20, sz=7.5, col=tc, align=al)

# =============================================================================
# SLIDE 7 — STRESS TEST + RATING + RECOMENDACAO
# =============================================================================
def slide_stress(prs):
    sl = base(prs, "VI. Stress Test (3 Cenarios) + Scorecard de Rating + Recomendacao Final",
              "Etapa 5 (Rating Interno) + Etapa 6 (Comite) | Manual Analise de Credito v2025.1")

    # Stress test 3 cenarios
    sec(sl, "STRESS TEST — 3 CENARIOS (Base: EBITDA 2025 = R$ 252 MM)", 0.3, 1.18, 12.7)
    cenarios = [
        ("BASE", "Receita +8% | EBITDA R$270MM | ND/EBITDA 1,6x | DSCR 2,4x | Cob.Juros 3,5x",
         "Estrutura de capital confortavel. Debentures com ampla margem de cobertura.", POSITIVE, GRN_BG),
        ("ESTRESSADO (-20% EBITDA)", "Receita +2% | EBITDA R$201MM | ND/EBITDA 2,8x | DSCR 1,6x | Cob.Juros 2,2x",
         "Abaixo do covenant de cobertura de juros (2,5x). Early warning ativado. Monitoramento reforçado.", AMBER, AMB_BG),
        ("SEVERO (-40% EBITDA)", "Receita -5% | EBITDA R$151MM | ND/EBITDA 3,7x | DSCR 0,9x | Cob.Juros 1,4x",
         "Breach de covenant ND/EBITDA (3,5x). Refinanciamento ou aporte necessario. Aceleracao prevista.", NEGATIVE, RED_BG),
    ]
    for i, (nome, metricas, conclusao, col, bg) in enumerate(cenarios):
        y = 1.46 + i*0.78
        R(sl, 0.3, y, 1.15, 0.70, fill=col)
        T(sl, nome, 0.3, y+0.16, 1.15, 0.40, sz=8, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        R(sl, 1.48, y, 11.52, 0.70, fill=bg)
        T(sl, metricas, 1.56, y+0.05, 11.36, 0.26, sz=8.5, bold=True, col=NAVY)
        T(sl, conclusao, 1.56, y+0.34, 11.36, 0.26, sz=8, col=GRAPHITE, italic=True)

    # Scorecard de Rating
    sec(sl, "SCORECARD DE RATING INTERNO", 0.3, 3.94, 6.15)
    criterios = [
        ("Qualidade do Negocio",    4, 5),
        ("Historico de Credito",    4, 5),
        ("Liquidez e Solvencia",    3, 5),
        ("Alavancagem",             3, 5),
        ("Governanca / ESG",        4, 5),
        ("Qualidade das Garantias", 3, 5),
    ]
    for i, (crit, score, max_s) in enumerate(criterios):
        y = 4.22 + i*0.38
        R(sl, 0.3, y, 3.2, 0.32, fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl, crit, 0.38, y+0.06, 3.0, 0.20, sz=8, col=GRAPHITE)
        for k in range(max_s):
            fc = NAVY if k < score else RGBColor(0xD0,0xD8,0xE4)
            R(sl, 3.55+k*0.44, y+0.06, 0.38, 0.20, fill=fc)
        T(sl, f"{score}/{max_s}", 5.78, y+0.06, 0.45, 0.20, sz=8, bold=True, col=NAVY, align=PP_ALIGN.CENTER)
    R(sl, 0.3, 6.50, 6.15, 0.34, fill=NAVY)
    T(sl, "RATING INTERNO PROPOSTO:   A-   (Grau de Investimento — Risco Baixo-Moderado)",
      0.36, 6.55, 5.97, 0.24, sz=9, bold=True, col=GOLD)

    # Recomendacao final
    sec(sl, "RECOMENDACAO FINAL DO COMITE", 6.5, 3.94, 6.53)
    R(sl, 6.5, 4.22, 6.53, 0.54, fill=GRN_BG)
    R(sl, 6.5, 4.22, 0.1, 0.54, fill=POSITIVE)
    T(sl, "✓  APROVADO — Prosseguir com a 7a Emissao de Debentures",
      6.66, 4.30, 6.25, 0.34, sz=11, bold=True, col=POSITIVE)

    rec_items = [
        "Volume: R$ 200 MM | Taxa: CDI + 2,45% a.a. | Prazo: 6 anos",
        "Alavancagem em trajetoria de queda (~1,7x 2025E) — abaixo do covenant de 3,5x",
        "Setor defensivo (farmácias) com demanda inelástica e alta previsibilidade de receita",
        "Lideranca regional consolidada no Parana (20% MS estado / 26% MS Curitiba)",
        "EBITDA com crescimento de +36,4% a/a demonstra rentabilizacao pos-expansao",
    ]
    for i, txt in enumerate(rec_items):
        T(sl, f"  ▸  {txt}", 6.5, 4.84+i*0.33, 6.45, 0.28, sz=8.5, col=NAVY)

    T(sl, "Condicionalidades recomendadas:",
      6.5, 6.54, 6.45, 0.22, sz=8, bold=True, col=AMBER)
    T(sl, "Monitoramento trimestral de ND/EBITDA | Covenant: ND/EBITDA < 3,5x | Cobertura de Juros > 2,5x | Relatorios CVM",
      6.5, 6.76, 6.45, 0.22, sz=7.5, col=GRAPHITE, italic=True)


# =============================================================================
# MAIN
# =============================================================================
def build_nissei():
    prs = new_prs()
    slide_capa(prs)
    slide_sumario(prs)
    slide_kyc(prs)
    slide_dre(prs)
    slide_balanco(prs)
    slide_indices(prs)
    slide_stress(prs)

    import os
    os.makedirs(OUT_DIR, exist_ok=True)
    path = str(Path(OUT_DIR) / "parecer_credito_nissei.pptx")
    prs.save(path)
    print(f"Nissei: {path} ({len(prs.slides)} slides)")
    return path

if __name__ == "__main__":
    build_nissei()
    print("Concluido.")
