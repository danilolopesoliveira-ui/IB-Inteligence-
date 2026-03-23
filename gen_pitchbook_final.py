"""
Pitch Book FINAL — Maximo de qualidade
Claude API (analise) + matplotlib/plotly (graficos) + Securato v2 (design)
Case: CRI BRZ Incorporadora — R$ 250M — Setor Imobiliario
6 blocos, 4 componentes/slide, fontes 14-16pt, callouts 44pt, graficos em todos os slides de dados
"""
import sys,os,json
sys.path.insert(0,os.path.dirname(__file__))
from dotenv import load_dotenv; load_dotenv()
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from datetime import datetime
from io import BytesIO
import anthropic
import matplotlib; matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
try:
    import plotly.graph_objects as go; PLOTLY=True
except: PLOTLY=False

OUT=Path("./templates/models"); OUT.mkdir(parents=True,exist_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# STEP 1 — Claude generates full analysis
# ═══════════════════════════════════════════════════════════════════════════
print("STEP 1: Claude gerando analise completa...\n")
client=anthropic.Anthropic()

PROMPT="""Voce e um analista senior de IB preparando um pitch book para:

EMPRESA: BRZ Incorporadora S.A.
SETOR: Incorporacao imobiliaria residencial e comercial — Sao Paulo, Rio, BH
OPERACAO: CRI (Certificado de Recebiveis Imobiliarios) — R$ 250 milhoes
RATING: A+ (S&P)
PRAZO: 6 anos, IPCA + 7,80% a.a., amortizacao linear a partir do 24o mes
LASTRO: Recebiveis de contratos de venda de unidades imobiliarias
GARANTIA: Alienacao fiduciaria de terrenos (R$ 380M avaliados) + cessao fiduciaria de recebiveis

Gere JSON com EXATAMENTE esta estrutura (dados REALISTAS para incorporadora brasileira mid-cap 2025):
{
  "company": {
    "name":"BRZ Incorporadora S.A.","sector":"Incorporacao imobiliaria residencial e comercial",
    "description":"3 frases descrevendo a empresa com dados concretos",
    "founded":"ano","hq":"cidade-UF","employees":numero,
    "differentials":["3 diferenciais com dados"],
    "governance":"resumo governanca",
    "landbank":"string descrevendo banco de terrenos com VGV potencial",
    "launches_history":{"2022":numero_VGV_M,"2023":numero,"2024":numero,"2025E":numero}
  },
  "financials": {
    "years":["2022","2023","2024","2025E","2026E"],
    "revenue":[nums R$M],"ebitda":[nums],"ebitda_margin":[decimais],
    "net_income":[nums],"net_debt":[nums],"dl_ebitda":[decimais],
    "coverage":[decimais],"fcf":[nums],"capex":[nums],
    "vgv_launched":[nums R$M],"vso":[decimais velocidade vendas]
  },
  "market": {
    "tam_brl_bi":numero,"tam_growth":"CAGR","brazil_position":"posicao Brasil",
    "drivers":["4 drivers com dados de fonte"],
    "kpis_setoriais":{"selic":"taxa","incc":"indice","fgts_funding":"dado","my_casa_vida":"dado"},
    "source":"fontes"
  },
  "deal": {
    "volume":250000000,"instrument":"CRI",
    "rate":"IPCA + 7,80% a.a.","tenor":"6 anos",
    "amort":"Linear a partir do 24o mes","payment":"Semestral",
    "lastro":"Recebiveis de contratos de venda — 1.200 unidades em 8 empreendimentos",
    "guarantee":"Alienacao fiduciaria de terrenos (R$ 380M) + cessao fiduciaria de recebiveis",
    "guarantee_coverage":"2.5x","rating":"A+ (S&P)",
    "registration":"RCVM 160 — automatico","securitizer":"True Securitizadora"
  },
  "comparables": [
    {"issuer":"nome REAL","rating":"","tenor":"","spread":"","volume":"","date":"","source":"ANBIMA"},
    ...6 emissoes REAIS de CRI imobiliario 2024-2025
  ],
  "covenants": {
    "dl_ebitda_max":4.0,"coverage_min":2.0,"dscr_min":1.2,
    "ltv_max":"70%","vso_min":"20%",
    "stress_bear":{"ebitda_drop":-15,"dl_ebitda":numero,"status":"OK/BREACH"},
    "stress_severe":{"ebitda_drop":-30,"dl_ebitda":numero,"status":"OK/BREACH"},
    "breach_threshold":"% queda"
  },
  "investors": [
    {"name":"instituicao REAL","type":"tipo","ticket":"R$ XM","score":"XX%","rationale":"justificativa"},
    ...8 investidores REAIS
  ],
  "risks": [
    {"risk":"risco","probability":"Alta/Media/Baixa","impact":"Alto/Medio/Baixo","mitigant":"mitigante"},
    ...5 riscos
  ],
  "headlines": {
    "exec_summary":"headline conclusivo exec summary",
    "company":"headline conclusivo perfil",
    "market":"headline conclusivo mercado",
    "financials":"headline conclusivo financeiro",
    "margins":"headline conclusivo margens",
    "credit":"headline conclusivo credito",
    "fcf":"headline conclusivo fluxo de caixa",
    "deal":"headline conclusivo estrutura",
    "guarantee":"headline conclusivo garantias",
    "pricing":"headline conclusivo pricing",
    "covenants":"headline conclusivo covenants",
    "investors":"headline conclusivo demanda",
    "timeline":"headline conclusivo cronograma",
    "risks":"headline conclusivo riscos"
  },
  "investment_thesis":"3 frases da tese de investimento"
}

REGRAS: dados REALISTAS para incorporadora brasileira 2025, comparaveis e investidores REAIS, headlines sao CONCLUSOES assertivas. Responda APENAS com JSON."""

r=client.messages.create(model="claude-sonnet-4-6",max_tokens=8000,messages=[{"role":"user","content":PROMPT}])
raw=r.content[0].text.strip()
if raw.startswith("```"): raw=raw.split("\n",1)[1].rsplit("```",1)[0]
D=json.loads(raw)
print(f"  Tokens: {r.usage.input_tokens}in/{r.usage.output_tokens}out")
print(f"  Empresa: {D['company']['name']} | Volume: R${D['deal']['volume']/1e6:.0f}M\n")

# ═══════════════════════════════════════════════════════════════════════════
# STEP 2 — Charts
# ═══════════════════════════════════════════════════════════════════════════
print("STEP 2: Gerando graficos...\n")
CN='#0A2342';CB='#1E5F8C';CG='#2D7A4F';CR='#C0392B';CGD='#C9A84C';CGR='#8492A6';CL='#E8ECF0'
fin=D['financials']; yrs=fin['years']

def mfig(w=6.8,h=4.0):
    fig,ax=plt.subplots(figsize=(w,h))
    ax.spines['top'].set_visible(False);ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(CL);ax.spines['bottom'].set_color(CL)
    ax.tick_params(colors=CGR,labelsize=10);ax.grid(axis='y',color=CL,linewidth=0.5)
    return fig,ax

def sf(fig):
    buf=BytesIO();fig.savefig(buf,format='png',dpi=220,bbox_inches='tight',facecolor='white');buf.seek(0);plt.close(fig);return buf

# 1. Revenue + EBITDA
fig,ax=mfig(); x=np.arange(len(yrs));bw=0.32
ax.bar(x-bw/2,fin['revenue'],bw,label='Receita Liquida',color=CN,alpha=0.9)
ax.bar(x+bw/2,fin['ebitda'],bw,label='EBITDA',color=CB,alpha=0.9)
for i,v in enumerate(fin['revenue']): ax.text(i-bw/2,v+max(fin['revenue'])*0.02,f'{v:,.0f}',ha='center',fontsize=9,color=CN,fontweight='bold')
for i,v in enumerate(fin['ebitda']): ax.text(i+bw/2,v+max(fin['revenue'])*0.02,f'{v:,.0f}',ha='center',fontsize=9,color=CB,fontweight='bold')
ax.set_xticks(x);ax.set_xticklabels(yrs,fontsize=11);ax.set_ylabel('R$ Milhoes',fontsize=10,color=CGR)
ax.legend(fontsize=10,frameon=False);ch_rev=sf(fig);print("  ch_rev OK")

# 2. Margins
fig,ax=mfig();mg=[m*100 for m in fin['ebitda_margin']]
ax.plot(yrs,mg,'o-',color=CN,linewidth=2.5,markersize=9);ax.fill_between(range(len(yrs)),mg,alpha=0.08,color=CN)
for i,v in enumerate(mg): ax.annotate(f'{v:.1f}%',(yrs[i],v),textcoords="offset points",xytext=(0,14),fontsize=10,ha='center',color=CN,fontweight='bold')
ax.set_ylim(min(mg)-5,max(mg)+10);ax.set_ylabel('Margem EBITDA (%)',fontsize=10,color=CGR)
ch_mg=sf(fig);print("  ch_margins OK")

# 3. Leverage dual axis
fig,ax1=mfig()
ax1.bar(yrs,fin['dl_ebitda'],color=CN,width=0.5,alpha=0.85,label='DL/EBITDA')
ax1.axhline(y=D['covenants']['dl_ebitda_max'],color=CR,linestyle='--',linewidth=1.5,label=f'Covenant ({D["covenants"]["dl_ebitda_max"]}x)')
for i,v in enumerate(fin['dl_ebitda']): ax1.text(i,v+0.06,f'{v:.1f}x',ha='center',fontsize=10,color=CN,fontweight='bold')
ax1.set_ylabel('DL/EBITDA (x)',fontsize=10,color=CGR);ax1.set_ylim(0,D['covenants']['dl_ebitda_max']+1.5)
ax2=ax1.twinx()
ax2.plot(yrs,fin['coverage'],'D-',color=CG,linewidth=2.5,markersize=8,label='Coverage')
for i,v in enumerate(fin['coverage']): ax2.annotate(f'{v:.1f}x',(yrs[i],v),textcoords="offset points",xytext=(0,12),fontsize=10,ha='center',color=CG,fontweight='bold')
ax2.set_ylabel('Coverage (x)',fontsize=10,color=CGR)
l1,la1=ax1.get_legend_handles_labels();l2,la2=ax2.get_legend_handles_labels()
ax1.legend(l1+l2,la1+la2,fontsize=9,frameon=False,loc='upper right')
for a in [ax1,ax2]: a.spines['top'].set_visible(False);a.tick_params(colors=CGR,labelsize=10)
ax1.spines['right'].set_color(CL);ch_lev=sf(fig);print("  ch_leverage OK")

# 4. FCF vs Capex
fig,ax=mfig()
ax.bar(np.arange(len(yrs))-0.2,fin['fcf'],0.4,label='FCF',color=CG,alpha=0.85)
ax.bar(np.arange(len(yrs))+0.2,fin['capex'],0.4,label='Capex',color=CB,alpha=0.6)
for i,v in enumerate(fin['fcf']): ax.text(i-0.2,v+max(fin['fcf'])*0.03,f'{v:,.0f}',ha='center',fontsize=9,color=CG,fontweight='bold')
ax.set_xticks(range(len(yrs)));ax.set_xticklabels(yrs,fontsize=11);ax.set_ylabel('R$ Milhoes',fontsize=10,color=CGR)
ax.legend(fontsize=10,frameon=False);ch_fcf=sf(fig);print("  ch_fcf OK")

# 5. VGV Lancado + VSO
fig,ax1=mfig()
vgv=fin.get('vgv_launched',fin['revenue']);vso=[v*100 for v in fin.get('vso',[0.3]*len(yrs))]
ax1.bar(yrs,vgv,color=CGD,width=0.5,alpha=0.8,label='VGV Lancado (R$M)')
for i,v in enumerate(vgv): ax1.text(i,v+max(vgv)*0.02,f'{v:,.0f}',ha='center',fontsize=9,color='#8B7330',fontweight='bold')
ax1.set_ylabel('VGV (R$ M)',fontsize=10,color=CGR)
ax2=ax1.twinx()
ax2.plot(yrs,vso,'s-',color=CN,linewidth=2.5,markersize=8,label='VSO (%)')
for i,v in enumerate(vso): ax2.annotate(f'{v:.0f}%',(yrs[i],v),textcoords="offset points",xytext=(0,12),fontsize=10,ha='center',color=CN,fontweight='bold')
ax2.set_ylabel('VSO (%)',fontsize=10,color=CGR)
l1,la1=ax1.get_legend_handles_labels();l2,la2=ax2.get_legend_handles_labels()
ax1.legend(l1+l2,la1+la2,fontsize=9,frameon=False);ax1.spines['right'].set_color(CL)
for a in [ax1,ax2]: a.spines['top'].set_visible(False);a.tick_params(colors=CGR,labelsize=10)
ch_vgv=sf(fig);print("  ch_vgv OK")

# 6. Stress test (plotly)
ch_stress=None
if PLOTLY:
    cov=D['covenants'];cur=fin['dl_ebitda'][-2] if len(fin['dl_ebitda'])>1 else fin['dl_ebitda'][-1]
    sc=['Base',f'Bear ({cov["stress_bear"]["ebitda_drop"]}%)',f'Stressed ({cov["stress_severe"]["ebitda_drop"]}%)',f'Covenant ({cov["dl_ebitda_max"]}x)']
    vl=[cur,cov['stress_bear']['dl_ebitda'],cov['stress_severe']['dl_ebitda'],cov['dl_ebitda_max']]
    cl=[CN,CGD,CR if cov['stress_severe']['status']=='BREACH' else CGD,'#888888']
    fig=go.Figure(go.Bar(x=sc,y=vl,marker_color=cl,text=[f'{v:.1f}x' for v in vl],textposition='outside',textfont_size=16))
    fig.add_hline(y=cov['dl_ebitda_max'],line_dash="dash",line_color=CR,line_width=2,annotation_text=f"Covenant {cov['dl_ebitda_max']}x",annotation_font_size=12)
    fig.update_layout(width=700,height=400,margin=dict(l=10,r=10,t=30,b=10),plot_bgcolor='white',paper_bgcolor='white',
        yaxis=dict(title='DL/EBITDA (x)',gridcolor=CL,title_font_size=12,tickfont_size=11),xaxis=dict(tickfont_size=11),font=dict(family='Calibri',color=CN))
    buf=BytesIO();fig.write_image(buf,format='png',scale=2);buf.seek(0);ch_stress=buf;print("  ch_stress OK (plotly)")

# 7. Guarantee waterfall
fig,ax=mfig(w=6.8,h=3.5)
g_labels=['Alienacao Fid.\nTerrenos','Cessao Fid.\nRecebiveis','TOTAL\nGarantias','Volume\nCRI']
g_vals=[380,D['deal']['volume']/1e6*float(D['deal']['guarantee_coverage'].replace('x',''))-380,380+D['deal']['volume']/1e6*float(D['deal']['guarantee_coverage'].replace('x',''))-380,D['deal']['volume']/1e6]
g_cols=[CN,CB,CG,CR]
bars=ax.barh(range(len(g_labels)),g_vals[:len(g_labels)],color=g_cols[:len(g_labels)],height=0.5,alpha=0.85)
for bar,v in zip(bars,g_vals): ax.text(bar.get_width()+3,bar.get_y()+bar.get_height()/2,f'R$ {v:.0f}M',va='center',fontsize=10,fontweight='bold',color=CN)
ax.axvline(x=D['deal']['volume']/1e6,color=CR,linestyle='--',linewidth=1.5,alpha=0.6)
ax.set_xlabel('R$ Milhoes',fontsize=10,color=CGR);ax.set_yticks(range(len(g_labels)));ax.set_yticklabels(g_labels,fontsize=10)
ch_gar=sf(fig);print("  ch_guarantee OK")
print()

# ═══════════════════════════════════════════════════════════════════════════
# STEP 3 — Build PPTX (Securato v2: 6 blocos, 4 componentes/slide)
# ═══════════════════════════════════════════════════════════════════════════
print("STEP 3: Construindo PPTX...\n")
prs=Presentation();prs.slide_width=Inches(13.33);prs.slide_height=Inches(7.5)
co=D['company'];dl=D['deal'];hl=D['headlines'];mkt=D['market']

# Paleta Securato v2
P=lambda r,g,b:RGBColor(r,g,b)
NAV=P(0x0A,0x23,0x42);INS=P(0x1E,0x5F,0x8C);GLD=P(0xC9,0xA8,0x4C);OWH=P(0xF5,0xF7,0xFA)
TXT=P(0x1A,0x1A,0x2E);MUT=P(0x84,0x92,0xA6);GRN=P(0x2D,0x7A,0x4F);AMB=P(0xC4,0x7A,0x1A);RED=P(0xC0,0x39,0x2B)
WHT=P(0xFF,0xFF,0xFF);GR3=P(0xE8,0xEC,0xF0)

def base(dark=False):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    bg=s.background.fill;bg.solid();bg.fore_color.rgb=NAV if dark else OWH
    if not dark:
        r=s.shapes.add_shape(1,0,0,Inches(13.33),Emu(22000));r.fill.solid();r.fill.fore_color.rgb=NAV;r.line.fill.background()
    r=s.shapes.add_shape(1,Inches(0.5),Inches(7.0),Inches(12.33),Emu(4000));r.fill.solid();r.fill.fore_color.rgb=MUT if dark else P(0xD0,0xD8,0xE0);r.line.fill.background()
    n=len(prs.slides)
    tb=s.shapes.add_textbox(Inches(0.5),Inches(7.04),Inches(8),Inches(0.25))
    p=tb.text_frame.paragraphs[0];p.text=f"Confidencial  |  IB-Agents  |  {co['name']}  |  CRI R$ {dl['volume']/1e6:.0f}M"
    p.font.size=Pt(8);p.font.color.rgb=WHT if dark else MUT;p.font.name="Calibri"
    tb=s.shapes.add_textbox(Inches(12.4),Inches(7.04),Inches(0.5),Inches(0.25))
    p=tb.text_frame.paragraphs[0];p.text=str(n);p.alignment=PP_ALIGN.RIGHT
    p.font.size=Pt(8);p.font.color.rgb=WHT if dark else MUT;p.font.name="Calibri"
    return s

def ttl(s,text):
    tb=s.shapes.add_textbox(Inches(0.5),Inches(0.2),Inches(12.3),Inches(0.75));tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=text;p.font.size=Pt(20);p.font.bold=True;p.font.color.rgb=NAV;p.font.name="Georgia"
    r=s.shapes.add_shape(1,Inches(0.5),Inches(0.92),Inches(12.33),Emu(4000));r.fill.solid();r.fill.fore_color.rgb=GLD;r.line.fill.background()

def txt(s,lines,x=0.5,y=1.15,w=6.0,h=5.0,sz=14):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));tf=tb.text_frame;tf.word_wrap=True
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.font.name="Calibri"
        if line.startswith("##"): p.text=line[2:].strip();p.font.size=Pt(sz+2);p.font.bold=True;p.font.color.rgb=INS;p.space_before=Pt(16)
        elif line.startswith("•"): p.text=line;p.font.size=Pt(sz);p.font.color.rgb=TXT;p.space_before=Pt(5)
        elif line.startswith("["): p.text=line;p.font.size=Pt(9);p.font.color.rgb=MUT;p.font.italic=True;p.space_before=Pt(4)
        elif line=="": p.text="";p.space_before=Pt(6)
        else: p.text=line;p.font.size=Pt(sz);p.font.color.rgb=TXT;p.space_before=Pt(5)

def co_box(s,val,label,x,y,w=2.8,h=1.3,color=NAV):
    r=s.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h));r.fill.solid();r.fill.fore_color.rgb=GR3;r.line.fill.background()
    tb=s.shapes.add_textbox(Inches(x+0.1),Inches(y+0.05),Inches(w-0.2),Inches(0.8))
    p=tb.text_frame.paragraphs[0];p.text=val;p.alignment=PP_ALIGN.CENTER;p.font.size=Pt(38);p.font.bold=True;p.font.color.rgb=color;p.font.name="Calibri"
    tb=s.shapes.add_textbox(Inches(x+0.1),Inches(y+0.88),Inches(w-0.2),Inches(0.35))
    p=tb.text_frame.paragraphs[0];p.text=label;p.alignment=PP_ALIGN.CENTER;p.font.size=Pt(10);p.font.color.rgb=MUT;p.font.name="Calibri"

def img(s,buf,x=0.5,y=1.15,w=7.0,h=4.3):
    if buf: s.shapes.add_picture(buf,Inches(x),Inches(y),Inches(w),Inches(h))

def tbl(s,data,x=0.5,y=1.15,w=12.3):
    rows=len(data);cols=len(data[0])
    t=s.shapes.add_table(rows,cols,Inches(x),Inches(y),Inches(w),Inches(rows*0.34)).table
    for ri,row in enumerate(data):
        for ci,v in enumerate(row):
            cell=t.cell(ri,ci);cell.text=str(v)
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(10);p.font.name="Calibri"
                p.font.color.rgb=WHT if ri==0 else TXT;p.font.bold=ri==0
            cell.fill.solid();cell.fill.fore_color.rgb=NAV if ri==0 else GR3 if ri%2==0 else OWH

# ═══════════════════════════════════════════════════════════════════════════
# BLOCO 1 — CREDENCIAL
# ═══════════════════════════════════════════════════════════════════════════

# SLIDE 1: CAPA (fundo escuro)
s=base(dark=True)
r=s.shapes.add_shape(1,Inches(0.8),Inches(2.8),Inches(4),Emu(24000));r.fill.solid();r.fill.fore_color.rgb=GLD;r.line.fill.background()
tb=s.shapes.add_textbox(Inches(0.8),Inches(0.8),Inches(11),Inches(1.0))
p=tb.text_frame.paragraphs[0];p.text=co['name'].upper();p.font.size=Pt(40);p.font.bold=True;p.font.color.rgb=WHT;p.font.name="Georgia"
tb=s.shapes.add_textbox(Inches(0.8),Inches(1.9),Inches(11),Inches(0.6))
p=tb.text_frame.paragraphs[0];p.text=f"Certificado de Recebiveis Imobiliarios — CRI  |  R$ {dl['volume']/1e6:.0f} Milhoes"
p.font.size=Pt(18);p.font.color.rgb=GLD;p.font.name="Calibri"
tb=s.shapes.add_textbox(Inches(0.8),Inches(3.2),Inches(11),Inches(0.5))
p=tb.text_frame.paragraphs[0];p.text=f"{dl['rate']}  |  {dl['tenor']}  |  Rating {dl['rating']}"
p.font.size=Pt(14);p.font.color.rgb=P(0xA0,0xB0,0xC8);p.font.name="Calibri"
tb=s.shapes.add_textbox(Inches(0.8),Inches(4.5),Inches(8),Inches(1.5));tf=tb.text_frame
for t in ["IB-Agents Intelligence Platform  |  Coordenador Lider","",datetime.now().strftime("%B %Y"),"Estritamente Privado e Confidencial — RCVM 160/2022"]:
    p=tf.add_paragraph();p.text=t;p.font.size=Pt(10);p.font.color.rgb=MUT;p.font.name="Calibri"

# SLIDE 2: DISCLAIMER
s=base();ttl(s,"Aviso Legal — Restricoes de Uso e Confidencialidade")
txt(s,["Este material foi preparado pelo IB-Agents Intelligence Platform com finalidade exclusivamente informativa e nao constitui oferta, convite ou solicitacao para aquisicao de valores mobiliarios nos termos da Resolucao CVM no 160/2022.",
"","As informacoes contidas neste documento sao baseadas em fontes publicas consideradas confiaveis. Nenhuma declaracao ou garantia, expressa ou implicita, e feita quanto a sua exatidao, completude ou atualidade.",
"","Este documento e confidencial e destinado exclusivamente a investidores profissionais e qualificados conforme a Resolucao CVM no 30/2021. Sua reproducao e expressamente proibida.",
"","[RCVM 160/2022 | RCVM 30/2021 | Codigo ANBIMA de Ofertas Publicas | Lei 6.404/1976 | Lei 6.385/1976]"],sz=14)

# SLIDE 3: EXECUTIVE SUMMARY (SCQA + 6 callouts)
s=base();ttl(s,hl['exec_summary'])
txt(s,[
    "## Tese de Investimento",
    D['investment_thesis'],
    "",
    "## Destaques",
    f"• Volume: R$ {dl['volume']/1e6:.0f}M em CRI | {dl['rate']} | {dl['tenor']} | {dl['guarantee']}",
    f"• {co['description']}",
],y=1.15,w=7.5,sz=14)
co_box(s,f"R${dl['volume']/1e6:.0f}M","Volume CRI",x=8.5,y=1.15,color=NAV)
co_box(s,dl['rating'].split('(')[0].strip(),"Rating S&P",x=8.5,y=2.6,color=GRN)
co_box(s,f"{fin['dl_ebitda'][-2]:.1f}x","DL/EBITDA",x=8.5,y=4.05,color=GRN)
co_box(s,dl['guarantee_coverage'],"Cobertura Garantias",x=11.3,y=1.15,color=NAV)
co_box(s,f"{fin['ebitda_margin'][-2]*100:.0f}%","Margem EBITDA",x=11.3,y=2.6,color=INS)
co_box(s,f"{fin['coverage'][-2]:.1f}x","Coverage",x=11.3,y=4.05,color=GRN)

# ═══════════════════════════════════════════════════════════════════════════
# BLOCO 2 — CONTEXTO
# ═══════════════════════════════════════════════════════════════════════════

# SLIDE 4: COMPANY OVERVIEW
s=base();ttl(s,hl['company'])
txt(s,[
    "## Visao Geral",co['description'],
    f"• Fundacao: {co['founded']} | Sede: {co['hq']} | {co['employees']} colaboradores",
    f"• {co.get('landbank','Banco de terrenos com VGV potencial significativo')}",
    "","## Diferenciais Competitivos",
]+[f"• {d}" for d in co['differentials']]+[
    "","## Governanca",f"• {co['governance']}",
],y=1.15,w=8.0,sz=14)
co_box(s,f"R${fin['revenue'][-2]:,.0f}M","Receita 2024",x=9.0,y=1.15,color=NAV)
co_box(s,f"{co['employees']:,}","Colaboradores",x=9.0,y=2.6,color=INS)

# SLIDE 5: MERCADO
s=base();ttl(s,hl['market'])
txt(s,[
    f"## TAM: R$ {mkt['tam_brl_bi']}Bi ({mkt['tam_growth']})",
    f"• {mkt['brazil_position']}",
    "","## Drivers",
]+[f"• {d}" for d in mkt['drivers']]+[
    f"","[Fonte: {mkt['source']}]",
],y=1.15,w=8.0,sz=14)
co_box(s,f"R${mkt['tam_brl_bi']}Bi","TAM",x=9.0,y=1.15,color=NAV)
co_box(s,mkt['tam_growth'],"CAGR Mercado",x=9.0,y=2.6,color=GRN)

# ═══════════════════════════════════════════════════════════════════════════
# BLOCO 3 — FINANCEIRO (4 slides com graficos)
# ═══════════════════════════════════════════════════════════════════════════

# SLIDE 6: RECEITA + EBITDA
s=base();ttl(s,hl['financials'])
img(s,ch_rev,x=0.5,y=1.15,w=7.2,h=4.5)
rcagr=((fin['revenue'][-2]/fin['revenue'][0])**(1/(len(yrs)-2))-1)*100
ecagr=((fin['ebitda'][-2]/fin['ebitda'][0])**(1/(len(yrs)-2))-1)*100
txt(s,["## Destaques",f"• Receita CAGR: {rcagr:.0f}%",f"• EBITDA CAGR: {ecagr:.0f}%",
    f"• Margem EBITDA: {fin['ebitda_margin'][-2]*100:.1f}%","","[DFs IFRS auditadas]"],x=8.0,y=1.15,w=5.0,sz=14)
co_box(s,f"{ecagr:.0f}%","CAGR EBITDA",x=10.5,y=4.8,w=2.3,color=GRN)

# SLIDE 7: MARGENS
s=base();ttl(s,hl['margins'])
img(s,ch_mg,x=0.5,y=1.15,w=7.2,h=4.5)
txt(s,["## Evolucao","• Margem EBITDA estavel/crescente","• Indica poder de precificacao","• Escala operacional gerando eficiencia"],x=8.0,y=1.15,w=5.0,sz=14)

# SLIDE 8: CREDITO + LEVERAGE
s=base();ttl(s,hl['credit'])
img(s,ch_lev,x=0.5,y=1.15,w=7.2,h=4.5)
cov=D['covenants'];folga=cov['dl_ebitda_max']-fin['dl_ebitda'][-2]
txt(s,["## Credito",f"• DL/EBITDA: {fin['dl_ebitda'][-2]:.1f}x (cov: {cov['dl_ebitda_max']}x — folga {folga:.1f}x)",
    f"• Coverage: {fin['coverage'][-2]:.1f}x (cov: {cov['coverage_min']}x)",
    "","## Stress Test",
    f"• Bear ({cov['stress_bear']['ebitda_drop']}%): {cov['stress_bear']['dl_ebitda']:.1f}x — {cov['stress_bear']['status']}",
    f"• Stressed ({cov['stress_severe']['ebitda_drop']}%): {cov['stress_severe']['dl_ebitda']:.1f}x — {cov['stress_severe']['status']}",
],x=8.0,y=1.15,w=5.0,sz=14)
co_box(s,f"{folga:.1f}x","Folga Covenant",x=10.5,y=5.0,w=2.3,color=GRN)

# SLIDE 9: VGV + VSO (KPI setorial)
s=base();ttl(s,"VGV lancado cresce consistentemente com velocidade de vendas acima da media do setor")
img(s,ch_vgv,x=0.5,y=1.15,w=7.2,h=4.5)
txt(s,["## KPIs Operacionais","• VGV = indicador-chave do setor imobiliario","• VSO (Vendas sobre Oferta) acima da media ABRAINC","• Landbank robusto para lancamentos futuros"],x=8.0,y=1.15,w=5.0,sz=14)

# SLIDE 10: FCF
s=base();ttl(s,hl['fcf'])
img(s,ch_fcf,x=0.5,y=1.15,w=7.2,h=4.5)
txt(s,["## Fluxo de Caixa",f"• FCF 2024: R$ {fin['fcf'][-2]:,.0f}M",f"• Capex: R$ {fin['capex'][-2]:,.0f}M","• Geracao sustenta servico da divida"],x=8.0,y=1.15,w=5.0,sz=14)

# ═══════════════════════════════════════════════════════════════════════════
# BLOCO 4 — TRANSACAO
# ═══════════════════════════════════════════════════════════════════════════

# SLIDE 11: TERM SHEET
s=base();ttl(s,hl['deal'])
tbl(s,[["Caracteristica","Detalhe"],
    ["Instrumento",f"CRI — {dl.get('securitizer','Securitizadora')}"],
    ["Volume",f"R$ {dl['volume']/1e6:.0f}.000.000"],["Remuneracao",dl['rate']],
    ["Prazo",dl['tenor']],["Amortizacao",dl['amort']],["Pgto Juros",dl['payment']],
    ["Lastro",dl['lastro']],["Garantia",dl['guarantee']],
    ["Cobertura",dl['guarantee_coverage']],["Rating",dl['rating']],
    ["Registro",dl['registration']],
],y=1.15)
co_box(s,dl['guarantee_coverage'],"Cobertura",x=10.0,y=5.3,w=2.8,color=GRN)

# SLIDE 12: GARANTIAS + CHART
s=base();ttl(s,hl.get('guarantee',"Garantia com cobertura de "+dl['guarantee_coverage']+" sobre o volume — protecao robusta"))
img(s,ch_gar,x=0.5,y=1.15,w=7.2,h=4.0)
txt(s,["## Estrutura",f"• {dl['guarantee']}","• Cobertura: "+dl['guarantee_coverage']+" sobre o volume","• Conta vinculada com fluxo segregado"],x=8.0,y=1.15,w=5.0,sz=14)

# SLIDE 13: COMPARAVEIS
s=base();ttl(s,hl['pricing'])
t=[["Emissor","Rating","Prazo","Spread","Volume","Data"]]
for c in D['comparables']: t.append([c['issuer'],c['rating'],c['tenor'],c['spread'],c['volume'],c['date']])
t.append([co['name']+" (proposta)",dl['rating'],dl['tenor'],dl['rate'],f"R$ {dl['volume']/1e6:.0f}M",datetime.now().strftime("%b/%y")])
tbl(s,t,y=1.15)
tb=s.shapes.add_textbox(Inches(0.5),Inches(6.5),Inches(12),Inches(0.3))
p=tb.text_frame.paragraphs[0];p.text="[Fonte: ANBIMA — Boletim Mercado Secundario + emissoes primarias 2024-2025]"
p.font.size=Pt(8);p.font.color.rgb=MUT;p.font.name="Calibri";p.font.italic=True

# SLIDE 14: STRESS TEST (plotly)
if ch_stress:
    s=base();ttl(s,hl['covenants'])
    img(s,ch_stress,x=0.5,y=1.15,w=7.2,h=4.5)
    txt(s,["## Cenarios",f"• Base: {fin['dl_ebitda'][-2]:.1f}x — OK",
        f"• Bear: {cov['stress_bear']['dl_ebitda']:.1f}x — {cov['stress_bear']['status']}",
        f"• Stressed: {cov['stress_severe']['dl_ebitda']:.1f}x — {cov['stress_severe']['status']}",
        f"• Breach: {cov['breach_threshold']}"],x=8.0,y=1.15,w=5.0,sz=14)

# ═══════════════════════════════════════════════════════════════════════════
# BLOCO 5 — EXECUCAO
# ═══════════════════════════════════════════════════════════════════════════

# SLIDE 15: INVESTIDORES
s=base();ttl(s,hl['investors'])
t=[["#","Instituicao","Tipo","Ticket","Score","Racional"]]
for i,inv in enumerate(D['investors']): t.append([str(i+1),inv['name'],inv['type'],inv['ticket'],inv['score'],inv.get('rationale','')])
tbl(s,t,y=1.15)

# SLIDE 16: RISCOS
s=base();ttl(s,hl['risks'])
t=[["Risco","Probabilidade","Impacto","Mitigante"]]
for ri in D['risks']: t.append([ri['risk'],ri['probability'],ri['impact'],ri['mitigant']])
tbl(s,t,y=1.15)
txt(s,["","## Parecer: RISCO BAIXO-MEDIO","Rating A+ sustentado por garantia real ("+dl['guarantee_coverage']+"), lastro diversificado e historico de execucao."],y=5.0,w=12,sz=14)

# ═══════════════════════════════════════════════════════════════════════════
# BLOCO 6 — ENCERRAMENTO
# ═══════════════════════════════════════════════════════════════════════════

# SLIDE 17: PRESTADORES + BASE LEGAL
s=base();ttl(s,"Prestadores de Servico e Base Legal Aplicavel")
tbl(s,[["Funcao","Instituicao"],
    ["Coordenador Lider","IB-Agents Intelligence Platform"],
    ["Securitizadora",dl.get('securitizer','True Securitizadora S.A.')],
    ["Cedente/Emissora",co['name']],
    ["Agente Fiduciario","Pentagonal S.A. DTVM"],
    ["Escriturador","Banco Bradesco S.A."],
    ["Assessor Juridico","Pinheiro Neto Advogados"],
    ["Rating",f"{dl['rating'].split('(')[1].replace(')','') if '(' in dl['rating'] else 'S&P'} ({dl['rating'].split('(')[0].strip()})"],
    ["Auditores","BDO Brasil"],
],y=1.15)
txt(s,["## Base Legal","• RCVM 160/2022 | RCVM 30/2021 | Codigo ANBIMA | Lei 6.404/1976 | Lei 6.385/1976 | Lei 9.514/1997 (CRI)"],y=5.0,w=12,sz=12)

# ═══════════════════════════════════════════════════════════════════════════
fn="pitchbook_final_cri_brz.pptx"
prs.save(str(OUT/fn))
print(f"\nGerado: {fn} ({len(prs.slides)} slides)")
print(f"Local: {(OUT/fn).resolve()}")
print(f"API: {r.usage.input_tokens}in/{r.usage.output_tokens}out (~${(r.usage.input_tokens*3/1e6+r.usage.output_tokens*15/1e6):.4f})")
