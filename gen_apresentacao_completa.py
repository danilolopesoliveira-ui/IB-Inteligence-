"""
Apresentacao Institucional Completa — Padrao IB de Mercado
Ferramentas: python-pptx + matplotlib + plotly (kaleido export)
Design: Farallon (branco, navy, serif)
Narrativa: Pyramid Principle + SCQA
Case: Debentures Eneva S.A. R$ 800M
"""
import os, sys
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
from datetime import datetime
from io import BytesIO

# Try to import chart libs
try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.ticker as mticker
    HAS_MPL = True
except ImportError:
    HAS_MPL = False

try:
    import plotly.graph_objects as go
    import plotly.io as pio
    HAS_PLOTLY = True
except ImportError:
    HAS_PLOTLY = False

OUT = Path("./templates/models")
OUT.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

# Farallon palette
NAV = RGBColor(0x0F,0x1F,0x3D); NV2 = RGBColor(0x1A,0x2B,0x4A)
GRF = RGBColor(0x2D,0x37,0x48); GR1 = RGBColor(0x71,0x80,0x96)
GR2 = RGBColor(0xA0,0xAE,0xC0); GR3 = RGBColor(0xE2,0xE8,0xF0)
WHT = RGBColor(0xFF,0xFF,0xFF); BLK = RGBColor(0x1A,0x20,0x2C)
RED = RGBColor(0xC5,0x30,0x30); GRN = RGBColor(0x27,0x67,0x49)

# Chart colors (matplotlib/plotly)
C_NAV = '#0F1F3D'; C_NV2 = '#1A2B4A'; C_GR = '#718096'; C_GR3 = '#E2E8F0'
C_GRN = '#276749'; C_RED = '#C53030'; C_BLU = '#3182CE'; C_GLD = '#C9A962'

# ── Helpers ──────────────────────────────────────────────────────────────────

def sl():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
    r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Emu(20000))
    r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
    r = s.shapes.add_shape(1, Inches(0.5), Inches(7.0), Inches(12.33), Emu(5000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()
    n = len(prs.slides)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = "Confidencial  |  IB-Agents  |  Eneva S.A."
    p.font.size = Pt(6.5); p.font.color.rgb = GR1; p.font.name = "Calibri"
    tb = s.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(6.5); p.font.color.rgb = GR1; p.font.name = "Calibri"
    return s

def hl(s, text):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.3), Inches(0.55))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAV; p.font.name = "Georgia"
    r = s.shapes.add_shape(1, Inches(0.5), Inches(0.76), Inches(12.33), Emu(4000))
    r.fill.solid(); r.fill.fore_color.rgb = GR2; r.line.fill.background()

def body(s, lines, x=0.5, y=1.0, w=12, h=5.5):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = "Calibri"
        if line.startswith("##"):
            p.text = line[2:].strip(); p.font.size = Pt(11); p.font.bold = True
            p.font.color.rgb = NAV; p.space_before = Pt(12)
        elif line.startswith("•"):
            p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRF; p.space_before = Pt(2)
        elif line.startswith("["):
            p.text = line; p.font.size = Pt(7.5); p.font.color.rgb = GR1
            p.font.italic = True; p.space_before = Pt(2)
        elif line == "": p.text = ""; p.space_before = Pt(3)
        else: p.text = line; p.font.size = Pt(9); p.font.color.rgb = GRF; p.space_before = Pt(3)

def tbl(s, data, x=0.5, y=3.8, w=12.3):
    rows = len(data); cols = len(data[0])
    t = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(rows*0.3)).table
    for r, row in enumerate(data):
        for c, v in enumerate(row):
            cell = t.cell(r, c); cell.text = str(v)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8); p.font.name = "Calibri"
                p.font.color.rgb = WHT if r == 0 else BLK; p.font.bold = r == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAV if r == 0 else GR3 if r % 2 == 0 else WHT

def kpis(s, items, y=1.0):
    n = len(items); bw = 12.0 / n; gap = 0.15
    for i, (val, label, color) in enumerate(items):
        x = 0.5 + i * (bw + gap * (1 if i < n-1 else 0))
        r = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(bw - gap), Inches(0.9))
        r.fill.solid(); r.fill.fore_color.rgb = GR3; r.line.fill.background()
        tb = s.shapes.add_textbox(Inches(x+0.15), Inches(y+0.05), Inches(bw-0.3), Inches(0.5))
        p = tb.text_frame.paragraphs[0]; p.text = val
        p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = color or NAV; p.font.name = "Calibri"
        tb = s.shapes.add_textbox(Inches(x+0.15), Inches(y+0.55), Inches(bw-0.3), Inches(0.3))
        p = tb.text_frame.paragraphs[0]; p.text = label
        p.font.size = Pt(8); p.font.color.rgb = GR1; p.font.name = "Calibri"

def add_chart_img(s, buf, x=0.5, y=1.0, w=5.5, h=3.5):
    """Insert matplotlib/plotly chart image into slide."""
    s.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))


# ── Chart generators ─────────────────────────────────────────────────────────

def chart_revenue_ebitda():
    if not HAS_MPL: return None
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    years = ['2020','2021','2022','2023','2024']
    rev = [3820, 4850, 5920, 7180, 8200]
    ebitda = [1580, 2180, 2650, 3320, 3800]
    x = range(len(years)); bw = 0.35
    ax.bar([i-bw/2 for i in x], rev, bw, label='Receita Liquida', color=C_NAV)
    ax.bar([i+bw/2 for i in x], ebitda, bw, label='EBITDA', color=C_BLU)
    for i, v in enumerate(rev): ax.text(i-bw/2, v+100, f'{v:,}', ha='center', fontsize=6.5, color=C_NAV)
    for i, v in enumerate(ebitda): ax.text(i+bw/2, v+100, f'{v:,}', ha='center', fontsize=6.5, color=C_BLU)
    ax.set_xticks(x); ax.set_xticklabels(years, fontsize=8)
    ax.set_ylabel('R$ Milhoes', fontsize=8, color=C_GR)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{x/1000:.0f}k'))
    ax.legend(fontsize=7, frameon=False)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(C_GR3); ax.spines['bottom'].set_color(C_GR3)
    ax.tick_params(colors=C_GR, labelsize=7)
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0)
    plt.close(fig)
    return buf

def chart_margins():
    if not HAS_MPL: return None
    fig, ax = plt.subplots(figsize=(5.5, 3.2))
    years = ['2020','2021','2022','2023','2024']
    mg_bruta = [39.3, 46.0, 46.8, 47.4, 48.0]
    mg_ebitda = [41.4, 44.9, 44.8, 46.2, 46.3]
    mg_liq = [10.0, 13.9, 14.5, 17.3, 18.5]
    ax.plot(years, mg_bruta, 'o-', color=C_NAV, linewidth=2, markersize=5, label='Margem Bruta')
    ax.plot(years, mg_ebitda, 's-', color=C_BLU, linewidth=2, markersize=5, label='Margem EBITDA')
    ax.plot(years, mg_liq, '^-', color=C_GRN, linewidth=2, markersize=5, label='Margem Liquida')
    for i, v in enumerate(mg_ebitda): ax.annotate(f'{v}%', (years[i], v), textcoords="offset points", xytext=(0,8), fontsize=6.5, ha='center', color=C_BLU)
    ax.set_ylabel('%', fontsize=8, color=C_GR); ax.set_ylim(0, 55)
    ax.legend(fontsize=7, frameon=False, loc='lower right')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(C_GR3); ax.spines['bottom'].set_color(C_GR3)
    ax.tick_params(colors=C_GR, labelsize=7); ax.grid(axis='y', color=C_GR3, linewidth=0.5)
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0)
    plt.close(fig)
    return buf

def chart_leverage():
    if not HAS_MPL: return None
    fig, ax1 = plt.subplots(figsize=(5.5, 3.2))
    years = ['2020','2021','2022','2023','2024']
    dl_ebitda = [2.78, 2.02, 1.77, 1.39, 1.18]
    cov = [3.04, 3.76, 4.08, 5.72, 7.31]
    ax1.bar(years, dl_ebitda, color=C_NAV, width=0.5, label='DL/EBITDA')
    ax1.axhline(y=3.5, color=C_RED, linestyle='--', linewidth=1, label='Covenant (3.5x)')
    for i, v in enumerate(dl_ebitda): ax1.text(i, v+0.08, f'{v}x', ha='center', fontsize=7, color=C_NAV, fontweight='bold')
    ax1.set_ylabel('DL/EBITDA (x)', fontsize=8, color=C_GR); ax1.set_ylim(0, 4.5)
    ax2 = ax1.twinx()
    ax2.plot(years, cov, 'D-', color=C_GRN, linewidth=2, markersize=5, label='EBITDA/Desp.Fin')
    for i, v in enumerate(cov): ax2.annotate(f'{v}x', (years[i], v), textcoords="offset points", xytext=(0,8), fontsize=6.5, ha='center', color=C_GRN)
    ax2.set_ylabel('Coverage (x)', fontsize=8, color=C_GR); ax2.set_ylim(0, 9)
    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1+lines2, labels1+labels2, fontsize=6.5, frameon=False, loc='upper right')
    for ax in [ax1, ax2]:
        ax.spines['top'].set_visible(False)
        ax.tick_params(colors=C_GR, labelsize=7)
    ax1.spines['left'].set_color(C_GR3); ax1.spines['bottom'].set_color(C_GR3); ax1.spines['right'].set_color(C_GR3)
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0)
    plt.close(fig)
    return buf

def chart_football_field():
    if not HAS_PLOTLY: return chart_football_mpl()
    methods = ['52-Week Range', 'EV/EBITDA Comps', 'P/E Comps', 'DCF', 'Precedent Txn']
    lows = [15.20, 16.50, 15.80, 18.50, 17.00]
    mids = [17.80, 19.20, 18.50, 21.30, 20.10]
    highs = [19.50, 22.00, 21.20, 24.10, 23.00]
    fig = go.Figure()
    for i, (m, lo, mi, hi) in enumerate(zip(methods, lows, mids, highs)):
        fig.add_trace(go.Bar(y=[m], x=[hi-lo], base=[lo], orientation='h',
            marker_color=C_NAV, opacity=0.3, showlegend=False, width=0.5))
        fig.add_trace(go.Scatter(y=[m], x=[mi], mode='markers',
            marker=dict(size=10, color=C_NAV, symbol='diamond'), showlegend=False))
    fig.add_vline(x=19.00, line_dash="dash", line_color=C_RED, line_width=2,
        annotation_text="Offer Price R$ 19,00", annotation_position="top right",
        annotation_font_size=9, annotation_font_color=C_RED)
    fig.update_layout(
        width=700, height=320, margin=dict(l=10,r=10,t=30,b=10),
        plot_bgcolor='white', paper_bgcolor='white',
        xaxis=dict(title='R$ / Acao', gridcolor=C_GR3, title_font_size=9, tickfont_size=8),
        yaxis=dict(tickfont_size=8), font=dict(family='Calibri', color=C_NAV),
    )
    buf = BytesIO(); fig.write_image(buf, format='png', scale=2); buf.seek(0)
    return buf

def chart_football_mpl():
    if not HAS_MPL: return None
    fig, ax = plt.subplots(figsize=(7, 3.2))
    methods = ['52-Week Range', 'EV/EBITDA Comps', 'P/E Comps', 'DCF', 'Precedent Txn']
    lows = [15.20, 16.50, 15.80, 18.50, 17.00]
    mids = [17.80, 19.20, 18.50, 21.30, 20.10]
    highs = [19.50, 22.00, 21.20, 24.10, 23.00]
    y = range(len(methods))
    for i in y:
        ax.barh(i, highs[i]-lows[i], left=lows[i], height=0.4, color=C_NAV, alpha=0.25)
        ax.plot(mids[i], i, 'D', color=C_NAV, markersize=7)
        ax.text(highs[i]+0.2, i, f'R$ {mids[i]:.2f}', va='center', fontsize=7, color=C_NAV)
    ax.axvline(x=19.00, color=C_RED, linestyle='--', linewidth=1.5, label='Offer R$ 19,00')
    ax.set_yticks(y); ax.set_yticklabels(methods, fontsize=8)
    ax.set_xlabel('R$ / Acao', fontsize=8, color=C_GR)
    ax.legend(fontsize=7, frameon=False)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(C_GR3); ax.spines['bottom'].set_color(C_GR3)
    ax.tick_params(colors=C_GR, labelsize=7); ax.grid(axis='x', color=C_GR3, linewidth=0.5)
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0)
    plt.close(fig)
    return buf

def chart_sensitivity_heatmap():
    if not HAS_PLOTLY: return chart_sensitivity_mpl()
    import numpy as np
    wacc_labels = ['11,0%','11,5%','12,0%','12,8%','13,5%','14,0%','15,0%']
    g_labels = ['3,0%','3,5%','4,0%','4,5%','5,0%']
    z = [[26.8,28.9,31.5,34.8,39.2],[24.2,25.9,28.0,30.6,33.8],[22.0,23.4,25.1,27.2,29.8],
         [19.2,20.3,21.6,23.2,25.1],[17.0,17.8,18.9,20.1,21.6],[15.6,16.3,17.2,18.2,19.5],
         [13.3,13.9,14.5,15.3,16.2]]
    text = [[f'R$ {v:.1f}' for v in row] for row in z]
    fig = go.Figure(data=go.Heatmap(z=z, x=g_labels, y=wacc_labels, text=text, texttemplate='%{text}',
        colorscale=[[0,C_RED],[0.5,'#FEFCE8'],[1,C_GRN]], showscale=False, textfont_size=9))
    fig.update_layout(width=600, height=320, margin=dict(l=10,r=10,t=30,b=10),
        xaxis=dict(title='Crescimento Perpetuidade (g)', title_font_size=9, tickfont_size=8, side='bottom'),
        yaxis=dict(title='WACC', title_font_size=9, tickfont_size=8, autorange='reversed'),
        font=dict(family='Calibri', color=C_NAV), plot_bgcolor='white', paper_bgcolor='white')
    buf = BytesIO(); fig.write_image(buf, format='png', scale=2); buf.seek(0)
    return buf

def chart_sensitivity_mpl():
    if not HAS_MPL: return None
    import numpy as np
    fig, ax = plt.subplots(figsize=(6, 3.2))
    z = np.array([[26.8,28.9,31.5,34.8,39.2],[24.2,25.9,28.0,30.6,33.8],[22.0,23.4,25.1,27.2,29.8],
         [19.2,20.3,21.6,23.2,25.1],[17.0,17.8,18.9,20.1,21.6],[15.6,16.3,17.2,18.2,19.5],
         [13.3,13.9,14.5,15.3,16.2]])
    im = ax.imshow(z, cmap='RdYlGn', aspect='auto')
    wacc = ['11,0%','11,5%','12,0%','12,8%','13,5%','14,0%','15,0%']
    g = ['3,0%','3,5%','4,0%','4,5%','5,0%']
    ax.set_xticks(range(5)); ax.set_xticklabels(g, fontsize=7)
    ax.set_yticks(range(7)); ax.set_yticklabels(wacc, fontsize=7)
    ax.set_xlabel('g perpetuidade', fontsize=8); ax.set_ylabel('WACC', fontsize=8)
    for i in range(7):
        for j in range(5):
            ax.text(j, i, f'{z[i,j]:.1f}', ha='center', va='center', fontsize=7, color='black' if z[i,j] > 18 else 'white')
    fig.tight_layout()
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', facecolor='white'); buf.seek(0)
    plt.close(fig)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

print("Gerando apresentacao completa com graficos...\n")

# 1. COVER
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill; bg.solid(); bg.fore_color.rgb = WHT
r = s.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(3.3))
r.fill.solid(); r.fill.fore_color.rgb = NAV; r.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
p = tb.text_frame.paragraphs[0]; p.text = "ENEVA S.A."
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHT; p.font.name = "Georgia"
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(0.6))
p = tb.text_frame.paragraphs[0]; p.text = "Emissao de Debentures  |  1a e 2a Series  |  R$ 800.000.000"
p.font.size = Pt(16); p.font.color.rgb = RGBColor(0xC0,0xD0,0xE8); p.font.name = "Calibri"
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11), Inches(0.4))
p = tb.text_frame.paragraphs[0]; p.text = "Rating AA+ (Fitch)  |  Perspectiva Estavel"
p.font.size = Pt(12); p.font.color.rgb = GR2; p.font.name = "Calibri"
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(8), Inches(1.5))
tf = tb.text_frame
for txt in ["Coordenador Lider: IB-Agents Intelligence Platform", datetime.now().strftime("%B %Y"), "Strictly Private and Confidential"]:
    p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(9); p.font.color.rgb = GR1; p.font.name = "Calibri"

# 2. EXEC SUMMARY
s = sl(); hl(s, "Fluxo previsivel e alavancagem de 1,2x sustentam emissao de R$ 800M a spreads competitivos")
body(s, [
    "## Situacao",
    "• Eneva: maior geradora termica a gas do Brasil (5,9 GW) — 92% receita contratada via PPAs de 15+ anos",
    "## Oportunidade",
    "• Mercado de debentures em alta (R$ 320Bi em 2024, +18%) — janela ideal para CDI+1,85% (5y) e IPCA+6,20% (7y)",
    "## Recomendacao",
    "• R$ 800M em 2 series: 1a (refi CDI+, bullet) + 2a (capex IPCA+, linear) — cessao fiduciaria 2,7x cobertura",
    "• Demanda estimada R$ 1,4Bi (1,75x) por 42 instituicoes — oversubscription praticamente garantida",
], y=0.95)
kpis(s, [("R$ 800M","Volume",NAV),("AA+","Rating Fitch",GRN),("1,2x","DL/EBITDA",GRN),("92%","Receita PPA",NAV),("1,75x","Oversub. Est.",NAV)], y=5.5)

# 3. FINANCIALS + CHART
s = sl(); hl(s, "EBITDA cresceu 20% a.a. nos ultimos 3 anos com margem expandindo de 44,9% para 46,3%")
buf = chart_revenue_ebitda()
if buf: add_chart_img(s, buf, x=0.5, y=1.0, w=6.0, h=3.5)
body(s, [
    "## Destaques (R$ M)",
    "• Receita 2024: R$ 8.200 (+14% YoY)",
    "• EBITDA 2024: R$ 3.800 (margem 46,3%)",
    "• Lucro Liq.: R$ 1.518 (+22% YoY)",
    "• CAGR EBITDA 3a: 20,3%",
    "• FCF Operac.: R$ 1.700",
    "",
    "[Fonte: Eneva — DFs IFRS auditadas PwC]",
], x=6.8, y=1.0, w=5.5)

# 4. MARGENS + CHART
s = sl(); hl(s, "Expansao consistente de margens reflete ganhos de escala e integracao vertical gas+energia")
buf = chart_margins()
if buf: add_chart_img(s, buf, x=0.5, y=1.0, w=6.0, h=3.5)
body(s, [
    "## Drivers de Margem",
    "• Integracao vertical: custo de gas 25% abaixo de competidores",
    "• Escala operacional: 5,9 GW com custos fixos diluidos",
    "• PPAs com reajuste IPCA: protecao contra inflacao de custos",
    "• Eficiencia: investimento em manutencao preditiva (-12% paradas)",
    "",
    "## Projecao 2025E",
    "• Receita: R$ 9.200M (+12%)",
    "• EBITDA: R$ 4.300M (margem 47%)",
], x=6.8, y=1.0, w=5.5)

# 5. CREDITO + CHART
s = sl(); hl(s, "Desalavancagem de 2,8x para 1,2x em 4 anos — folga de 2,3x no covenant principal")
buf = chart_leverage()
if buf: add_chart_img(s, buf, x=0.5, y=1.0, w=6.0, h=3.5)
body(s, [
    "## Metricas de Credito (dez/24)",
    "• DL/EBITDA: 1,18x (cov: <=3,5x)",
    "• EBITDA/Desp.Fin: 4,8x (cov: >=2,5x)",
    "• DSCR: 2,45x (cov: >=1,3x)",
    "• Capacidade adic.: R$ 8,8Bi",
    "",
    "## Stress Test",
    "• Bear (-15%): DL/EBITDA 1,39x — OK",
    "• Stressed (-30%): 1,69x — OK",
    "• Breach: queda de -66% no EBITDA",
], x=6.8, y=1.0, w=5.5)

# 6. ESTRUTURA DA EMISSAO
s = sl(); hl(s, "Duas series complementares otimizam custo e perfil de amortizacao")
tbl(s, [
    ["Caracteristica", "1a Serie", "2a Serie"],
    ["Volume", "R$ 480.000.000", "R$ 320.000.000"],
    ["Prazo", "5 anos (Abr/2030)", "7 anos (Abr/2032)"],
    ["Remuneracao", "CDI + 1,85% a.a.", "IPCA + 6,20% a.a."],
    ["Pagamento Juros", "Semestral", "Semestral"],
    ["Amortizacao", "Bullet", "Linear (a partir 5o ano)"],
    ["Garantia", "Cessao fid. recebiveis PPA", "Cessao fid. recebiveis PPA"],
    ["Rating", "AA+ (Fitch)", "AA+ (Fitch)"],
    ["Registro", "RCVM 160 — automatico", "RCVM 160 — automatico"],
    ["Custodia", "B3 (CETIP21)", "B3 (CETIP21)"],
], y=1.0)

# 7. PRICING + COMPS
s = sl(); hl(s, "Spread de CDI+1,85% posiciona Eneva na mediana do setor eletrico AA+")
tbl(s, [
    ["Emissor", "Rating", "Prazo", "Spread", "Volume", "Data"],
    ["CPFL Energia", "AA+", "5y", "CDI+1,78%", "R$ 1,2Bi", "Jan/25"],
    ["Sabesp", "AA+", "5y", "CDI+1,70%", "R$ 1,5Bi", "Mar/25"],
    ["Neoenergia", "AA", "5y", "CDI+1,92%", "R$ 1,5Bi", "Jan/25"],
    ["Equatorial", "AA+", "7y", "IPCA+6,35%", "R$ 2,0Bi", "Mar/25"],
    ["Taesa", "AAA", "7y", "CDI+1,45%", "R$ 800M", "Fev/25"],
    ["Engie Brasil", "AAA", "10y", "IPCA+5,80%", "R$ 1,0Bi", "Fev/25"],
    ["ENEVA (proposta)", "AA+", "5+7y", "1,85%/6,20%", "R$ 800M", "Abr/25"],
], y=1.0)
tb = s.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(12), Inches(0.3))
p = tb.text_frame.paragraphs[0]; p.text = "[Fonte: ANBIMA — Boletim de Mercado Secundario de Debentures, marco/2025]"
p.font.size = Pt(7); p.font.color.rgb = GR1; p.font.name = "Calibri"; p.font.italic = True

# 8. FOOTBALL FIELD
s = sl(); hl(s, "Valuation por multiplas metodologias confirma range de R$ 17,50-21,30 por acao")
buf = chart_football_field()
if buf: add_chart_img(s, buf, x=0.5, y=1.0, w=7.0, h=3.5)
body(s, [
    "## Metodologias",
    "• 52-Week Range: R$ 15,20-19,50",
    "• EV/EBITDA Comps (7,2x): R$ 16,50-22,00",
    "• P/E Comps (13,5x): R$ 15,80-21,20",
    "• DCF (WACC 12,8%): R$ 18,50-24,10",
    "• Precedent Transactions: R$ 17,00-23,00",
    "",
    "• Media ponderada: R$ 19,50/acao",
    "• Cotacao ENEV3: R$ 18,20 (+7%)",
], x=7.8, y=1.0, w=5)

# 9. SENSITIVITY HEATMAP
s = sl(); hl(s, "Preco de R$ 19-21 e robusto mesmo em cenarios conservadores de WACC e crescimento")
buf = chart_sensitivity_heatmap()
if buf: add_chart_img(s, buf, x=0.5, y=1.0, w=6.5, h=3.5)
body(s, [
    "## Leitura do Heatmap",
    "• Centro (12,8% / 4,0%): R$ 21,60",
    "• Cenario conservador (14% / 3%): R$ 15,60",
    "• Cenario otimista (11% / 5%): R$ 39,20",
    "",
    "## Premissas do DCF",
    "• WACC: 12,8% (CAPM Brasil)",
    "• g perpetuidade: 4,0% (PIB nominal LP)",
    "• Projecao: 5 anos + terminal value",
    "• Capex/Receita: 25% (guidance)",
], x=7.3, y=1.0, w=5.5)

# 10. COVENANTS
s = sl(); hl(s, "Pacote de covenants com folga de 2,3x — robusto em cenario de stress -30% EBITDA")
body(s, [
    "## Covenants Financeiros (semestral)",
    "• DL/EBITDA <= 3,5x (atual: 1,2x — folga 2,3x — breach em queda -66%)",
    "• EBITDA/Desp.Fin >= 2,5x (atual: 4,8x)",
    "• DSCR >= 1,3x (atual: 2,45x)",
    "",
    "## Nao-Financeiros",
    "• Cross-default (>R$150M) | Change of control | Restricted Payments (DL/EBITDA>3,0x)",
    "• Alienacao ativos essenciais (>10% AT) vedada sem anuencia | Manutencao seguros",
], y=0.95)
tbl(s, [
    ["Cenario", "EBITDA (R$M)", "DL/EBITDA", "Status", "EBITDA/Fin", "Status"],
    ["Base", "3.800", "1,18x", "OK (folga 2,3x)", "4,8x", "OK"],
    ["Bear (-15%)", "3.230", "1,39x", "OK (folga 2,1x)", "4,1x", "OK"],
    ["Stressed (-30%)", "2.660", "1,69x", "OK (folga 1,8x)", "3,4x", "OK"],
    ["Breach", "1.286 (-66%)", "3,50x", "BREACH", "2,5x", "BREACH"],
], y=4.5)

# 11. DISTRIBUICAO
s = sl(); hl(s, "Demanda de R$ 1,4Bi por 42 instituicoes — oversubscription de 1,75x garante pricing otimo")
tbl(s, [
    ["#", "Instituicao", "Tipo", "AUM (R$Bi)", "Ticket", "Score"],
    ["1", "Itau Asset", "Asset", "780", "R$ 80M", "95%"],
    ["2", "BTG Pactual", "Asset", "320", "R$ 60M", "92%"],
    ["3", "Previ (BB)", "F. Pensao", "250", "R$ 100M", "90%"],
    ["4", "XP Investimentos", "Banco", "150", "R$ 40M", "88%"],
    ["5", "Kinea", "Asset", "85", "R$ 30M", "85%"],
    ["6", "Funcef", "F. Pensao", "110", "R$ 50M", "82%"],
    ["7", "Verde Asset", "Asset", "42", "R$ 25M", "80%"],
    ["8", "SPX Capital", "Asset", "55", "R$ 20M", "78%"],
    ["—", "Outros (34)", "Diversos", "—", "R$ 495M", "60-75%"],
    ["", "TOTAL", "", "", "R$ 1.400M", "1,75x"],
], y=1.0)

# 12. CRONOGRAMA
s = sl(); hl(s, "Bookbuilding em abril abre janela antes do ciclo de corte de SELIC no 2S25")
tbl(s, [
    ["Etapa", "Periodo", "Status"],
    ["Mandato e kick-off", "15/Fev/2025", "Concluido"],
    ["Due diligence", "18/Fev — 10/Mar", "Concluido"],
    ["Rating Fitch", "20/Fev", "Concluido — AA+"],
    ["Escritura de emissao", "01-15/Mar", "Em andamento"],
    ["Aprovacao CA", "18/Mar", "Pendente"],
    ["Registro CVM", "20/Mar", "Pendente"],
    ["Bookbuilding", "25/Mar — 02/Abr", "—"],
    ["Pricing", "02/Abr", "—"],
    ["Liquidacao (D+2)", "04/Abr", "—"],
    ["Inicio negociacao B3", "07/Abr", "—"],
], y=1.0)

# 13. RISCOS
s = sl(); hl(s, "Tres riscos principais mitigados por protecoes estruturais — rating AA+ estavel")
tbl(s, [
    ["Risco", "Prob.", "Impacto", "Mitigante"],
    ["Concentracao MA (70%)", "Media", "Alto", "Diversificacao RR/CE + seguros"],
    ["Reservas de gas", "Baixa", "Alto", "Certificacao DeGolyer 15+ anos"],
    ["Regulatorio ANEEL", "Media", "Medio", "PPAs nao retroativos"],
    ["Liquidez secundaria", "Media", "Baixo", "Base diversificada 42 invest."],
    ["Variacao CDI/IPCA", "Alta", "Medio", "1a serie hedge natural CDI"],
    ["Refinanciamento", "Baixa", "Alto", "Perfil escalonado + FCF R$1,7Bi"],
], y=1.0)
body(s, [
    "",
    "## Parecer Global: RISCO BAIXO",
    "FCF previsivel (92% PPA) + alavancagem confortavel (1,2x) + garantia real (2,7x) + rating AA+ + track record 5 emissoes.",
], y=4.8)

# 14. CONTATOS
s = sl(); hl(s, "Prestadores de Servico e Base Legal")
tbl(s, [
    ["Funcao", "Instituicao"],
    ["Coordenador Lider", "IB-Agents Intelligence Platform"],
    ["Emissora — RI", "Eneva S.A. (ri@eneva.com.br)"],
    ["Agente Fiduciario", "Pentagonal S.A. DTVM"],
    ["Escriturador", "Banco Bradesco S.A."],
    ["Assessor Juridico (Emissor)", "Machado Meyer"],
    ["Assessor Juridico (Coord.)", "Pinheiro Neto"],
    ["Rating", "Fitch Ratings"],
    ["Auditores", "PwC"],
], y=1.0)
body(s, [
    "## Base Legal",
    "• RCVM 160/2022 | RCVM 30/2021 | Codigo ANBIMA | Lei 6.404/1976 (arts. 52-74) | Lei 6.385/1976",
], y=4.5)

# ═══════════════════════════════════════════════════════════════════════════════
fn = "apresentacao_institucional_completa.pptx"
prs.save(str(OUT / fn))
print(f"Gerado: {fn} ({len(prs.slides)} slides)")
print(f"Local: {(OUT / fn).resolve()}")
n_charts = sum(1 for x in [HAS_MPL, HAS_PLOTLY] if x)
print(f"Graficos: matplotlib={'OK' if HAS_MPL else 'N/A'}, plotly={'OK' if HAS_PLOTLY else 'N/A'}")
