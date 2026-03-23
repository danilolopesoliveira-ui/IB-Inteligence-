"""
Gerador de Pareceres IB — DCM e ECM
COBERTURA COMPLETA dos Manuais Guia:
  - Manual_Analise_Credito_v4_1.docx (7 Etapas + Memorando I-VII + Guia Perfeccionista Sec.12)
  - Manual_ECM_IB.docx (Secoes 1-13: KYC, Valuation, DD, Regulacao, Bookbuilding, Flags)
  - Modelo_Analise_Credito_DCM.xlsx (abas: PREMISSAS, DRE, BALANCO, FLUXO, DIVIDA, INDICES, SENSIBILIDADE)
  - Modelo_ECM_IB.xlsx (abas: PREMISSAS, DRE, BALANCO, WACC, DCF, MULTIPLOS, OFERTA, RETORNO, SENSIBILIDADE)
Design: Farallon-inspired | Padrao BTG Pactual / Goldman Sachs
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

NAVY     = RGBColor(0x0F, 0x1F, 0x3D)
NAVY_MED = RGBColor(0x1A, 0x2B, 0x4A)
GRAPHITE = RGBColor(0x2D, 0x37, 0x48)
MUTED    = RGBColor(0x71, 0x80, 0x96)
LIGHT_BG = RGBColor(0xE2, 0xE8, 0xF0)
BORDER   = RGBColor(0xA0, 0xAE, 0xC0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
NEGATIVE = RGBColor(0xC5, 0x30, 0x30)
POSITIVE = RGBColor(0x27, 0x67, 0x49)
GOLD     = RGBColor(0xD4, 0xA8, 0x53)
AMBER    = RGBColor(0xD9, 0x7B, 0x06)
RED_BG   = RGBColor(0xFF, 0xED, 0xED)
AMB_BG   = RGBColor(0xFF, 0xF3, 0xCD)

OUT_DIR  = r"C:\Users\dloge\OneDrive\Área de Trabalho\Teste 1\ib-agents\frontend\public\knowledge"

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs): return prs.slide_layouts[6]

def R(sl, l, t, w, h, fill=None, line=None):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s

def T(sl, text, l, t, w, h, sz=9, bold=False, col=None, align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align; r = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = col or GRAPHITE
    return tb

def base(prs, headline, sub=None, footer=None):
    sl = prs.slides.add_slide(blank(prs))
    R(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    R(sl, 0, 0, 13.33, 1.05, fill=NAVY)
    T(sl, headline, 0.3, 0.1, 12.7, 0.72, sz=17, bold=True, col=WHITE, font="Georgia")
    if sub: T(sl, sub, 0.3, 0.76, 12.7, 0.25, sz=8, col=LIGHT_BG)
    R(sl, 0, 1.05, 13.33, 0.02, fill=BORDER)
    R(sl, 0, 7.15, 13.33, 0.35, fill=LIGHT_BG)
    ft = footer or "Documento Confidencial — Uso Interno | Manual de Analise de Credito v2025.1 | Repositorio de Premissas Base v1.0 (Mar/2026)"
    T(sl, ft, 0.3, 7.17, 12.5, 0.18, sz=7, col=MUTED, italic=True)
    return sl

def sec(sl, text, l, t, w, col=NAVY):
    R(sl, l, t, w, 0.23, fill=col)
    T(sl, text, l+0.08, t+0.03, w-0.12, 0.18, sz=8, bold=True, col=WHITE)

def TH(sl, cols, l, t, cw, rh=0.26):
    for j, c in enumerate(cols):
        R(sl, l+j*cw, t, cw-0.02, rh, fill=NAVY)
        T(sl, c, l+j*cw+0.04, t+0.04, cw-0.08, rh-0.08, sz=7.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

def TR(sl, vals, l, t, cw, rh=0.26, bg=WHITE, tc=GRAPHITE, bold=False):
    for j, v in enumerate(vals):
        R(sl, l+j*cw, t, cw-0.02, rh, fill=bg)
        T(sl, v, l+j*cw+0.04, t+0.03, cw-0.08, rh-0.06, sz=8, bold=bold, col=tc, align=PP_ALIGN.CENTER)

def KPI(sl, l, t, w, h, val, label, col=NAVY):
    R(sl, l, t, w, h, fill=LIGHT_BG)
    T(sl, val, l+0.07, t+0.06, w-0.14, h*0.52, sz=18, bold=True, col=col, align=PP_ALIGN.CENTER)
    T(sl, label, l+0.07, t+h*0.58, w-0.14, h*0.38, sz=7.5, col=MUTED, align=PP_ALIGN.CENTER)

def flag_row(sl, icon, ico_col, text, detail, l, t, w):
    R(sl, l, t, 0.28, 0.28, fill=ico_col)
    T(sl, icon, l, t+0.02, 0.28, 0.24, sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    T(sl, text, l+0.32, t+0.02, w*0.45, 0.22, sz=8.5, bold=True, col=GRAPHITE)
    T(sl, detail, l+0.32+w*0.45+0.1, t+0.02, w*0.5, 0.22, sz=8, col=MUTED, italic=True)

# =============================================================================
# DCM — MEMORANDO DE CREDITO CORPORATIVO (10 slides)
# Cobertura: Manual_Analise_Credito_v4_1.docx — todos os pontos
#            Modelo_Analise_Credito_DCM.xlsx — 7 abas
# =============================================================================
def build_dcm():
    prs = new_prs()
    FT_DCM = "Memorando de Credito | Manual Analise Credito v2025.1 (7 Etapas) | Repositorio Premissas Base v1.0"

    # 1 — CAPA ----------------------------------------------------------------
    sl = prs.slides.add_slide(blank(prs))
    R(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    R(sl, 0, 5.3, 13.33, 2.2, fill=NAVY_MED)
    R(sl, 0.32, 0.9, 0.07, 4.5, fill=GOLD)
    T(sl, "MEMORANDO DE CRÉDITO", 0.55, 0.9, 11, 0.65, sz=28, bold=True, col=WHITE, font="Georgia")
    T(sl, "Analise de Credito Corporativo — Padrao BTG Pactual | Itau BBA", 0.55, 1.6, 9, 0.36, sz=14, col=GOLD, font="Georgia", italic=True)
    R(sl, 0.55, 2.1, 5, 0.03, fill=GOLD)
    T(sl, "[NOME DA EMPRESA S.A.] | CNPJ: [XX.XXX.XXX/0001-XX]", 0.55, 2.28, 10, 0.36, sz=13, bold=True, col=WHITE)
    T(sl, "Produto: [CCB / Debenture / CRI / CRA / NCE]  |  Volume: R$[___]M  |  Prazo: [___] meses  |  Taxa: [Indexador]+[___]% a.a.", 0.55, 2.72, 10, 0.3, sz=10, col=LIGHT_BG)
    T(sl, "Setor: [___]  |  Rating Interno: [A/B/C/D/E]  |  Rating Agencia: [AA-/A+/BBB+] — [Fitch/S&P/Moody's]", 0.55, 3.08, 10, 0.28, sz=9.5, col=MUTED)
    T(sl, "Analista: ________________  |  Data: ___/___/______  |  Nivel de Comite: [Gerente Senior / Local / Ampliado / Executivo]", 0.55, 3.44, 10, 0.28, sz=9, col=MUTED)
    T(sl, "ETAPAS: 1.KYC  →  2.Qualitativa  →  3.Financeira  →  4.Estrutura  →  5.Rating  →  6.Comite  →  7.Monitoramento", 0.55, 3.88, 10, 0.28, sz=9, col=MUTED)
    T(sl, "CONFIDENCIAL — Uso Interno Restrito. Destinado exclusivamente aos membros do Comite de Credito autorizado.\nSua reproducao ou distribuicao sem autorizacao expressa da Diretoria de Credito e vedada.",
      0.32, 5.45, 12.5, 0.5, sz=8, col=MUTED, italic=True)

    # 2 — I. SUMARIO EXECUTIVO ------------------------------------------------
    sl = base(prs, "I. Sumario Executivo — Recomendacao, Big Idea e Parametros da Operacao",
              "Manual Sec. 10 — Memorando de Credito | Etapas 1-7 consolidadas", FT_DCM)
    for i, (lbl, col) in enumerate([("✓  APROVADO", POSITIVE), ("~  APROVADO COM RESSALVAS", AMBER), ("✗  REPROVADO", NEGATIVE)]):
        R(sl, 0.3+i*3.25, 1.18, 3.1, 0.38, fill=LIGHT_BG)
        R(sl, 0.3+i*3.25, 1.18, 0.06, 0.38, fill=col)
        T(sl, lbl, 0.44+i*3.25, 1.24, 2.8, 0.26, sz=9.5, bold=True, col=GRAPHITE)
    T(sl, "[ ] Marcar a recomendacao:", 9.9, 1.24, 3.1, 0.26, sz=9, col=MUTED)

    sec(sl, "BIG IDEA — TESE DE CREDITO (1 frase — Pyramid Principle)", 0.3, 1.65, 12.7)
    R(sl, 0.3, 1.9, 12.7, 0.46, fill=LIGHT_BG)
    T(sl, "[Empresa lider em [setor] com EBITDA de R$[X]M e DSCR de [X]x — operacao de [instrumento] de R$[X]M a [indexador+spread]% a.a., com garantias de [X]x de cobertura, alinhada ao apetite de risco da instituicao.]",
      0.45, 1.94, 12.4, 0.4, sz=9.5, col=NAVY, italic=True)

    sec(sl, "DADOS DA OPERACAO", 0.3, 2.44, 12.7)
    fields = [("Empresa / CNPJ","[Nome] | [CNPJ]","Rating Interno","[A/B/C]"),
              ("Instrumento","[CCB/Deb/CRI/CRA/NCE]","Rating Agencia","[AA-] — [Fitch]"),
              ("Volume","R$ [___]M","Prazo / Carencia","[___]m / [___]m"),
              ("Indexador + Spread","[CDI/IPCA/PRE]+[X,XX]%","Amortizacao","[Bullet/SAC/PRICE]"),
              ("Garantias","[Tipo] — [X,X]x cobertura","Finalidade","[Uso exato dos recursos]"),
              ("Nivel Comite","[Gerente/Local/Ampliado/Exec]","Exposicao Total Grupo","R$ [___]M"),]
    for i, (k1,v1,k2,v2) in enumerate(fields):
        y = 2.72+i*0.4
        R(sl, 0.3, y, 2.8, 0.36, fill=LIGHT_BG); T(sl, k1, 0.38, y+0.07, 2.65, 0.24, sz=7.5, bold=True, col=NAVY)
        R(sl, 3.13, y, 3.0, 0.36, fill=WHITE);   T(sl, v1, 3.21, y+0.07, 2.88, 0.24, sz=8.5, col=GRAPHITE)
        R(sl, 6.3, y, 2.8, 0.36, fill=LIGHT_BG); T(sl, k2, 6.38, y+0.07, 2.65, 0.24, sz=7.5, bold=True, col=NAVY)
        R(sl, 9.13, y, 3.87, 0.36, fill=WHITE);  T(sl, v2, 9.21, y+0.07, 3.72, 0.24, sz=8.5, col=GRAPHITE)

    sec(sl, "PRINCIPAIS RISCOS E MITIGANTES (3 prioritarios)", 0.3, 5.15, 12.7)
    for i, (r, m) in enumerate([("Risco 1: [Nome]","[Mitigante especifico]"),("Risco 2: [Nome]","[Mitigante especifico]"),("Risco 3: [Nome]","[Mitigante especifico]")]):
        T(sl, f"  ▸ {r}", 0.3, 5.43+i*0.38, 5.5, 0.34, sz=9, bold=True, col=NEGATIVE)
        T(sl, f"→ {m}", 5.85, 5.43+i*0.38, 7.15, 0.34, sz=9, col=POSITIVE)

    # 3 — II. KYC + 5 Cs + ESG -----------------------------------------------
    sl = base(prs, "II. Triagem KYC, 5 Cs do Credito, Analise de Gestao e ESG",
              "Manual Etapas 1 (KYC) e 2 (Qualitativa) | Framework: 5 Cs + Porter + ESG + Governanca", FT_DCM)

    sec(sl, "CHECKLIST KYC — ETAPA 1 (ELIMINATORIA)", 0.3, 1.18, 6.2)
    kyc = [("Empresa ativa min. 3 anos","Fundacao: [___]","[ ]"),
           ("Faturamento > R$30MM","Receita LTM: R$[___]M","[ ]"),
           ("Sem falencia/recuperacao judicial","Confirmado: JUCEA/JUCESP","[ ]"),
           ("Nao consta COAF/OFAC/ONU/BACEN","Checado em: [data]","[ ]"),
           ("Balanco auditado 3 exercicios","Auditor: [___] | CVM: [reg]","[ ]"),
           ("Documentacao KYC completa","UBO identificado: [Sim/Nao]","[ ]"),]
    for i, (item, detalhe, status) in enumerate(kyc):
        y = 1.45+i*0.36
        R(sl, 0.3, y, 0.28, 0.28, fill=POSITIVE)
        T(sl, status, 0.3, y+0.03, 0.28, 0.22, sz=9, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        T(sl, item, 0.64, y+0.04, 3.2, 0.22, sz=9, bold=True, col=NAVY)
        T(sl, detalhe, 3.88, y+0.04, 2.45, 0.22, sz=8.5, col=MUTED)

    sec(sl, "5 Cs DO CREDITO — ETAPA 2", 6.65, 1.18, 6.38)
    cinco = [("Carater","Historico pagamentos, reputacao, compliance, referencias","[ ] Alto  [ ] Medio  [ ] Baixo"),
             ("Capacidade","FCO, EBITDA, DSCR — capacidade de servir a divida","DSCR: [__]x  |  ICR: [__]x"),
             ("Capital","Alavancagem, PL, fontes alternativas, subordinacao","Div/EBITDA: [__]x  |  PL: R$[__]M"),
             ("Colateral","Qualidade, liquidez, executabilidade das garantias","Cobertura: [__]x  |  LTV: [__]%"),
             ("Condicoes","Macro, ciclo setorial, sazonalidade, regulatorio","[ ] Fav  [ ] Neutro  [ ] Adverso"),]
    for i, (c, desc, val) in enumerate(cinco):
        y = 1.45+i*0.68
        R(sl, 6.65, y, 0.6, 0.58, fill=NAVY)
        T(sl, c[0], 6.65, y+0.1, 0.6, 0.38, sz=18, bold=True, col=GOLD, align=PP_ALIGN.CENTER, font="Georgia")
        T(sl, c, 7.3, y+0.04, 2.0, 0.22, sz=8.5, bold=True, col=NAVY)
        T(sl, desc, 7.3, y+0.27, 2.5, 0.3, sz=8, col=GRAPHITE)
        R(sl, 9.85, y, 3.15, 0.58, fill=LIGHT_BG)
        T(sl, val, 9.95, y+0.15, 2.95, 0.28, sz=8.5, bold=True, col=NAVY, align=PP_ALIGN.CENTER)

    sec(sl, "ANALISE ESG (Manual Sec. 4.4) — Rating ESG Interno", 0.3, 3.66, 6.2)
    esg = [("Ambiental","Passivos ambientais, licencamentos, riscos climaticos/transicao energetica","[ ] OK  [ ] Pendencia  [ ] Red Flag"),
           ("Social","Relacoes trabalhistas, comunidade, diversidade, condicoes laborais","[ ] OK  [ ] Pendencia  [ ] Red Flag"),
           ("Governanca","Controles internos, Lei 12.846 anticorrupcao, auditoria interna, LGPD","[ ] OK  [ ] Pendencia  [ ] Red Flag"),]
    for i, (e, d, s) in enumerate(esg):
        y = 3.93+i*0.42
        R(sl, 0.3, y, 0.75, 0.36, fill=POSITIVE)
        T(sl, e, 0.3, y+0.07, 0.75, 0.22, sz=7, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
        T(sl, d, 1.1, y+0.07, 4.0, 0.24, sz=8.5, col=GRAPHITE)
        R(sl, 5.14, y, 1.33, 0.36, fill=LIGHT_BG)
        T(sl, s, 5.18, y+0.07, 1.27, 0.22, sz=7, col=NAVY, align=PP_ALIGN.CENTER)

    sec(sl, "ANALISE DE GESTAO E GOVERNANCA CORPORATIVA (Manual Sec. 4.3)", 0.3, 5.24, 6.2)
    for i, g in enumerate(["CEO/CFO: track record, tempo no cargo — [___] anos  |  Formacao: [___]",
                            "Conselho: [__] membros, [__]% independentes — [Adequado/Insuficiente para RCVM]",
                            "Risco pessoa-chave: [Alto/Medio/Baixo]  |  Plano de sucessao: [Sim/Nao]",
                            "Plano estrategico alinhado a geracao de caixa: [Sim/Nao/Parcial]"]):
        T(sl, f"  ▸  {g}", 0.3, 5.5+i*0.36, 5.9, 0.33, sz=9, col=GRAPHITE)

    sec(sl, "ANALISE DE INDUSTRIA — 5 FORCAS DE PORTER (Manual Sec. 4.2.1)", 6.65, 3.66, 6.38)
    porter = [("Rivalidade setorial","Concentracao, guerra de precos","[ ] Alta  [ ] Media  [ ] Baixa"),
              ("Poder fornecedores","Insumos estrategicos, repasse","[ ] Alto  [ ] Medio  [ ] Baixo"),
              ("Poder clientes","Concentracao, churn, contratos","[ ] Alto  [ ] Medio  [ ] Baixo"),
              ("Novos entrantes","Barreiras regulatorias, CAPEX","[ ] Alta  [ ] Media  [ ] Baixa"),
              ("Substitutos","Disrupcao tecnologica, substituicao","[ ] Alto  [ ] Medio  [ ] Baixo"),]
    for i, (forca, desc, nota) in enumerate(porter):
        y = 3.93+i*0.36
        T(sl, forca, 6.65, y+0.05, 1.9, 0.22, sz=8.5, bold=True, col=NAVY)
        T(sl, desc, 8.6, y+0.05, 2.1, 0.22, sz=8, col=GRAPHITE)
        R(sl, 10.75, y, 2.25, 0.32, fill=LIGHT_BG)
        T(sl, nota, 10.82, y+0.05, 2.1, 0.22, sz=7.5, col=NAVY)

    # 4 — III. ANALISE FINANCEIRA — DRE E INDICES ----------------------------
    sl = base(prs, "III. Analise Financeira — DRE, Balanco, Fluxo de Caixa e Indices de Credito",
              "Manual Etapa 3 | Modelo Excel: Abas 1.PREMISSAS → 2.DRE → 3.BALANCO → 4.FLUXO → 5.DIVIDA → 6.INDICES", FT_DCM)

    sec(sl, "PREMISSAS MACRO E OPERACIONAIS (Aba 1 — Repositorio Base v1.0)", 0.3, 1.18, 12.7)
    pmacro = [("CDI (% a.a.)","Fonte: BACEN Focus","[___]%","[___]%","[___]%","[___]%"),
              ("IPCA (% a.a.)","Fonte: BACEN Focus","[___]%","[___]%","[___]%","[___]%"),
              ("USD/BRL","Fonte: BACEN Focus","[___]","[___]","[___]","[___]"),
              ("Cresc. Receita","Premissa central","[___]%","[___]%","[___]%","[___]%"),
              ("Margem EBITDA","Calculado DRE","[___]%","[___]%","[___]%","[___]%"),]
    hcols = ["Premissa","Fonte","2022A","2023A","2024A","2025E"]
    cw_p = 12.7/len(hcols)
    TH(sl, hcols, 0.3, 1.44, cw_p, rh=0.24)
    for i, row in enumerate(pmacro):
        TR(sl, row, 0.3, 1.68+i*0.28, cw_p, rh=0.26, bg=LIGHT_BG if i%2==0 else WHITE)

    sec(sl, "DRE CONSOLIDADA (Aba 2) + INDICES DE CREDITO (Aba 6)", 0.3, 3.2, 12.7)
    anos = ["2022A","2023A","2024A","LTM","2025E","CAGR"]
    cw_d = 12.7/(len(anos)+1)
    TH(sl, ["Indicador (R$M)"]+anos, 0.3, 3.46, cw_d)
    dre_rows = [
        (["Receita Liquida","[__]","[__]","[__]","[__]","[__]","[__]%"], WHITE, GRAPHITE, False),
        (["Margem Bruta %","[__]%","[__]%","[__]%","[__]%","[__]%","+[_]pp"], LIGHT_BG, GRAPHITE, False),
        (["EBITDA Ajustado","[__]","[__]","[__]","[__]","[__]","[__]%"], RGBColor(0xE8,0xF0,0xFF), NAVY, True),
        (["Margem EBITDA %","[__]%","[__]%","[__]%","[__]%","[__]%","+[_]pp"], RGBColor(0xE8,0xF0,0xFF), POSITIVE, False),
        (["Resultado Financeiro","([__])","([__])","([__])","([__])","([__])","—"], WHITE, NEGATIVE, False),
        (["Lucro Liquido","[__]","[__]","[__]","[__]","[__]","[__]%"], WHITE, GRAPHITE, False),
        (["FCO (Fluxo Caixa Op.)","[__]","[__]","[__]","[__]","[__]","—"], LIGHT_BG, GRAPHITE, False),
        (["FCO / EBITDA","[__]%","[__]%","[__]%","[__]%","[__]%","min 60%"], LIGHT_BG, POSITIVE, False),
        (["Div.Liq./EBITDA","[__]x","[__]x","[__]x","[__]x","[__]x","max [__]x"], WHITE, GRAPHITE, False),
        (["DSCR (cob.servico div.)","[__]x","[__]x","[__]x","[__]x","[__]x","min 1,5x"], WHITE, NAVY, True),
        (["ICR (EBITDA/Juros)","[__]x","[__]x","[__]x","[__]x","[__]x","min 2,0x"], WHITE, GRAPHITE, False),
    ]
    for i, (row, bg, tc, bd) in enumerate(dre_rows):
        TR(sl, row, 0.3, 3.72+i*0.32, cw_d, rh=0.3, bg=bg, tc=tc, bold=bd)

    T(sl, "Fontes: DFs auditadas [Auditor] (20XX-20XX) | LTM: periodo intermediario + trailing 12M | EBITDA reconstruido pelo analista independentemente | Covenants: Aba 6 do Modelo Excel",
      0.3, 6.96, 12.5, 0.18, sz=7, col=MUTED, italic=True)

    # 5 — III-B. ANALISE FORENSE DRE (Guia Perfeccionista Sec.12.1) ----------
    sl = base(prs, "III-B. Analise Forense da DRE — Red e Amber Flags (Manual Sec. 12.1)",
              "Guia Perfeccionista: 'Numeros nao mentem — mas podem ser arranjados. Leia as NOTAS EXPLICATIVAS antes do balanco.'", FT_DCM)

    T(sl, "Principio Fundamental: Cruzar DRE x DFC x BP. Se FCO/EBITDA < 60%, ha problema de qualidade de lucro. Recalcular EBITDA do zero a partir do EBIT auditado.",
      0.3, 1.18, 12.7, 0.26, sz=9, col=NAVY, bold=True)

    flags_dre = [
        (NEGATIVE,"🔴","RECEITA — Qualidade e Concentracao",
         "Cresc. receita x cresc. FCO | Concentracao clientes | PMR tendencia",
         "FCO negativo com receita crescendo | Cliente unico >30% | Mudanca de politica de reconhecimento sem justificativa"),
        (NEGATIVE,"🔴","CPV/CMV — Margem Bruta e Pressao de Custos",
         "Margem bruta historica e tendencia vs. setor",
         "Margem caindo >3pp em 2 trimestres | CPV crescendo mais rapido que receita | Margem muito acima do setor (agressividade contabil)"),
        (AMBER,"🟠","SG&A — Alavancagem Operacional",
         "SG&A/Receita historico | Crescimento absoluto vs. receita",
         "SG&A >25% receita em distribuicao/atacado | Crescendo >15%aa sem expansao | Itens 'nao recorrentes' que se repetem"),
        (NEGATIVE,"🔴","EBITDA — Ajustes e Qualidade",
         "EBITDA reportado vs. EBITDA caixa | Natureza dos ajustes",
         "Ajustado >20% vs. IFRS | Equiv. patrim. positiva incluida | Ajustes recorrentes >2 anos | 'Outros' >15% do EBITDA"),
        (NEGATIVE,"🔴","RESULTADO FINANCEIRO — Custo da Divida",
         "Desp. financeira/EBITDA | Composicao (juros, cambio, multas, IOF)",
         "Taxa implicita (Desp.Fin./Div.Bruta) muito acima do CDI+spread negociado | Multas/juros mora crescentes | IOF desproporcional"),
        (AMBER,"🟠","IR/CSLL — Aliquota Efetiva e Diferidos",
         "Aliquota efetiva vs. 34% | Ativo/Passivo fiscal diferido",
         "Aliquota <15% sem explicacao clara | Ativo diferido >20% PL sem perspectiva de realizacao | Planejamento fiscal agressivo"),
        (NEGATIVE,"🔴","QUALIDADE DO LUCRO — FCO vs. Lucro Liquido",
         "IQL = FCO/LL | Accruals (LL - FCO) | Dividendos vs. FCO",
         "IQL < 0,7x por 2 anos | Accruals positivos e crescentes | LL positivo com FCO negativo | Dividendos pagos com divida"),
    ]
    for i, (ico_col, icon, titulo, analise, alert) in enumerate(flags_dre):
        col_n = i % 2
        row_n = i // 2
        if i == 6: col_n = 0; row_n = 3
        l = 0.3 + col_n * 6.55
        t = 1.5 + row_n * 1.38
        w = 6.2
        R(sl, l, t, w, 1.28, fill=RED_BG if ico_col == NEGATIVE else AMB_BG)
        R(sl, l, t, w, 0.26, fill=ico_col)
        T(sl, f"{icon} {titulo}", l+0.08, t+0.04, w-0.15, 0.2, sz=8, bold=True, col=WHITE)
        T(sl, f"O que analisar: {analise}", l+0.08, t+0.3, w-0.15, 0.3, sz=8, col=GRAPHITE)
        T(sl, f"Sinal de alerta: {alert}", l+0.08, t+0.63, w-0.15, 0.55, sz=7.5, col=ico_col, bold=True)

    # 6 — III-C. ANALISE FORENSE BALANCO (Guia Sec.12.2) ----------------------
    sl = base(prs, "III-C. Analise Forense do Balanco Patrimonial — Red e Amber Flags (Manual Sec. 12.2)",
              "Principio: verificar restricoes do caixa, aging de recebiveis, obsolescencia de estoques, CAPEX vs. depreciacao e partes relacionadas", FT_DCM)

    flags_bp = [
        (NEGATIVE,"🔴","CAIXA — Qualidade e Restricoes",
         "Composicao do caixa | Restricoes de uso | Caixa minimo operacional (30-45 dias receita)",
         "Caixa elevado mas FCO negativo | Caixa restrito >30% sem disclosure | Queda abrupta entre trimestres"),
        (NEGATIVE,"🔴","CONTAS A RECEBER — Qualidade e PDD",
         "PMR historico e tendencia | Aging (vencidos >60/90/180d) | Cobertura PDD | Concentracao devedor",
         "PMR crescendo >15d/ano | PDD caindo com inadimplencia crescente | Recebíveis partes relacionadas >10%"),
        (AMBER,"🟠","ESTOQUES — Giro e Obsolescencia",
         "PME historico vs. setor | Provisao obsolescencia | Metodo avaliacao (PEPS/custo medio)",
         "PME crescendo sem crescimento de receita | Provisao zerada em desaceleracao | Write-down subsequente"),
        (NEGATIVE,"🔴","IMOBILIZADO — CAPEX e Impairment",
         "Intensidade CAPEX | Vida util dos ativos | Teste de impairment | Intangiveis nao amortizaveis",
         "CAPEX <deprec. por >2 anos (desinvestimento) | CAPEX >3x deprec. sem expansao | Goodwill >30% sem impairment"),
        (AMBER,"🟠","PARTES RELACIONADAS — Transacoes",
         "Saldos a receber/pagar | Natureza das transacoes | Condicoes de mercado (arm's length)",
         "Emprestimos a socios/coligadas | Management fee excessivo | Transacoes sem substancia economica"),
        (AMBER,"🟠","LIQUIDEZ E CAPITAL DE GIRO",
         "Liquidez corrente/seca | CCC (DSO+DIO-DPO) | Working capital vs. necessidade operacional",
         "Liquidez corrente <1,0x | CCC deteriorando sem expansao | Capital de giro negativo crescente"),
    ]
    for i, (ico_col, icon, titulo, analise, alert) in enumerate(flags_bp):
        col_n = i % 2; row_n = i // 2
        l = 0.3 + col_n * 6.55; t = 1.18 + row_n * 1.6; w = 6.2
        R(sl, l, t, w, 1.5, fill=RED_BG if ico_col == NEGATIVE else AMB_BG)
        R(sl, l, t, w, 0.26, fill=ico_col)
        T(sl, f"{icon} {titulo}", l+0.08, t+0.04, w-0.15, 0.2, sz=8, bold=True, col=WHITE)
        T(sl, f"O que analisar: {analise}", l+0.08, t+0.3, w-0.15, 0.35, sz=8, col=GRAPHITE)
        T(sl, f"Sinal de alerta: {alert}", l+0.08, t+0.7, w-0.15, 0.7, sz=7.5, col=ico_col, bold=True)

    # 7 — V. ESTRUTURA DA OPERACAO -------------------------------------------
    sl = base(prs, "V. Estrutura da Operacao — Parametros, Garantias, Covenants e Precificacao",
              "Manual Etapa 4 | Modelo Excel Aba 5.DIVIDA | Politica de Garantias e LTV | Pricing", FT_DCM)

    sec(sl, "PARAMETROS FORMAIS", 0.3, 1.18, 6.2)
    params = [("Modalidade","[CCB/NCE/Deb/CRI/CRA/FIDC/Bond]"),("Volume","R$[___]M  |  Moeda: [BRL/USD]"),
              ("Prazo","[___]m total  |  Carencia: [___]m"),("Indexador","[CDI+/IPCA+/PRE/SELIC]"),
              ("Spread","[X,XX]% a.a. — base 252 DU"),("Amortizacao","[Bullet/SAC/PRICE/Custom]"),
              ("Pagto Juros","[Mensal/Trim./Semestral]"),("Finalidade","[Descricao exata — sem desvio]"),
              ("Regime","[RCVM 160 art.3o / CVM 400]"),("Comissao","[___]% flat do volume emitido"),]
    for i, (k,v) in enumerate(params):
        y=1.45+i*0.36; R(sl,0.3,y,2.1,0.32,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,k,0.38,y+0.06,2.0,0.22,sz=7.5,bold=True,col=NAVY)
        T(sl,v,2.44,y+0.06,3.9,0.22,sz=8.5,col=GRAPHITE)

    sec(sl, "GARANTIAS E LTV (Manual Sec. 6.2 — Politica de Garantias)", 6.65, 1.18, 6.38)
    TH(sl, ["Tipo de Garantia","Valor Bruto","Haircut","Valor Liq.","Cobertura"], 6.65, 1.44, cw=6.38/5, rh=0.25)
    for i, row in enumerate([["[Recebiveis — diversificados]","R$[__]M","20%","R$[__]M","[__]x"],
                              ["[Imovel rural c/ benfeitorias]","R$[__]M","30%","R$[__]M","[__]x"],
                              ["[Equipamentos industriais]","R$[__]M","45%","R$[__]M","[__]x"],
                              ["[Aval pessoal solidario]","R$[__]M","—","R$[__]M","[__]x"],
                              ["TOTAL","R$[__]M","—","R$[__]M","[__]x"]]):
        bg = NAVY if i==4 else (LIGHT_BG if i%2==0 else WHITE)
        tc = WHITE if i==4 else GRAPHITE
        TR(sl, row, 6.65, 1.69+i*0.28, cw=6.38/5, rh=0.26, bg=bg, tc=tc, bold=(i==4))

    sec(sl, "COVENANTS FINANCEIROS (Manual Sec. 6.3)", 6.65, 3.1, 6.38)
    TH(sl, ["Covenant","Limite","Cure Period","Consequencia"], 6.65, 3.36, cw=6.38/4, rh=0.25)
    covs = [["Div.Liq./EBITDA","max [___]x","[__] dias","Venc. antecipado"],
            ["ICR (EBITDA/Juros)","min [___]x","[__] dias","Venc. antecipado"],
            ["DSCR","min [___]x","[__] dias","Cash sweep"],
            ["FCO/EBITDA","min [___]%","[__] dias","Waiver comite"],
            ["Dividendos","max [__]% LL se Div/EBITDA>[__]x","—","Restricao automatica"]]
    for i, row in enumerate(covs):
        TR(sl, row, 6.65, 3.61+i*0.3, cw=6.38/4, rh=0.28, bg=LIGHT_BG if i%2==0 else WHITE)

    sec(sl, "PRECIFICACAO (Manual Sec. 6.4)", 0.3, 5.06, 6.2)
    T(sl, "Taxa = Custo de Captacao + Custo Capital Regulatorio + PD Esperada + Margem\n"
          "Ref. NTN-B [prazo]: IPCA+[__]%  |  DI Futuro [venc.]: [__]%aa\n"
          "Spread mercado ([Rating]): [Index]+[__]%  |  Proposto: [Index]+[__]%\n"
          "Taxa implicita dívida atual: [__]% aa  |  Prêmio sobre Tesouro: [__]pp",
      0.3, 5.32, 5.9, 0.72, sz=9, col=GRAPHITE)
    sec(sl, "CROSS-SELL E HISTORICO COMERCIAL", 0.3, 6.12, 6.2)
    T(sl, "Exposicao atual grupo: R$[__]M  |  Historico pagamentos: [Pontual/Atrasos]  |  Receita cross-sell: R$[__]M/ano\n"
          "Wallet share: [__]%  |  Outras instituicoes na operacao: [Sim/Nao] — [Nomes]",
      0.3, 6.38, 5.9, 0.4, sz=9, col=GRAPHITE)

    # 8 — VI. STRESS TEST + CENARIOS (Aba 7 Sensibilidade) --------------------
    sl = base(prs, "VI. Stress Test, Cenarios e Analise de Sensibilidade — Aba 7 do Modelo Excel",
              "Manual Etapa 3 Sec. 5.4 | 4 Cenarios: Base / Otimista / Estressado / Downside Severo | Break-even + Early Warnings", FT_DCM)

    sec(sl, "PREMISSAS POR CENARIO (Repositorio Base v1.0 — Modulo 7)", 0.3, 1.18, 12.7)
    cen_h = ["Premissa","Unidade","BASE","OTIMISTA","ESTRESSADO","DOWNSIDE SEVERO"]
    cen_w = 12.7/len(cen_h)
    TH(sl, cen_h, 0.3, 1.44, cen_w)
    cen_rows = [["Crescimento Receita","%","[+8%]","[+12%]","[0%]","[-20%]"],
                ["Margem EBITDA","%","[19%]","[22%]","[16%]","[12%]"],
                ["Taxa CDI","% a.a.","[___]","[-100bps]","[+115bps]","[+335bps]"],
                ["CAPEX / Receita","%","[4%]","[3,5%]","[4,5%]","[5,5%]"],
                ["BRL/USD","R$/US$","[5,78]","[5,40]","[6,20]","[6,80]"]]
    for i, row in enumerate(cen_rows):
        TR(sl, row, 0.3, 1.7+i*0.3, cen_w, rh=0.28, bg=LIGHT_BG if i%2==0 else WHITE)

    sec(sl, "RESULTADOS POR CENARIO — METRICAS DE CREDITO", 0.3, 3.27, 12.7)
    res_h = ["Metrica","Base","Otimista","Estressado","Downside","Covenant","Break-even EBITDA"]
    res_w = 12.7/len(res_h)
    TH(sl, res_h, 0.3, 3.53, res_w)
    res_rows = [
        (["Receita (R$M)","[__]","[__]","[__]","[__]","—","—"],WHITE,GRAPHITE),
        (["EBITDA (R$M)","[__]","[__]","[__]","[__]","—","R$[__]M"],RGBColor(0xE8,0xF0,0xFF),NAVY),
        (["Div.Liq./EBITDA","[__]x","[__]x","[__]x","[__]x","max [__]x","[__]x"],WHITE,GRAPHITE),
        (["DSCR","[__]x","[__]x","[__]x","[__]x","min [__]x","[__]x"],LIGHT_BG,NAVY),
        (["ICR","[__]x","[__]x","[__]x","[__]x","min [__]x","[__]x"],WHITE,GRAPHITE),
        (["Covenants OK?","SIM","SIM","SIM","[SIM/NAO]","—","Queda max:[__]%"],WHITE,POSITIVE),
    ]
    for i, (row,bg,tc) in enumerate(res_rows):
        TR(sl, row, 0.3, 3.79+i*0.33, res_w, rh=0.31, bg=bg, tc=tc, bold=(row[0]=="DSCR"))

    sec(sl, "ANALISE DE SENSIBILIDADE — VARIAVEIS CRITICAS", 0.3, 5.77, 12.7)
    T(sl, "Break-even EBITDA (DSCR=1,0x): R$[__]M = queda de [__]% vs. LTM — margem de seguranca: [__]pp\n"
          "Sensib. Receita: queda de 10% → DSCR cai [__]x | Ponto critico: -[__]% de receita\n"
          "Sensib. Cambial: [Aplicavel/NA] — para cada R$0,10 apreciacao BRL: EBITDA cai R$[__]M\n"
          "Sensib. Juros: para cada +100bps CDI/IPCA: despesa financeira aumenta R$[__]M/ano",
      0.3, 6.02, 8.5, 0.7, sz=9.5, col=GRAPHITE)

    sec(sl, "INDICADORES EARLY WARNING (Manual Sec. 9.2)", 8.9, 5.77, 4.1)
    ews = ["Receita cai >20% vs. trimestre anterior","Margem EBITDA cai >5pp em 2 trim.",
           "Div.Liq./EBITDA sobe >1,0x sem aprovacao","Atraso pagto fornecedores >15 dias"]
    for i, e in enumerate(ews):
        T(sl, f"  ⚠  {e}", 8.9, 6.05+i*0.33, 4.0, 0.3, sz=8.5, col=AMBER)

    # 9 — VII. RATING + RECOMENDACAO -----------------------------------------
    sl = base(prs, "VII. Scorecard de Rating, Alçada de Aprovacao e Recomendacao Final",
              "Manual Etapas 5 (Scorecard 70/30) e 6 (Comite de Credito) | Escala A → E", FT_DCM)

    sec(sl, "SCORECARD DE CREDITO — 70% QUANTITATIVO + 30% QUALITATIVO (Manual Sec. 7.1)", 0.3, 1.18, 12.7)
    sc_h = ["Criterio","Peso","Nota (1-10)","Score Pond.","Benchmark","Comentario"]
    sc_w = 12.7/len(sc_h)
    TH(sl, sc_h, 0.3, 1.44, sc_w)
    sc_rows = [
        (["EBITDA / Juros (ICR)","15%","[__]","[__]",">2,5x=10","[Comentario do analista]"],LIGHT_BG),
        (["Div. Liq. / EBITDA","15%","[__]","[__]","<2,0x=10","[Comentario do analista]"],WHITE),
        (["DSCR","15%","[__]","[__]",">2,0x=10","[Comentario do analista]"],LIGHT_BG),
        (["FCO / EBITDA","10%","[__]","[__]",">75%=10","[Comentario do analista]"],WHITE),
        (["Cresc. Receita","10%","[__]","[__]",">15%aa=10","[Comentario do analista]"],LIGHT_BG),
        (["Margem EBITDA","5%","[__]","[__]","Bench.setor","[Comentario do analista]"],WHITE),
        (["Carater / Historico","10%","[__]","[__]","Sem ocorr.","[Comentario do analista]"],LIGHT_BG),
        (["Qualidade Gestao","8%","[__]","[__]","TR>5a","[Comentario do analista]"],WHITE),
        (["Posic. Mercado","7%","[__]","[__]","Top3","[Comentario do analista]"],LIGHT_BG),
        (["Garantias/LTV","5%","[__]","[__]","Cob.>1,5x=10","[Comentario do analista]"],WHITE),
        (["SCORE TOTAL","100%","—","[___]","—",""],RGBColor(0xE8,0xF0,0xFF)),
    ]
    for i, (row, bg) in enumerate(sc_rows):
        bd = row[0] == "SCORE TOTAL"
        TR(sl, row, 0.3, 1.7+i*0.34, sc_w, rh=0.32, bg=bg, tc=NAVY if bd else GRAPHITE, bold=bd)

    sec(sl, "ESCALA DE RATING E ALCADA DE APROVACAO (Manual Sec. 7.2 e 8.1)", 0.3, 5.52, 12.7)
    escalas = [("A","Score 8,5-10",POSITIVE,"Gerente Senior"),("B","Score 7,0-8,4",RGBColor(0x27,0x67,0x49),"Comite Local (3 votantes)"),
               ("C","Score 5,5-6,9",AMBER,"Comite Ampliado (Dir.+Risco)"),("D","Score 3,0-5,4",NEGATIVE,"Comite Executivo (VP+Dir.)"),("E","Score <3,0",RGBColor(0x7B,0x00,0x00),"REPROVADO")]
    for i, (r, s, c, a) in enumerate(escalas):
        R(sl, 0.3+i*2.51, 5.78, 2.41, 0.28, fill=c)
        T(sl, f"Rating {r} — {s}", 0.4+i*2.51, 5.8, 2.25, 0.22, sz=7.5, bold=True, col=WHITE)
        T(sl, a, 0.4+i*2.51, 6.1, 2.25, 0.22, sz=7.5, col=GRAPHITE)

    R(sl, 0.3, 6.42, 12.7, 0.34, fill=LIGHT_BG)
    T(sl, f"Rating calculado: [___]  |  Score: [__]/10  |  Alcada: [_______________]  |  Analista: _________________  |  Data: ___/___/______",
      0.4, 6.48, 9.5, 0.24, sz=8.5, bold=True, col=NAVY)
    T(sl, "[ ] APROVADO  [ ] APROVADO C/ RESSALVAS  [ ] REPROVADO", 10.0, 6.48, 3.0, 0.24, sz=9, bold=True, col=NAVY)

    # 10 — MONITORAMENTO (Etapa 7) -------------------------------------------
    sl = base(prs, "Etapa 7 — Monitoramento Pos-Aprovacao e Plano de Acao para Creditos Problemáticos",
              "Manual Sec. 9 — Frequencia de revisao por rating | Early Warnings | Watch List | PCLD", FT_DCM)

    sec(sl, "FREQUENCIA DE REVISAO POR RATING (Manual Sec. 9.1)", 0.3, 1.18, 7.0)
    TH(sl, ["Rating","Revisao","Monitoramento Covenants","Reuniao com Gestao"], 0.3, 1.44, cw=7.0/4)
    mon_rows = [["A","Anual","Trimestral","Se covenant em alerta"],
                ["B","Semestral","Trimestral","Semestral"],
                ["C","Trimestral","Mensal","Trimestral"],
                ["D/E (Watch List)","Mensal","Mensal","Mensal — obrigatorio"]]
    for i, row in enumerate(mon_rows):
        TR(sl, row, 0.3, 1.7+i*0.36, cw=7.0/4, rh=0.34, bg=LIGHT_BG if i%2==0 else WHITE, tc=NEGATIVE if "Watch" in row[0] else GRAPHITE)

    sec(sl, "INDICADORES DE EARLY WARNING (Manual Sec. 9.2)", 0.3, 3.23, 7.0)
    for i, (flag, nivel) in enumerate([
        ("Atraso pagamentos a fornecedores >15 dias", AMBER),
        ("Queda receita >20% vs. trimestre anterior (sem sazonalidade)", AMBER),
        ("Deterioracao EBITDA margin >5pp em 2 trimestres", NEGATIVE),
        ("Aumento Div.Liq./EBITDA >1,0x sem justificativa aprovada", NEGATIVE),
        ("Violacao ou proxied violation de qualquer covenant", NEGATIVE),
        ("Noticias negativas: investigacoes, autuacoes, perda cliente-chave", AMBER),
        ("Mudanca nao comunicada no controle acionario ou gestao senior", NEGATIVE),]):
        T(sl, f"  ⚠  {flag}", 0.3, 3.5+i*0.34, 6.7, 0.3, sz=9, col=nivel)

    sec(sl, "PLANO DE ACAO — CREDITOS PROBLEMÁTICOS (Manual Sec. 9.3)", 7.5, 1.18, 5.55)
    T(sl, "Creditos Watch List ou Rating D/E exigem Plano de Acao em 5 dias uteis da deteccao:", 7.5, 1.45, 5.45, 0.26, sz=8.5, col=NAVY, bold=True)
    plano = ["(i) Diagnostico da deterioracao — causa raiz identificada",
             "(ii) Proposta de reestruturacao ou execucao de garantias",
             "(iii) Estimativa de perda esperada: LGD / PD / EAD",
             "(iv) Provisao PCLD recomendada — alinhada ao BACEN 2682",
             "(v) Prazo de resolucao e proximo checkpoint com Diretoria"]
    for i, p in enumerate(plano):
        T(sl, f"  {p}", 7.5, 1.75+i*0.4, 5.45, 0.36, sz=9, col=GRAPHITE)

    sec(sl, "DOCUMENTACAO OBRIGATORIA PARA O COMITE (Manual Sec. 8.2)", 7.5, 3.9, 5.55)
    docs = ["Memorando de Credito completo (este documento)",
            "Scorecard de Rating com fundamentacao individual",
            "Laudos de avaliacao das garantias (max. 90 dias)",
            "Parecer juridico sobre garantias e estrutura",
            "Relatorio de compliance e KYC atualizado",
            "Modelo financeiro Excel com 3 cenarios + sensibilidade",
            "Ata do Comite com votacao e voto dissidente fundamentado"]
    for i, d in enumerate(docs):
        T(sl, f"  ✓  {d}", 7.5, 4.18+i*0.38, 5.45, 0.34, sz=9, col=GRAPHITE)

    sec(sl, "CLAUSULAS DE MONITORAMENTO OBRIGATORIAS", 0.3, 5.32, 7.0)
    cls = ["Entrega de DFs auditadas em ate [__] dias apos encerramento do exercicio",
           "Entrega de balancetes trimestrais em [__] dias apos encerramento do trimestre",
           "Comunicado imediato de qualquer fato relevante material",
           "Certificado de conformidade de covenants assinado pelo CFO — trimestral",
           "Autorizacao previa para: (i) alienacao de ativos >R$[__]M, (ii) novo endividamento >R$[__]M, (iii) mudanca de controle"]
    for i, c in enumerate(cls):
        T(sl, f"  ▸  {c}", 0.3, 5.6+i*0.32, 6.8, 0.3, sz=9, col=GRAPHITE)

    path = str(Path(OUT_DIR) / "parecer_credito_dcm.pptx")
    prs.save(path)
    print(f"DCM: {path} ({len(prs.slides)} slides)")
    return path


# =============================================================================
# ECM — EQUITY OPINION (11 slides)
# Cobertura: Manual_ECM_IB.docx — Secoes 1-13 completas
#            Modelo_ECM_IB.xlsx — 11 abas
# =============================================================================
def build_ecm():
    prs = new_prs()
    FT_ECM = "Equity Opinion ECM | Manual ECM IB v1.0 (Secoes 1-13) | Repositorio Premissas Base v1.0 (Mar/2026)"

    # 1 — CAPA ----------------------------------------------------------------
    sl = prs.slides.add_slide(blank(prs))
    R(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    R(sl, 0, 5.2, 13.33, 2.3, fill=NAVY_MED)
    R(sl, 0.32, 0.9, 0.07, 4.2, fill=GOLD)
    T(sl, "EQUITY OPINION", 0.55, 0.9, 10, 0.62, sz=28, bold=True, col=WHITE, font="Georgia")
    T(sl, "Analise de Valuation e Estrutura de Oferta — ECM", 0.55, 1.58, 9, 0.35, sz=14, col=GOLD, font="Georgia", italic=True)
    R(sl, 0.55, 2.05, 4.8, 0.03, fill=GOLD)
    T(sl, "[NOME DA EMPRESA] | [IPO / FOLLOW-ON / BLOCK TRADE]", 0.55, 2.22, 10, 0.38, sz=13, bold=True, col=WHITE)
    T(sl, "Faixa: R$[__]—R$[__] por acao  |  Equity Value: R$[__]Bi  |  EV/EBITDA NTM: [__]x  |  Free Float: [__]%", 0.55, 2.68, 10, 0.3, sz=10, col=LIGHT_BG)
    T(sl, "Segmento: [Novo Mercado/Nivel 2]  |  Regime: [RCVM 160 art.3o / CVM 400]  |  Coord.: [Nome]", 0.55, 3.04, 10, 0.28, sz=9.5, col=MUTED)
    T(sl, "Analista: ________________  |  MD: ________________  |  Data: ___/___/______  |  Status: [Pre-Marketing/Bookbuilding]", 0.55, 3.4, 10, 0.28, sz=9, col=MUTED)
    T(sl, "COBERTURA: Sec.1 Intro | Sec.2 Participantes | Sec.3 Tipos | Sec.4 Processo 8 Fases | Sec.5 Valuation | Sec.6 Pricing | Sec.7 DD | Sec.8 Regulacao | Sec.9 Modelo | Sec.10 Risco | Sec.11 Otimizacao | Sec.12 Flags", 0.55, 3.84, 10, 0.28, sz=8.5, col=MUTED)
    T(sl, "CONFIDENCIAL — Destinado a membros do Comite de Pricing e Investidores Qualificados (RCVM 160/2022).\nNao constitui oferta publica de valores mobiliarios.",
      0.32, 5.35, 12.5, 0.5, sz=8, col=MUTED, italic=True)

    # 2 — I. SUMARIO + FOOTBALL FIELD -----------------------------------------
    sl = base(prs, "I. Sumario Executivo — Equity Story, Football Field e Recomendacao",
              "Manual Sec. 5.1 e 5.4 (Football Field) + Sec. 6.1 (Faixa de Preco) | Pyramid Principle", FT_ECM)

    T(sl, "EQUITY STORY — BIG IDEA (Pyramid Principle — 1 frase conclusiva):", 0.3, 1.18, 12.7, 0.22, sz=8.5, bold=True, col=NAVY)
    R(sl, 0.3, 1.42, 12.7, 0.44, fill=LIGHT_BG)
    T(sl, "[Empresa lider em [setor] com CAGR de [X]%, margem EBITDA de [X]% e share de [X]% em mercado de R$[X]Bi — oferta de R$[X]Bi ao EV/EBITDA NTM de [X]x com desconto de [X]% vs. peers, gerando upside de [X]% ao preco justo.]",
      0.45, 1.46, 12.4, 0.38, sz=9.5, col=NAVY, italic=True)

    sec(sl, "FOOTBALL FIELD — FAIXA DE VALUATION POR METODO (Manual Sec. 5.4 | Aba 6-7 Excel)", 0.3, 1.95, 7.2)
    metodos = [("DCF (WACC [__]%, g [__]%, TV [__]%)","Base",NAVY),
               ("EV/EBITDA Peers LTM ([__]x P25-P75)","Relativo",NAVY_MED),
               ("EV/EBITDA Peers NTM ([__]x P25-P75)","Relativo",NAVY_MED),
               ("P/E Peers NTM ([__]x P25-P75)","Relativo",GRAPHITE),
               ("Transacoes Precedentes ([__]x mediana)","Transacao",MUTED),
               ("DDM — Dividend Discount Model (g=[__]%)","Dividendo",MUTED),]
    scl_min, scl_max = 2.0, 11.0; scale_w = 6.5
    for i, (label, tipo, col) in enumerate(metodos):
        y = 2.24+i*0.46
        T(sl, label, 0.3, y+0.1, 2.85, 0.28, sz=7.5, col=GRAPHITE)
        T(sl, tipo, 3.18, y+0.1, 0.75, 0.28, sz=7.5, col=MUTED)
        bar_l = 4.0 + (3.0-scl_min)/(scl_max-scl_min)*scale_w
        bar_w = 1.8/(scl_max-scl_min)*scale_w
        R(sl, bar_l, y+0.06, bar_w, 0.3, fill=col)
        T(sl, "R$[min]", bar_l-0.5, y+0.1, 0.48, 0.24, sz=7, col=MUTED, align=PP_ALIGN.RIGHT)
        T(sl, "R$[max]", bar_l+bar_w+0.03, y+0.1, 0.48, 0.24, sz=7, col=MUTED)
    R(sl, 4.0+(3.5-scl_min)/(scl_max-scl_min)*scale_w, 2.2, 0.05, 2.9, fill=GOLD)
    T(sl, "Faixa proposta", 4.0+(3.5-scl_min)/(scl_max-scl_min)*scale_w-0.6, 5.18, 0.65, 0.3, sz=7, col=GOLD, bold=True)

    sec(sl, "PARAMETROS DA OFERTA (Aba 8 Excel)", 7.7, 1.95, 5.3)
    oferta_s = [("Faixa Indicativa","R$[__] — R$[__] por acao"),("Equity Value","R$[__]Bi — R$[__]Bi"),
                ("Enterprise Value","R$[__]Bi — R$[__]Bi"),("EV/EBITDA NTM","[__]x — [__]x | Peers: [__]x"),
                ("P/E NTM","[__]x — [__]x | Peers: [__]x"),("Upside vs. Fair Value","[__]% — [__]%"),
                ("Volume Base","R$[__]Bi | c/ Greenshoe: R$[__]Bi"),("Free Float pos-oferta","[__]% base | [__]% c/ greenshoe"),
                ("Oferta Primaria","R$[__]M — [__]M acoes novas"),("Oferta Secundaria","R$[__]M — [__]M acoes vendedores"),]
    for i, (k,v) in enumerate(oferta_s):
        y=2.24+i*0.38; R(sl,7.7,y,2.1,0.34,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,k,7.78,y+0.06,2.0,0.24,sz=7.5,bold=True,col=NAVY); T(sl,v,9.83,y+0.06,3.12,0.24,sz=8.5,col=GRAPHITE)

    sec(sl, "3 RAZOES PARA INVESTIR (Pyramid Principle)", 0.3, 5.92, 12.7)
    for i, r in enumerate(["[Razao 1 — Headline conclusivo: ex: 'CAGR de 22% sustentado por expansao contratada — ROI de 28%aa']",
                            "[Razao 2 — Headline conclusivo: ex: 'Margem EBITDA de 28% acima do peer group — pricing power comprovado em ciclos adversos']",
                            "[Razao 3 — Headline conclusivo: ex: 'Valuation com desconto de 15% vs. peers — upside de 35% ao preco justo do DCF']"]):
        R(sl,0.3,6.2+i*0.43,0.32,0.32,fill=GOLD)
        T(sl,str(i+1),0.3,6.2+i*0.43,0.32,0.32,sz=14,bold=True,col=NAVY,align=PP_ALIGN.CENTER,font="Georgia")
        T(sl,r,0.7,6.25+i*0.43,12.2,0.3,sz=9,col=GRAPHITE)

    # 3 — II. PARTICIPANTES + TIPOS DE OFERTA --------------------------------
    sl = base(prs, "II. Participantes da Operacao e Tipos de Oferta ECM",
              "Manual Sec. 2 (Mapa de Agentes) + Sec. 3 (Taxonomia Completa) + Sec. 8 (Regulacao B3)", FT_ECM)

    sec(sl, "MAPA DE PARTICIPANTES (Manual Sec. 2.1)", 0.3, 1.18, 6.2)
    partic = [("Emissor / Ofertante","[Nome da Empresa]  |  Acionistas vendedores: [Nomes]"),
              ("Coordenador Lider","[Nome do Banco]  |  Mandato: [Exclusivo/Conjunto]"),
              ("Coordenadores Consorciados","[Nome 1], [Nome 2]  |  Papel: co-distribuicao"),
              ("Agente Estabilizador","[Nome]  |  Greenshoe: ate [__]% | Prazo: 30 dias"),
              ("Escriturador / Custodiante","[Nome]  |  CNPJ: [___]"),
              ("Auditor Independente","[Nome]  |  Registro CVM: [___]"),
              ("Advogados","[Empresa X] (emissor) | [Empresa Y] (bancos)"),
              ("CVM / B3","Registro: [numero]  |  Segmento: [Novo Mercado]  |  Ticker: [____]"),]
    for i, (k,v) in enumerate(partic):
        y=1.45+i*0.36; R(sl,0.3,y,2.0,0.32,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,k,0.38,y+0.06,1.88,0.22,sz=7.5,bold=True,col=NAVY); T(sl,v,2.34,y+0.06,3.78,0.22,sz=8.5,col=GRAPHITE)

    sec(sl, "TIPO DE OPERACAO (Manual Sec. 3)", 6.65, 1.18, 6.38)
    tipos = [("[ ]","IPO — Initial Public Offering","1a oferta de acoes | Oferta primaria + secundaria | RCVM 160 registro CVM"),
             ("[ ]","Follow-on — Oferta Subsequente","ICVM 476 (restrita, <5d) ou ICVM 160 ampla (>5d) | Acao ja listada"),
             ("[ ]","Block Trade — Aceleracao de Livro","Overnight, sem prospecto completo | Desconto 3-7% VWAP | Risco banco"),
             ("[ ]","Greenshoe — Lote Suplementar","Ate 15% do lote base | Prazo 30 dias | Mecanismo de estabilizacao"),
             ("[ ]","Deb. Conversiveis / CRI-CRA c/ Part.","Instrumento hibrido | Modelar diluicao e impacto no EPS"),]
    for i, (chk, nome, desc) in enumerate(tipos):
        y=1.45+i*0.72
        R(sl,6.65,y,6.38,0.66,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,chk,6.72,y+0.2,0.3,0.28,sz=11,col=NAVY)
        T(sl,nome,7.07,y+0.06,2.5,0.24,sz=9,bold=True,col=NAVY)
        T(sl,desc,7.07,y+0.33,5.9,0.28,sz=8.5,col=GRAPHITE)

    sec(sl, "REGULACAO E SEGMENTOS B3 (Manual Sec. 8)", 0.3, 4.87, 12.7)
    T(sl, "Marco regulatorio principal: RCVM 160/2022 | Lei 6.385/76 | Lei 6.404/76 | ANBIMA Codigo de Ofertas", 0.3, 5.13, 12.5, 0.22, sz=8.5, col=NAVY, bold=True)
    segs = [("Novo Mercado","Padrao institucional maximo","So acoes ON | FF min 25% | Tag-along 100% | Min 5 conselheiros (20% indep.)"),
            ("Nivel 2","Acoes ON e PN com direitos","Tag-along 100% ON+PN | Arbitragem obrigatorio | Maior flexibilidade"),
            ("Nivel 1","Disclosure aprimorado","FF min 25% | Reunioes publicas anuais | Liquidez minima"),
            ("EGEM (Emerging)","Empresas menores","FF min 15% | Procedimento simplificado | Isentos de alguns requisitos"),]
    for i, (seg, label, req) in enumerate(segs):
        R(sl,0.3+i*3.2,5.38,3.1,0.26,fill=NAVY)
        T(sl,f"{seg} — {label}",0.4+i*3.2,5.4,2.9,0.22,sz=7.5,bold=True,col=WHITE)
        T(sl,req,0.4+i*3.2,5.68,2.9,0.42,sz=7.5,col=GRAPHITE)

    sec(sl, "COMPLIANCE — PERIODOS CRITICOS (Manual Sec. 8.3 e 4.1)", 0.3, 6.18, 12.7)
    T(sl, "Quiet period: [data inicio] — sem comunicacoes publicas sobre oferta  |  Periodo vedacao negociacao administradores: 15d antes resultados\n"
          "Lock-up controladores: 180 dias pos-pricing  |  Lock-up administradores: 90 dias  |  Research embargo: 40 dias pos-pricing",
      0.3, 6.43, 12.5, 0.42, sz=9, col=GRAPHITE)

    # 4 — III. DESEMPENHO FINANCEIRO ------------------------------------------
    sl = base(prs, "III. Desempenho Financeiro — DRE, Balanco, Fluxo e KPIs Operacionais",
              "Manual Sec. 7.2 (DD Financeira) | Modelo Excel Abas 2.DRE + 3.BALANCO + 4.FLUXO | Qualidade dos Lucros", FT_ECM)

    anos_e = ["2022A","2023A","2024A","LTM","2025E","2026E"]
    cw_e = 12.7/(len(anos_e)+1)
    TH(sl, ["Indicador (R$M)"]+anos_e, 0.3, 1.18, cw_e, rh=0.26)
    dre_e = [
        (["Receita Liquida","[__]","[__]","[__]","[__]","[__]","[__]"],WHITE,GRAPHITE,False),
        (["Margem Bruta %","[__]%","[__]%","[__]%","[__]%","[__]%","[__]%"],LIGHT_BG,POSITIVE,False),
        (["EBITDA Ajustado","[__]","[__]","[__]","[__]","[__]","[__]"],RGBColor(0xE8,0xF0,0xFF),NAVY,True),
        (["Margem EBITDA %","[__]%","[__]%","[__]%","[__]%","[__]%","[__]%"],RGBColor(0xE8,0xF0,0xFF),POSITIVE,False),
        (["D&A","([__])","([__])","([__])","([__])","([__])","([__])"],WHITE,GRAPHITE,False),
        (["Lucro Liquido","[__]","[__]","[__]","[__]","[__]","[__]"],WHITE,GRAPHITE,False),
        (["FCO","[__]","[__]","[__]","[__]","[__]","[__]"],LIGHT_BG,GRAPHITE,False),
        (["IQL = FCO/LL","[__]x","[__]x","[__]x","[__]x","[__]x","[__]x"],LIGHT_BG,POSITIVE,False),
        (["CAPEX / Receita","[__]%","[__]%","[__]%","[__]%","[__]%","[__]%"],WHITE,GRAPHITE,False),
        (["CAPEX / D&A","[__]x","[__]x","[__]x","[__]x","[__]x","[__]x"],WHITE,MUTED,False),
        (["Div.Liq./EBITDA","[__]x","[__]x","[__]x","[__]x","[__]x","[__]x"],WHITE,GRAPHITE,False),
    ]
    for i, (row,bg,tc,bd) in enumerate(dre_e):
        TR(sl, row, 0.3, 1.44+i*0.32, cw_e, rh=0.3, bg=bg, tc=tc, bold=bd)

    sec(sl, "KPIs OPERACIONAIS DO SETOR (Manual Sec. 7.3)", 0.3, 5.08, 12.7)
    for i, (k,v,bench) in enumerate([("[KPI 1 — ex: SSS / CAGR vol. / vidas cobertas]","[Valor] [Un]","Bench: [___]"),
                                      ("[KPI 2 — ex: Ticket medio / NPS / Market share]","[Valor] [Un]","Bench: [___]"),
                                      ("[KPI 3 — ex: Sinistralidade / Churn / DSO]","[Valor] [Un]","Bench: [___]"),
                                      ("[KPI 4 — ex: ROIC / Giro ativo / CAPEX/Receita]","[Valor] [Un]","Bench: [___]")]):
        l=0.3+i*3.15; R(sl,l,5.34,3.0,0.86,fill=LIGHT_BG)
        T(sl,k,l+0.1,5.38,2.8,0.28,sz=8,bold=True,col=NAVY)
        T(sl,v,l+0.1,5.66,2.8,0.28,sz=16,bold=True,col=NAVY,align=PP_ALIGN.CENTER)
        T(sl,bench,l+0.1,5.96,2.8,0.2,sz=7.5,col=MUTED,align=PP_ALIGN.CENTER)

    T(sl, "Premissas projecoes: Repositorio Base v1.0 Modulo 1 (IPCA/CDI/cambio) + Modulo 4 (setoriais) | IQL=FCO/LL: avalia qualidade do lucro (Manual Sec. 7.2) | CAPEX/D&A <1,0x = desinvestimento (Manual Sec. 12.2)",
      0.3, 6.28, 12.5, 0.22, sz=7, col=MUTED, italic=True)

    # 5 — IV. VALUATION -------------------------------------------------------
    sl = base(prs, "IV. Valuation — DCF, Multiplos, DDM e Analise de Sensibilidade",
              "Manual Sec. 5.2 (DCF) + 5.3 (Multiplos) + 5.5 (DDM) | Modelo Excel Abas 5.WACC + 6.DCF + 7.MULTIPLOS + 10.SENSIBILIDADE", FT_ECM)

    sec(sl, "DCF — WACC E COMPONENTES (Manual Sec. 5.2.1 | Aba 5)", 0.3, 1.18, 6.2)
    wacc = [("Taxa livre de risco (Rf)","NTN-B [prazo] = IPCA + [__]% (Modulo 1.4 Repos. Base)"),
            ("Premio risco mercado (ERP)","[__]% (Damodaran Emerging Markets Jan/2026)"),
            ("Premio pais (CDS Brasil 5a)","[__]% (Bloomberg [data])"),
            ("Beta desalavancado (βU)","[__] (Damodaran setor [___] Jan/2026)"),
            ("Beta alavancado (βL)","[__] (ajustado pela estrutura de capital alvo [D/E])"),
            ("Ke (Custo Equity)","[__]% a.a. (nominal BRL) = Rf + βL x ERP + Pais"),
            ("Kd (Custo Divida pre-IR)","[__]% a.a. (taxa captacao atual da empresa)"),
            ("Aliquota IR/CSLL","34% (padrao) | [__]% se regime especial"),
            ("WACC Nominal BRL","[__]% a.a. (Bench. setor: Repositorio Base Modulo 5.2)"),
            ("Taxa g (perpetuidade)","[__]% a.a. (Repositorio Base Modulo 5.3 — max 4,5%)"),
            ("Terminal Value / EV total","[__]% do EV (>80% = baixa visibilidade — rever premissas)"),
            ("Equity Value DCF","R$[__]Bi  |  Preco/acao implicito: R$[__]"),]
    for i, (k,v) in enumerate(wacc):
        y=1.45+i*0.36; bg=LIGHT_BG if i%2==0 else WHITE
        if k in ["WACC Nominal BRL","Equity Value DCF"]: bg=RGBColor(0xE8,0xF0,0xFF)
        R(sl,0.3,y,2.6,0.32,fill=bg); T(sl,k,0.38,y+0.06,2.5,0.22,sz=7.5,bold=(k in ["WACC Nominal BRL","Equity Value DCF"]),col=NAVY)
        T(sl,v,2.94,y+0.06,3.48,0.22,sz=8,col=NAVY if k in ["WACC Nominal BRL","Equity Value DCF"] else GRAPHITE,bold=(k in ["WACC Nominal BRL","Equity Value DCF"]))

    sec(sl, "ANALISE DE MULTIPLOS — PEERS LISTADOS (Manual Sec. 5.3 | Aba 7)", 6.65, 1.18, 6.38)
    TH(sl, ["Empresa","Mkt Cap","EV/EBITDA LTM","EV/EBITDA NTM","P/E NTM","EV/Rec."], 6.65, 1.44, cw=6.38/6, rh=0.25)
    peers = [["[Peer 1]","R$[__]Bi","[__]x","[__]x","[__]x","[__]x"],
             ["[Peer 2]","R$[__]Bi","[__]x","[__]x","[__]x","[__]x"],
             ["[Peer 3]","R$[__]Bi","[__]x","[__]x","[__]x","[__]x"],
             ["[Peer 4]","R$[__]Bi","[__]x","[__]x","[__]x","[__]x"],
             ["P25","—","[__]x","[__]x","[__]x","[__]x"],
             ["Mediana","—","[__]x","[__]x","[__]x","[__]x"],
             ["P75","—","[__]x","[__]x","[__]x","[__]x"],
             ["[EMPRESA ALVO]","N/A","[__]x","[__]x","[__]x","[__]x"]]
    for i, row in enumerate(peers):
        is_tgt="ALVO" in row[0]; is_med="Mediana" in row[0]; is_p="P2" in row[0] or "P7" in row[0]
        bg=RGBColor(0xE8,0xF0,0xFF) if is_tgt else (NAVY_MED if is_med else (LIGHT_BG if is_p or i%2==0 else WHITE))
        tc=WHITE if is_med else (NAVY if is_tgt else GRAPHITE)
        TR(sl, row, 6.65, 1.69+i*0.3, cw=6.38/6, rh=0.28, bg=bg, tc=tc, bold=(is_tgt or is_med))

    sec(sl, "DDM (Manual Sec. 5.5) — aplicavel se empresa com historico de dividendos", 6.65, 4.14, 6.38)
    T(sl, "Formula Gordon: P = D1 / (Ke - g)  |  D1 = proximo dividendo anual esperado\n"
          "Ke = [__]%  |  D1 = R$[__]/acao  |  g = [__]% (taxa crescimento sustentavel dividendos)\n"
          "Preco DDM implicito: R$[__]/acao  |  [Aplicavel / N/A para este emissor]",
      6.65, 4.4, 6.3, 0.55, sz=9, col=GRAPHITE)

    sec(sl, "SENSIBILIDADE DCF — WACC x g (Aba 10 Excel)", 0.3, 4.66, 6.2)
    T(sl, "WACC →", 0.3, 4.92, 0.85, 0.24, sz=8, bold=True, col=NAVY)
    wrange = ["[W-1%]","[W base]","[W+1%]"]
    for j,w in enumerate(wrange): T(sl,w,1.2+j*1.55,4.92,1.48,0.24,sz=8.5,bold=True,col=NAVY,align=PP_ALIGN.CENTER)
    grange = ["[g-0,5%]","[g base]","[g+0,5%]"]
    T(sl, "g ↓", 0.3, 5.18, 0.5, 0.24, sz=8, bold=True, col=NAVY)
    for i,g in enumerate(grange):
        T(sl,g,0.3,5.2+i*0.28,1.0,0.24,sz=8,col=NAVY)
        for j in range(3):
            bg=POSITIVE if (i==0 and j==0) else (LIGHT_BG if (i==1 and j==1) else (RGBColor(0xE8,0xF0,0xFF) if i+j<3 else RED_BG))
            R(sl,1.2+j*1.55,5.18+i*0.28,1.48,0.26,fill=bg)
            T(sl,"R$[__]",1.2+j*1.55,5.2+i*0.28,1.48,0.22,sz=9,col=WHITE if bg==POSITIVE else NAVY,align=PP_ALIGN.CENTER,bold=(i==1 and j==1))

    T(sl, "Fontes: Damodaran Jan/2026 | Bloomberg [data] | BACEN Focus | Repositorio Base v1.0 Modulo 5 | Pares: B3 + Research BTG/XP/Itau",
      0.3, 6.02, 6.2, 0.22, sz=7.5, col=MUTED, italic=True)

    sec(sl, "RETORNO AO INVESTIDOR (Manual Sec. 11 | Aba 9)", 6.65, 4.76, 6.38)
    ret = [("Retorno Base (1 ano)","[__]% | P/E saida: [__]x"),("Retorno Otimista","[__]% | Multiplo re-rating: +[__]x"),
           ("Retorno Downside","[__]% | Protecao: dividendo [__]%"),("Payback (anos)","[__] anos"),
           ("Eficiencia tributaria","JCP: [__]% do payout | IR dividendo: [isento/tributavel]"),]
    for i, (k,v) in enumerate(ret):
        y=5.02+i*0.38; R(sl,6.65,y,2.5,0.34,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,k,6.73,y+0.06,2.38,0.22,sz=7.5,bold=True,col=NAVY); T(sl,v,9.18,y+0.06,3.82,0.22,sz=8.5,col=GRAPHITE)

    # 6 — V. DUE DILIGENCE ---------------------------------------------------
    sl = base(prs, "V. Due Diligence — Financeira, Negocios e Legal/Regulatoria",
              "Manual Sec. 7 (3 tipos de DD) — pre-requisito para avanco ao bookbuilding", FT_ECM)

    sec(sl, "7.2 DUE DILIGENCE FINANCEIRA", 0.3, 1.18, 6.2)
    ddf = ["DFs auditadas 3-5 anos — auditor independente registrado CVM",
           "Reconciliacao GAAP vs. IFRS vs. metricas gerenciais (Adjusted EBITDA)",
           "IQL = FCO/LL: avalia qualidade dos lucros — alerta se <0,7x por 2 anos",
           "Accruals = LL - FCO: positivos e crescentes = risco contabil",
           "Analise de working capital: variacao vs. crescimento da receita",
           "CAPEX: manutencao vs. crescimento | CAPEX/D&A vs. benchmark setorial",
           "Analise de divida: covenants, vencimentos, moeda, taxa fixa vs. flutuante",
           "Partes relacionadas: arm's length, disclosure, management fee",
           "Provisoes e contingencias: PCLD, trabalhistas, fiscais, ambientais",
           "Reconhecimento de receita: politica, mudancas recentes, agressividade"]
    for i, item in enumerate(ddf):
        T(sl, f"  ✓  {item}", 0.3, 1.46+i*0.36, 5.9, 0.32, sz=9, col=GRAPHITE)

    sec(sl, "7.3 DUE DILIGENCE DE NEGOCIOS", 6.65, 1.18, 6.38)
    ddn = ["Posicionamento competitivo: market share, barreiras de entrada, moat",
           "Modelo de receita: recorrente vs. pontual, concentracao clientes, ticket medio",
           "Churn rate: [__]%  |  NPS: [__]  |  LTV/CAC: [__]x (min 3x)",
           "Analise operacional: capacidade instalada, eficiencia, CAPEX manutencao",
           "Gestao e governanca: track record, incentivos LP, estrutura de conselho",
           "Analise setorial: ciclo, regulacao, pressao de precos, commoditizacao",
           "Pipeline de produtos/servicos/expansao — cronograma e investimentos",
           "Principais riscos do negocio: [Risco 1] | [Risco 2] | [Risco 3]"]
    for i, item in enumerate(ddn):
        T(sl, f"  ✓  {item}", 6.65, 1.46+i*0.38, 6.28, 0.34, sz=9, col=GRAPHITE)

    sec(sl, "7.4 DUE DILIGENCE LEGAL E REGULATORIA", 0.3, 5.0, 6.2)
    ddl = ["Estrutura societaria: ON/PN, tag-along, poison pill, direitos controladores",
           "Contingencias: provisoes vs. riscos possiveis (trabalhistas, fiscais, ambientais, civeis)",
           "Propriedade intelectual: patentes, marcas, licencas — risco expiração/litígio",
           "Contratos relevantes: change of control, exclusividade, take-or-pay",
           "Compliance: LGPD, ANTT, ANVISA, BACEN — autorizacoes e licencas em vigor"]
    for i, item in enumerate(ddl):
        T(sl, f"  ✓  {item}", 0.3, 5.27+i*0.38, 5.9, 0.34, sz=9, col=GRAPHITE)

    sec(sl, "STATUS DA DUE DILIGENCE", 6.65, 5.0, 6.38)
    TH(sl, ["Area","Status","Responsavel","Prazo"], 6.65, 5.27, cw=6.38/4, rh=0.25)
    dd_status = [["Financeira","[ ] Completa  [ ] Em andamento","[Nome]","[Data]"],
                 ["Negocios","[ ] Completa  [ ] Em andamento","[Nome]","[Data]"],
                 ["Legal","[ ] Completa  [ ] Em andamento","[Advogados]","[Data]"],
                 ["Tributaria","[ ] Completa  [ ] Em andamento","[Advogados]","[Data]"],
                 ["Ambiental","[ ] Completa  [ ] N/A","[Consultor]","[Data]"]]
    for i, row in enumerate(dd_status):
        TR(sl, row, 6.65, 5.52+i*0.3, cw=6.38/4, rh=0.28, bg=LIGHT_BG if i%2==0 else WHITE)

    # 7 — VI. ESTRUTURA DA OFERTA + BOOKBUILDING ------------------------------
    sl = base(prs, "VI. Estrutura da Oferta, Bookbuilding e Timing",
              "Manual Sec. 4 (8 Fases) + Sec. 6 (Precificacao e Bookbuilding) + Sec. 11 (Timing e Otimizacao)", FT_ECM)

    sec(sl, "ESTRUTURA FORMAL (Manual Sec. 3 e 4)", 0.3, 1.18, 6.2)
    est = [("Tipo","[IPO / Follow-on / Block Trade]"),("Regime","[RCVM 160 art.3o / CVM 400]"),
           ("Segmento B3","[Novo Mercado / Nivel 2 / Nivel 1 / EGEM]"),("Volume Base","R$[__]Bi — [__]M acoes"),
           ("Primaria","R$[__]M — [__]M acoes novas emitidas"),("Secundaria","R$[__]M — [__]M acoes dos vendedores"),
           ("Greenshoe","Ate 15% — R$[__]M | Agente: [Nome] | Prazo: 30d"),("Lote Adicional","Ate 20% — R$[__]M"),
           ("Free Float base","[__]%  |  c/ Greenshoe: [__]%"),("Lock-up","Controladores: 180d | Adm.: 90d"),]
    for i, (k,v) in enumerate(est):
        y=1.45+i*0.36; R(sl,0.3,y,2.0,0.32,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,k,0.38,y+0.06,1.88,0.22,sz=7.5,bold=True,col=NAVY); T(sl,v,2.34,y+0.06,3.78,0.22,sz=8.5,col=GRAPHITE)

    sec(sl, "BOOKBUILDING — MECÂNICA E QUALIDADE (Manual Sec. 6.2-6.4)", 6.65, 1.18, 6.38)
    bb = [("Definicao da Faixa","Football Field + pre-marketing com [__] anchor investors | Desconto [10-15%] vs. fair value para garantir upside"),
          ("Coleta de Ordens","Strike bids = prioridade maxima | Limit bids = proporcional | Meta: >3x oversubscribed"),
          ("Criterios Alocacao","Qualidade (long-only > momentum) | Relacao banco | Geo: [__]% local / [__]% estrangeiro"),
          ("Sinais Qualidade","Livro >5x = muito forte | >3x = bom | <2x = atencao | <1,5x = reavaliar pricing"),
          ("Greenshoe / Estabilizacao","Coordenador compra no secundario se preco cair | Limite: at 15% / 30 dias"),
          ("Pricing Final","Comite de Pricing: [__] membros | Criterio: demanda + perfil investidores + condicoes macro"),]
    for i, (k,v) in enumerate(bb):
        y=1.45+i*0.67; R(sl,6.65,y,6.38,0.6,fill=LIGHT_BG if i%2==0 else WHITE)
        T(sl,k,6.73,y+0.04,1.8,0.22,sz=8,bold=True,col=NAVY); T(sl,v,8.58,y+0.04,4.38,0.52,sz=8.5,col=GRAPHITE)

    sec(sl, "TIMING DA OFERTA — INDICADORES DE JANELA (Manual Sec. 11.1)", 0.3, 5.23, 6.2)
    timing = [("✓","Ibovespa em alta com volatilidade (VIX local) abaixo de 20",POSITIVE),
              ("✓","IPOs recentes precificados acima do meio da faixa",POSITIVE),
              ("✓","Fluxo estrangeiro positivo (EMBI+ comprimindo)",POSITIVE),
              ("⚠","Ausencia de COPOM, eleicoes ou eventos geopoliticos nos 30 dias seguintes",AMBER),
              ("✗","VIX local >25 ou EMBI+ >250bps = FECHAR janela",NEGATIVE),]
    for i, (icon,item,col) in enumerate(timing):
        T(sl, f"  {icon}  {item}", 0.3, 5.5+i*0.36, 5.9, 0.32, sz=9, col=col)

    sec(sl, "CRONOGRAMA (8 FASES — Manual Sec. 4)", 6.65, 5.23, 6.38)
    fases = [("Fase 1","Mandato + Kick-off","[Data]"),("Fase 2","Pre-marketing + Due Diligence","[Data]"),
             ("Fase 3","Protocolo CVM + Prospecto","[Data]"),("Fase 4","Roadshow ([__] cidades)","[Data]"),
             ("Fase 5","Bookbuilding ([__] dias)","[Data]"),("Fase 6","Pricing + Alocacao","[Data]"),
             ("Fase 7","Liquidacao (D+2)","[Data]"),("Fase 8","Pos-oferta + Pesquisa","D+40")]
    for i, (fase,desc,data) in enumerate(fases):
        col_n=i%2; row_n=i//2; l=6.65+col_n*3.25; t=5.5+row_n*0.4
        R(sl,l,t,3.1,0.36,fill=NAVY if col_n==0 else NAVY_MED)
        T(sl,f"{fase}: {desc}",l+0.08,t+0.06,2.2,0.22,sz=7.5,bold=True,col=WHITE)
        T(sl,data,l+2.32,t+0.08,0.72,0.2,sz=7,col=GOLD)

    # 8 — VII. RED + YELLOW FLAGS + RECOMENDACAO -----------------------------
    sl = base(prs, "VII. Red Flags, Yellow Flags e Recomendacao Final de Pricing",
              "Manual Sec. 12.1 (Red Flags — nao esta pronta) + Sec. 12.2 (Yellow Flags — disclosure) + Comite de Pricing", FT_ECM)

    sec(sl, "RED FLAGS — EMPRESA NAO ESTA PRONTA (>=2 flags = comunicar MD) — Manual Sec. 12.1", 0.3, 1.18, 12.7)
    T(sl, "Presenca de 2 ou mais red flags DEVE ser comunicada ao MD antes de avancar para o protocolo CVM.", 0.3, 1.44, 12.5, 0.2, sz=8, col=NEGATIVE, bold=True)
    red_flags = [("FCO/LL < 0,7x por 2 anos (lucros sem caixa)","[ ] Presente  [ ] Ausente"),
                 ("EBITDA Adjusted >20% vs. EBITDA IFRS","[ ] Presente  [ ] Ausente"),
                 ("Litigio material >10% do PL nao provisionado","[ ] Presente  [ ] Ausente"),
                 ("Concentracao de cliente >30% da receita","[ ] Presente  [ ] Ausente"),
                 ("Mudanca de auditor ou politica contabil em 12m","[ ] Presente  [ ] Ausente"),
                 ("CAPEX < D&A por >2 anos (desinvestimento)","[ ] Presente  [ ] Ausente"),
                 ("Partes relacionadas >15% do ativo sem disclosure","[ ] Presente  [ ] Ausente"),
                 ("Oferta >60% secundaria sem plano de uso de recursos","[ ] Presente  [ ] Ausente"),
                 ("Ebitda TV >80% do EV (baixa visibilidade)","[ ] Presente  [ ] Ausente"),
                 ("Peers inadequados (risco/crescimento muito diferente)","[ ] Presente  [ ] Ausente"),]
    for i, (flag, status) in enumerate(red_flags):
        col_n=i%2; row_n=i//2; l=0.3+col_n*6.45; y=1.68+row_n*0.38
        R(sl,l,y,0.28,0.3,fill=NEGATIVE)
        T(sl,"!",l,y+0.03,0.28,0.24,sz=10,bold=True,col=WHITE,align=PP_ALIGN.CENTER)
        T(sl,flag,l+0.33,y+0.06,4.15,0.22,sz=8.5,col=GRAPHITE)
        R(sl,l+4.52,y,1.85,0.3,fill=LIGHT_BG)
        T(sl,status,l+4.57,y+0.06,1.78,0.22,sz=8,col=NAVY,bold=True)

    sec(sl, "YELLOW FLAGS — PONTOS DE ATENCAO PARA DISCLOSURE NO PROSPECTO (Manual Sec. 12.2)", 0.3, 3.68, 12.7)
    yellow = [("Margem EBITDA abaixo da mediana do peer group por >2 anos","Incluir no Fator de Risco com texto especifico e quantificado"),
              ("Crescimento receita acima do setor sem explicacao clara","Risk de market share artificial — disclosure de premissas"),
              ("SG&A como % receita crescendo mais rapido que receita","Deseconomias de escala — incluir plano de eficiencia"),
              ("CAPEX historico <D&A — subinvestimento no negocio","Risco de deterioracao operacional futura — quantificar CAPEX requerido"),
              ("NRR (Net Revenue Retention) <100% em SaaS/recorrente","Churn implicito — disclosure com dados historicos de cohorte"),]
    for i, (flag, acao) in enumerate(yellow):
        y=3.95+i*0.42
        R(sl,0.3,y,0.28,0.34,fill=AMBER)
        T(sl,"~",0.3,y+0.05,0.28,0.24,sz=10,bold=True,col=WHITE,align=PP_ALIGN.CENTER)
        T(sl,flag,0.63,y+0.06,6.2,0.22,sz=8.5,bold=True,col=GRAPHITE)
        T(sl,f"→  {acao}",6.88,y+0.06,6.1,0.22,sz=8.5,col=POSITIVE)

    sec(sl, "RECOMENDACAO FINAL — COMITE DE PRICING (Manual Sec. 6.1)", 0.3, 5.96, 12.7)
    T(sl, "[ ]  AVANCAR — Faixa indicativa: R$[__] — R$[__] | Preco recomendado: R$[__] (midpoint) | EV/EBITDA NTM: [__]x | Upside: [__]%\n"
          "[ ]  AVANCAR COM AJUSTE — [Descrever ajuste: reduzir faixa / resolver red flag / aguardar mercado]\n"
          "[ ]  ADIAR — Janela de mercado desfavoravel ou red flags materiais identificados",
      0.3, 6.22, 9.0, 0.75, sz=10, col=NAVY)
    T(sl, "Analista: ________________\nMD Responsavel: ________________\nData: ___/___/______",
      9.5, 6.22, 3.5, 0.75, sz=9, col=GRAPHITE)

    path = str(Path(OUT_DIR) / "equity_opinion_ecm.pptx")
    prs.save(path)
    print(f"ECM: {path} ({len(prs.slides)} slides)")
    return path


if __name__ == "__main__":
    import os; os.makedirs(OUT_DIR, exist_ok=True)
    build_dcm()
    build_ecm()
    print("Concluido.")
