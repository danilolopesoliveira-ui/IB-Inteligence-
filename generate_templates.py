"""
Generate professional IB templates — Farallon-inspired design.
Palette: Clean white, deep navy, thin lines, serif titles, institutional.
Run: python generate_templates.py
"""
from pathlib import Path
from datetime import datetime

OUT = Path("./templates/models")
OUT.mkdir(parents=True, exist_ok=True)

# Farallon palette — light, institutional, ultra-clean
BG  = (0xFF, 0xFF, 0xFF)  # white background
NAV = (0x0F, 0x1F, 0x3D)  # deep navy (primary accent)
NV2 = (0x1A, 0x2B, 0x4A)  # navy lighter
GRF = (0x4A, 0x55, 0x68)  # graphite (body text)
GR1 = (0x71, 0x80, 0x96)  # mid gray
GR2 = (0xA0, 0xAE, 0xC0)  # light gray (borders)
GR3 = (0xE2, 0xE8, 0xF0)  # very light gray (table alt rows)
WHT = (0xFF, 0xFF, 0xFF)
BLK = (0x1A, 0x20, 0x2C)  # near-black for titles

def gen_pptx(title, company, subtitle, slides_data, filename):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    C = lambda r,g,b: RGBColor(r,g,b)

    def mk(title_text, body=None, tbl=None, is_cover=False):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = C(*BG)
        n = len(prs.slides)
        # Thin navy line at top
        s = sl.shapes.add_shape(1, 0, 0, Inches(13.33), Emu(18000))
        s.fill.solid(); s.fill.fore_color.rgb = C(*NAV); s.line.fill.background()
        # Footer line
        s = sl.shapes.add_shape(1, Inches(0.6), Inches(7.05), Inches(12.1), Emu(6000))
        s.fill.solid(); s.fill.fore_color.rgb = C(*GR2); s.line.fill.background()
        # Footer text
        tb = sl.shapes.add_textbox(Inches(0.6), Inches(7.1), Inches(8), Inches(0.25))
        p = tb.text_frame.paragraphs[0]; p.text = "Confidencial  |  IB-Agents Intelligence Platform"
        p.font.size = Pt(7); p.font.color.rgb = C(*GR1); p.font.name = "Calibri"
        # Page number
        tb = sl.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.25))
        p = tb.text_frame.paragraphs[0]; p.text = str(n); p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(7); p.font.color.rgb = C(*GR1); p.font.name = "Calibri"

        if is_cover:
            # Navy block at top (taller)
            s = sl.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(3.2))
            s.fill.solid(); s.fill.fore_color.rgb = C(*NAV); s.line.fill.background()
            return sl

        # Title
        tb = sl.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.5))
        p = tb.text_frame.paragraphs[0]; p.text = title_text
        p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = C(*NAV); p.font.name = "Georgia"
        # Thin line under title
        s = sl.shapes.add_shape(1, Inches(0.6), Inches(0.9), Inches(12.1), Emu(6000))
        s.fill.solid(); s.fill.fore_color.rgb = C(*GR2); s.line.fill.background()

        if body:
            tb = sl.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12), Inches(5.5))
            tf = tb.text_frame; tf.word_wrap = True
            for i, line in enumerate(body):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.font.name = "Calibri"
                if line.startswith("##"):
                    p.text = line[2:].strip(); p.font.size = Pt(11); p.font.bold = True
                    p.font.color.rgb = C(*NAV); p.space_before = Pt(14)
                elif line.startswith("•"):
                    p.text = line; p.font.size = Pt(9); p.font.color.rgb = C(*GRF)
                    p.space_before = Pt(2); p.level = 1
                elif line.startswith("["):
                    p.text = line; p.font.size = Pt(8); p.font.color.rgb = C(*GR1)
                    p.font.italic = True; p.space_before = Pt(2)
                elif line == "":
                    p.text = ""; p.space_before = Pt(4)
                else:
                    p.text = line; p.font.size = Pt(9); p.font.color.rgb = C(*GRF)
                    p.space_before = Pt(3)

        if tbl:
            y = Inches(4.0) if body else Inches(1.15)
            rows = len(tbl); cols = len(tbl[0])
            t = sl.shapes.add_table(rows, cols, Inches(0.6), y, Inches(12), Inches(rows*0.32)).table
            for r, row in enumerate(tbl):
                for c, val in enumerate(row):
                    cell = t.cell(r, c); cell.text = str(val)
                    for pp in cell.text_frame.paragraphs:
                        pp.font.size = Pt(8); pp.font.name = "Calibri"
                        pp.font.color.rgb = C(*WHT) if r == 0 else C(*BLK)
                        pp.font.bold = r == 0
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C(*NAV) if r == 0 else C(*GR3) if r % 2 == 0 else C(*BG)
        return sl

    # ── COVER ──
    sl = mk("", is_cover=True)
    # Company name on navy band
    tb = sl.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(10), Inches(0.8))
    p = tb.text_frame.paragraphs[0]; p.text = company
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = C(*WHT); p.font.name = "Georgia"
    # Subtitle
    tb = sl.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(10), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = subtitle
    p.font.size = Pt(14); p.font.color.rgb = C(0xC0,0xD0,0xE8); p.font.name = "Calibri"
    # Title below navy band
    tb = sl.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.size = Pt(12); p.font.color.rgb = C(*GRF); p.font.name = "Calibri"
    # Date and confidential
    tb = sl.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(8), Inches(1))
    tf = tb.text_frame
    for txt in [f"{datetime.now().strftime('%B %Y')}", "IB-Agents Intelligence Platform", "Strictly Private and Confidential"]:
        p = tf.add_paragraph(); p.text = txt; p.font.size = Pt(9); p.font.color.rgb = C(*GR1)
        p.font.name = "Calibri"

    # ── CONTENT SLIDES ──
    for s_title, s_body, s_tbl in slides_data:
        mk(s_title, s_body, s_tbl)

    prs.save(str(OUT / filename))
    print(f"  PPTX: {filename} ({len(prs.slides)} slides)")

# ═════════════════════════════════════════════════════════════════════════════

def gen_xlsx(name, sheets, filename):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    wb = Workbook()
    HF = Font(name="Calibri", size=9, bold=True, color="FFFFFF")
    HFL = PatternFill("solid", fgColor="0F1F3D")
    NF = Font(name="Calibri", size=9, color="1A202C")
    BF = Font(name="Calibri", size=9, bold=True, color="0F1F3D")
    BRD = Border(bottom=Side(style="thin", color="CBD5E0"))
    for idx, (sn, data) in enumerate(sheets):
        ws = wb.active if idx == 0 else wb.create_sheet()
        ws.title = sn; ws.sheet_properties.tabColor = "0F1F3D"
        for r, row in enumerate(data):
            for c, v in enumerate(row):
                cell = ws.cell(row=r+1, column=c+1, value=v)
                cell.font = HF if r == 0 else NF; cell.border = BRD
                if r == 0: cell.fill = HFL; cell.alignment = Alignment(horizontal="center")
                if isinstance(v, float) and abs(v) < 1 and v != 0: cell.number_format = '0.0%'
                elif isinstance(v, (int, float)) and abs(v) >= 100: cell.number_format = '#,##0'
                ws.column_dimensions[chr(65+min(c,25))].width = max(len(str(v))+5, 15)
    wb.save(str(OUT / filename))
    print(f"  XLSX: {filename} ({len(wb.sheetnames)} abas)")

# ═════════════════════════════════════════════════════════════════════════════

def gen_pdf(title, sections, filename):
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.lib.colors import HexColor
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(OUT / filename), pagesize=A4)
        w, h = A4; pg = [0]
        def footer():
            c.setStrokeColor(HexColor("#A0AEC0")); c.setLineWidth(0.3)
            c.line(2*cm, 1.3*cm, w-2*cm, 1.3*cm)
            c.setFont("Helvetica", 6); c.setFillColor(HexColor("#718096"))
            c.drawString(2*cm, 0.8*cm, "Confidencial  |  IB-Agents Intelligence Platform")
            c.drawRightString(w-2*cm, 0.8*cm, f"{pg[0]}")
        # Cover
        # Cover — white background with navy header band
        pg[0] = 0; c.setFillColor(HexColor("#FFFFFF")); c.rect(0, 0, w, h, fill=1)
        # Navy band at top
        c.setFillColor(HexColor("#0F1F3D")); c.rect(0, h-6*cm, w, 6*cm, fill=1)
        # Title on navy
        c.setFont("Times-Bold", 22); c.setFillColor(HexColor("#FFFFFF"))
        c.drawString(2*cm, h-3.5*cm, title.split("—")[0].strip()[:50])
        if "—" in title:
            c.setFont("Times-Roman", 12); c.setFillColor(HexColor("#C0D0E8"))
            c.drawString(2*cm, h-4.5*cm, title.split("—",1)[1].strip()[:60])
        # Thin line separator
        c.setStrokeColor(HexColor("#A0AEC0")); c.setLineWidth(0.5)
        c.line(2*cm, h-7*cm, w-2*cm, h-7*cm)
        c.setFont("Helvetica", 9); c.setFillColor(HexColor("#4A5568"))
        c.drawString(2*cm, h-7.8*cm, f"{datetime.now().strftime('%B %Y')}  |  IB-Agents Intelligence Platform")
        c.drawString(2*cm, h-8.4*cm, "Strictly Private and Confidential")
        c.showPage()
        # Content pages
        for sec_t, paras in sections:
            pg[0] += 1
            c.setFillColor(HexColor("#FFFFFF")); c.rect(0, 0, w, h, fill=1)
            # Navy thin bar at top
            c.setFillColor(HexColor("#0F1F3D")); c.rect(0, h-0.4*cm, w, 0.4*cm, fill=1)
            # Section title
            c.setFont("Times-Bold", 13); c.setFillColor(HexColor("#0F1F3D"))
            c.drawString(2*cm, h-1.5*cm, sec_t)
            # Thin line under title
            c.setStrokeColor(HexColor("#A0AEC0")); c.setLineWidth(0.3)
            c.line(2*cm, h-1.8*cm, w-2*cm, h-1.8*cm)
            footer()
            y = h - 2.5*cm
            for p in paras:
                if y < 3*cm: c.showPage(); pg[0] += 1; c.setFillColor(HexColor("#FFFFFF")); c.rect(0,0,w,h,fill=1); c.setFillColor(HexColor("#0F1F3D")); c.rect(0,h-0.4*cm,w,0.4*cm,fill=1); footer(); y = h - 1.5*cm
                if p.startswith("##"):
                    y -= 0.3*cm; c.setFont("Times-Bold", 10); c.setFillColor(HexColor("#0F1F3D"))
                    c.drawString(2*cm, y, p[2:].strip()); y -= 0.5*cm
                elif p.startswith("•"):
                    c.setFont("Helvetica", 8.5); c.setFillColor(HexColor("#4A5568"))
                    words = p.split(); line = ""; lines = []
                    for word in words:
                        if c.stringWidth(line + " " + word, "Helvetica", 8.5) > 15*cm:
                            lines.append(line); line = word
                        else: line = (line + " " + word).strip()
                    if line: lines.append(line)
                    for l in lines: c.drawString(2.5*cm, y, l); y -= 0.38*cm
                elif p == "": y -= 0.25*cm
                else:
                    c.setFont("Helvetica", 9); c.setFillColor(HexColor("#2D3748"))
                    words = p.split(); line = ""; lines = []
                    for word in words:
                        if c.stringWidth(line + " " + word, "Helvetica", 9) > 16*cm:
                            lines.append(line); line = word
                        else: line = (line + " " + word).strip()
                    if line: lines.append(line)
                    for l in lines: c.drawString(2*cm, y, l); y -= 0.4*cm
            c.showPage()
        c.save()
        print(f"  PDF:  {filename} ({pg[0]+1} paginas)")
    except ImportError:
        print(f"  SKIP: {filename} (reportlab not installed)")

# ═════════════════════════════════════════════════════════════════════════════
# GENERATE
# ═════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Gerando templates — Design Apollo...\n")

    # 1. DCM Pitch Book
    gen_pptx("Oferta Publica de Debentures", "ENEVA S.A.", "R$ 800.000.000 | 1a e 2a Series | AA+ (Fitch)", [
        ("Sumario Executivo", [
            "## Destaques da Operacao",
            "• Emissor: Eneva S.A. — maior geradora termica privada a gas natural do Brasil (5,9 GW)",
            "• Volume: R$ 800M em debentures simples, quirografaria com cessao fiduciaria de recebiveis PPA",
            "• 1a Serie: R$ 480M, 5 anos, CDI+1,85% a.a. | 2a Serie: R$ 320M, 7 anos, IPCA+6,20% a.a.",
            "• Rating: AA+ (Fitch) — perspectiva estavel | 5a emissao — track record solido",
            "• Garantia: cessao fiduciaria de recebiveis de PPAs (contratos 15+ anos, receita fixa)",
            "• Destinacao: 60% refinanciamento + 40% expansao Parnaiba VI (380 MW)",
            "",
            "## Racional de Investimento",
            "• 92% da receita contratada via PPAs de longo prazo — fluxo previsivel",
            "• Net Debt/EBITDA de 1,2x (covenant: 3,5x) — folga de 2,3x",
            "• Operacao integrada gas+energia — menor custo da regiao Norte/Nordeste",
            "• Demanda estimada: R$ 1,4Bi (1,75x oversubscription)",
        ], [["Metrica", "1a Serie", "2a Serie", "Ref. Mercado"],
            ["Spread", "CDI+1,85%", "IPCA+6,20%", "AA+: 1,60-2,00%"],
            ["Prazo", "5 anos", "7 anos", "—"],
            ["Duration", "4,6 anos", "5,8 anos", "—"],
            ["Amortizacao", "Bullet", "Linear (5o ano+)", "—"]]),
        ("Panorama Setorial", [
            "## Setor Eletrico Brasileiro — Geracao Termica",
            "• Capacidade instalada: 198 GW (2024) — termica: 28% da matriz",
            "• Crescimento demanda: +3,2% a.a. (EPE 2025-2030) | PIB 2025E: +2,5%",
            "• Gas natural: combustivel de transicao — menor emissao vs carvao/oleo",
            "• Novo Marco do Gas (Lei 14.134/2021) — abertura do mercado, beneficio para Eneva",
            "",
            "## Ambiente Macro",
            "• SELIC: 13,25% | Projecao dez/25: 12,50% (Focus) | IPCA: 4,8%",
            "• Mercado de debentures: R$ 320Bi emitidos em 2024 (+18% YoY)",
            "• Spread medio AA+ (5y): 1,78% (ANBIMA mar/25)",
            "[Fonte: ANEEL, EPE, BCB Focus, ANBIMA — marco/2025]",
        ], None),
        ("Perfil do Emissor", [
            "## Eneva S.A. — Visao Geral",
            "• Listada B3 (ENEV3) — Novo Mercado | Market cap: R$ 10,6Bi",
            "• Operacao integrada: exploracao gas + geracao energia na Bacia do Parnaiba (MA)",
            "• 92% receita contratada PPAs (15-25 anos) — previsibilidade de caixa",
            "",
            "## Financials Consolidados (R$ M)",
        ], [["", "2021", "2022", "2023", "2024", "CAGR 3a"],
            ["Receita Liq.", "4.850", "5.920", "7.180", "8.200", "19,1%"],
            ["EBITDA", "2.180", "2.650", "3.320", "3.800", "20,3%"],
            ["Margem EBITDA", "44,9%", "44,8%", "46,2%", "46,3%", "—"],
            ["Lucro Liquido", "850", "1.050", "1.380", "1.600", "23,4%"],
            ["Divida Liquida", "4.400", "4.700", "4.600", "4.500", "—"],
            ["DL/EBITDA", "2,0x", "1,8x", "1,4x", "1,2x", "—"],
            ["EBITDA/Desp.Fin", "3,8x", "3,5x", "4,2x", "4,8x", "—"],
            ["FCF Operacional", "980", "1.150", "1.520", "1.700", "—"]]),
        ("Estrutura da Emissao", [], [
            ["Caracteristica", "1a Serie", "2a Serie"],
            ["Volume", "R$ 480.000.000", "R$ 320.000.000"],
            ["Prazo", "5 anos (mar/2030)", "7 anos (mar/2032)"],
            ["Remuneracao", "CDI + 1,85% a.a.", "IPCA + 6,20% a.a."],
            ["Pgto Juros", "Semestral", "Semestral"],
            ["Amortizacao", "Bullet", "Linear (a partir 5o ano)"],
            ["Especie", "Quirografaria c/ cessao fid.", "Quirografaria c/ cessao fid."],
            ["Rating", "AA+ (Fitch)", "AA+ (Fitch)"],
            ["Publico", "Investidores qualificados", "Investidores qualificados"],
            ["Registro", "RCVM 160 — automatico", "RCVM 160 — automatico"],
            ["Ag. Fiduciario", "Pentagonal DTVM", "Pentagonal DTVM"],
            ["Liquidacao", "B3 (CETIP21)", "B3 (CETIP21)"]]),
        ("Pricing & Spread Analysis", [
            "## Posicionamento na Curva de Credito",
            "• Base: curva DI futuro (B3) + premio de credito por rating/setor/prazo",
            "• 1a Serie (CDI+1,85%, 5y): em linha com mediana do setor (1,78-1,95%)",
            "• 2a Serie (IPCA+6,20%, 7y): desconto de 15bps vs NTN-B + spread medio",
            "• Premio de nova emissao (NIP): ~5-10bps para garantir demanda",
        ], [["Emissor", "Rating", "Prazo", "Spread", "Volume", "Data"],
            ["CPFL Energia", "AA+", "5y", "CDI+1,78%", "R$ 1,2Bi", "Jan/25"],
            ["Taesa", "AAA", "7y", "CDI+1,45%", "R$ 800M", "Fev/25"],
            ["Neoenergia", "AA", "5y", "CDI+1,92%", "R$ 1,5Bi", "Jan/25"],
            ["Equatorial", "AA+", "7y", "IPCA+6,35%", "R$ 2,0Bi", "Mar/25"],
            ["Engie Brasil", "AAA", "10y", "IPCA+5,80%", "R$ 1,0Bi", "Fev/25"],
            ["ENEVA (esta)", "AA+", "5+7y", "CDI+1,85/IPCA+6,20", "R$ 800M", "Mar/25"]]),
        ("Covenants", [
            "## Covenants Financeiros (teste semestral)",
            "• Divida Liquida / EBITDA: <= 3,5x (atual: 1,2x — folga de 2,3x)",
            "• EBITDA / Desp. Financeira: >= 2,5x (atual: 4,8x)",
            "• DSCR: >= 1,3x",
            "",
            "## Covenants Nao-Financeiros",
            "• Cross-default: inadimplencia em dividas > R$ 150M",
            "• Change of control: vencimento antecipado",
            "• Restricted Payments: dividendos > minimo legal somente se DL/EBITDA <= 3,0x",
            "• Alienacao de ativos essenciais: vedada sem anuencia (> 10% ativo total)",
            "",
            "## Eventos de Vencimento Antecipado",
            "• Inadimplencia juros/principal > 5 dias uteis",
            "• Descumprimento covenant nao curado em 60 dias",
            "• Pedido recuperacao judicial/falencia",
            "• Rebaixamento rating abaixo de A-",
        ], None),
        ("Mapa de Distribuicao", [
            "## 42 instituicoes mapeadas — Demanda estimada: R$ 1,4Bi",
        ], [["Instituicao", "Tipo", "AUM (R$Bi)", "Ticket Est.", "Score"],
            ["Itau Asset", "Asset", "780", "R$ 80M", "95%"],
            ["BTG Pactual Asset", "Asset", "320", "R$ 60M", "92%"],
            ["Previ (BB)", "F. Pensao", "250", "R$ 100M", "90%"],
            ["XP Investimentos", "Banco", "150", "R$ 40M", "88%"],
            ["Kinea", "Asset", "85", "R$ 30M", "85%"],
            ["Funcef", "F. Pensao", "110", "R$ 50M", "82%"],
            ["Verde Asset", "Asset", "42", "R$ 25M", "80%"],
            ["SPX Capital", "Asset", "55", "R$ 20M", "78%"],
            ["Outros (34)", "Diversos", "—", "R$ 495M", "60-75%"]]),
        ("Cronograma", [], [
            ["Etapa", "Data", "Status"],
            ["Mandato e kick-off", "15/Fev/2025", "Concluido"],
            ["Due diligence", "18/Fev — 10/Mar", "Concluido"],
            ["Rating (Fitch)", "20/Fev", "Concluido — AA+"],
            ["Escritura de emissao", "01-15/Mar", "Em andamento"],
            ["Aprovacao CA", "18/Mar", "Pendente"],
            ["Registro CVM", "20/Mar", "Pendente"],
            ["Bookbuilding", "25/Mar — 02/Abr", "—"],
            ["Pricing", "02/Abr", "—"],
            ["Liquidacao (D+2)", "04/Abr", "—"],
            ["Inicio negociacao", "07/Abr", "—"]]),
        ("Fatores de Risco", [
            "## Riscos do Emissor",
            "• Concentracao geografica: 70% da capacidade na Bacia do Parnaiba (MA)",
            "• Reservas de gas: dependencia de reservas provadas e provaveis",
            "• Regulatorio: alteracoes nas regras ANEEL de contratacao de energia",
            "",
            "## Riscos da Emissao",
            "• Credito: capacidade de pagamento atrelada ao FCF operacional",
            "• Liquidez: mercado secundario pode ter liquidez limitada",
            "• Mercado: variacao CDI/IPCA impacta custo efetivo da divida",
            "",
            "## Mitigantes",
            "• 92% receita contratada PPA (15+ anos) — fluxo previsivel",
            "• Rating AA+ estavel — baixa probabilidade de default",
            "• Cessao fiduciaria de recebiveis como garantia",
            "• 5 emissoes anteriores bem-sucedidas",
        ], None),
    ], "pitch_book_debentures_infra.pptx")

    # 2. ECM Pitch Book
    gen_pptx("Oferta Publica de Acoes — Follow-on", "SLC AGRICOLA S.A.", "R$ 1.200.000.000 | Primario + Secundario", [
        ("Investment Highlights", [
            "## Equity Story",
            "• #1 em escala: 680.000 ha — maior produtora individual de graos do Brasil",
            "• Custo/ha 12% abaixo da media regional — escala gera vantagem competitiva",
            "• 70% receita indexada USD — hedge cambial natural",
            "• FCF yield 8,5% | Dividend yield 4,2% | Payout 50%",
            "• Negocia 6,8x EV/EBITDA vs 7,5x mediana peers (-9% desconto)",
            "",
            "## Uso dos Recursos (Primario R$ 720M)",
            "• 50% expansao area (+80.000 ha ate 2027)",
            "• 30% tecnologia (agricultura de precisao, IA)",
            "• 20% reducao divida (DL/EBITDA → 1,1x)",
        ], None),
        ("Valuation — Football Field", [
            "## Price Range por Acao (R$)",
            "• EV/EBITDA LTM (7,5x): R$ 17,80 | NTM (6,5x): R$ 19,20",
            "• P/L NTM (12x): R$ 18,50 | DCF (WACC 13,2%, g 4%): R$ 21,30",
            "• Media ponderada: R$ 19,50 | Cotacao SLCE3: R$ 17,20 (+13% upside)",
            "",
            "## Offer Price Range",
            "• Floor: R$ 17,50 (desconto 10% vs fair value)",
            "• Mid: R$ 19,00 (desconto 5% — padrao follow-on)",
            "• Cap: R$ 20,50 (proximo ao DCF)",
        ], [["Metodologia", "Low", "Mid", "High", "Peso"],
            ["EV/EBITDA LTM", "16,20", "17,80", "19,40", "20%"],
            ["EV/EBITDA NTM", "17,50", "19,20", "20,90", "25%"],
            ["P/L NTM", "16,80", "18,50", "20,20", "15%"],
            ["DCF", "18,50", "21,30", "24,10", "30%"],
            ["Precedent M&A", "18,00", "20,10", "22,20", "10%"],
            ["Ponderado", "17,50", "19,50", "21,80", "100%"]]),
        ("Estrutura da Oferta", [], [
            ["Caracteristica", "Detalhe"],
            ["Tipo", "Follow-on (primario + secundario)"],
            ["Volume Total", "R$ 1.200.000.000"],
            ["Primario", "R$ 720M (60%) — caixa empresa"],
            ["Secundario", "R$ 480M (40%) — acionistas"],
            ["Acoes Ofertadas", "63,2M (ao mid R$ 19,00)"],
            ["Diluicao Primaria", "8,2%"],
            ["Free Float Pos", "42,5% (atende NM 25%)"],
            ["Greenshoe (15%)", "9,5M acoes (R$ 180M)"],
            ["Lock-up Control.", "180 dias"],
            ["Registro", "RCVM 160 — ANBIMA"]]),
        ("Bookbuilding Strategy", [
            "## Cronograma",
            "• Reservas: 25/Mar — 05/Abr (8 DU) | Pricing: 07/Abr | Liquidacao: D+2",
            "",
            "## Alocacao Sugerida",
            "• Institucional local: 45% (fundos, previdencia)",
            "• Institucional internacional: 30% (Reg S / 144A)",
            "• Varejo: 15% (reservas em corretoras)",
            "• Ancoras: 10% (GIC, Norges Bank, Capital Group)",
        ], None),
        ("ESG", [
            "## Ambiental",
            "• 32% area preservada (acima exigido) | 0,42 tCO2e/ton (-30% vs media)",
            "• 100% rastreabilidade RTRS | Meta: -25% emissoes/ton ate 2030 (SBTi)",
            "",
            "## Social",
            "• 8.500 diretos + 12.000 safristas | 1.200 jovens formados",
            "",
            "## Governanca",
            "• Novo Mercado (100% ON, tag along) | 55,6% conselheiros independentes",
            "• Comites: Auditoria, Remuneracao, Sustentabilidade",
        ], None),
    ], "pitch_book_ecm_follow_on.pptx")

    # 3. Financial Model XLSX
    gen_xlsx("Modelo Financeiro DCF", [
        ("Premissas", [["Parametro", "Valor", "Fonte"],
            ["Rf (NTN-B 2035)", 0.115, "Tesouro Direto mar/25"],
            ["Beta (desalav.)", 0.75, "Damodaran Elec. Utilities EM"],
            ["ERP", 0.065, "Damodaran EM + Country Risk"],
            ["Ke", 0.1638, "CAPM"], ["Kd pre-tax", 0.135, "Media emissoes"],
            ["IR/CS", 0.34, "Legislacao"], ["Kd after-tax", 0.0891, ""],
            ["D/(D+E)", 0.40, "Target mercado"], ["WACC", 0.128, ""],
            ["g perpetuidade", 0.04, "PIB nominal LP"]]),
        ("Historico", [["R$ M", "2020", "2021", "2022", "2023", "2024"],
            ["Receita Liq.", 3820, 4850, 5920, 7180, 8200],
            ["EBITDA", 1580, 2180, 2650, 3320, 3800],
            ["Mg EBITDA", 0.414, 0.449, 0.448, 0.462, 0.463],
            ["Lucro Liq.", 383, 673, 858, 1241, 1518],
            ["Div. Liquida", 4400, 4400, 4700, 4600, 4500],
            ["DL/EBITDA", 2.78, 2.02, 1.77, 1.39, 1.18]]),
        ("Projecoes", [["R$ M", "2025E", "2026E", "2027E", "2028E", "2029E"],
            ["Receita", 9184, 10102, 11012, 11893, 12725],
            ["EBITDA", 4316, 4799, 5286, 5709, 6108],
            ["NOPAT", 2121, 2367, 2616, 2826, 3024],
            ["Capex", -2296, -2526, -2753, -2973, -3181],
            ["FCFF", 849, 981, 1111, 1210, 1303]]),
        ("DCF", [["Componente", "R$ M"],
            ["PV FCFs", 3758], ["PV Terminal Value", 8412],
            ["Enterprise Value", 16278], ["(-) Div. Liquida", -4500],
            ["Equity Value", 11918], ["Preco/acao", 20.55]]),
        ("Sensibilidade", [["WACC\\g", "3,0%", "3,5%", "4,0%", "4,5%", "5,0%"],
            ["11,0%", 26.8, 28.9, 31.5, 34.8, 39.2],
            ["12,0%", 22.0, 23.4, 25.1, 27.2, 29.8],
            ["12,8%", 19.2, 20.3, 21.6, 23.2, 25.1],
            ["14,0%", 15.6, 16.3, 17.2, 18.2, 19.5]]),
        ("Credito", [["Metrica", "2022", "2023", "2024", "Covenant"],
            ["DL/EBITDA", 1.77, 1.39, 1.18, "<=3,5x"],
            ["EBITDA/Desp.Fin", 4.08, 5.72, 7.31, ">=2,5x"],
            ["DSCR", 1.85, 2.10, 2.45, ">=1,3x"]]),
        ("Comps", [["Empresa", "EV/EBITDA", "P/E", "EV/Rev", "DL/EBITDA", "Mg EBITDA"],
            ["CPFL", 6.80, 11.20, 1.45, 2.10, 0.380],
            ["Engie", 7.20, 13.50, 2.80, 1.90, 0.420],
            ["Taesa", 8.50, 15.80, 5.20, 3.20, 0.850],
            ["Equatorial", 7.80, 14.20, 1.95, 2.80, 0.340],
            ["Mediana", 7.20, 13.50, 1.95, 2.50, 0.380],
            ["Eneva", 3.96, 6.95, 1.84, 1.18, 0.463]]),
    ], "modelo_financeiro_dcf_credit.xlsx")

    # 4. Bond Pricing XLSX
    gen_xlsx("Bond Pricing", [
        ("Pricing", [["", "1a Serie", "2a Serie"],
            ["Volume", 480000000, 320000000], ["Indexador", "CDI+", "IPCA+"],
            ["Spread", 0.0185, 0.0620], ["YTM", 0.1500, 0.1100],
            ["Prazo", "5 anos", "7 anos"], ["PU", 998.45, 1012.30],
            ["Duration", 4.60, 5.80], ["DV01", 4.00, 5.30]]),
        ("Comps", [["Emissor", "Rating", "Prazo", "Spread", "Vol (R$M)", "Coord."],
            ["CPFL", "AA+", "5y", "CDI+1,78%", 1200, "Itau BBA"],
            ["Taesa", "AAA", "7y", "CDI+1,45%", 800, "BTG"],
            ["Neoenergia", "AA", "5y", "CDI+1,92%", 1500, "Bradesco"],
            ["Equatorial", "AA+", "7y", "IPCA+6,35%", 2000, "XP"],
            ["Sabesp", "AA+", "5y", "CDI+1,70%", 1500, "BTG"]]),
        ("Sensibilidade", [["Cenario", "Spread", "YTM", "PU", "Var%"],
            ["-50bps", "1,35%", "14,50%", 1020.8, "+2,2%"],
            ["BASE", "1,85%", "15,00%", 998.5, "0,0%"],
            ["+50bps", "2,35%", "15,50%", 976.5, "-2,2%"]]),
    ], "modelo_bond_pricing.xlsx")

    # 5-9. PDFs
    gen_pdf("Memorando de Informacoes — Debentures Eneva S.A.", [
        ("Sumario da Oferta", ["Emissao de R$ 800M em debentures simples, 1a e 2a series, rating AA+ (Fitch).",
            "","## Termos Principais","• 1a Serie: R$ 480M, CDI+1,85%, 5 anos, bullet","• 2a Serie: R$ 320M, IPCA+6,20%, 7 anos, linear",
            "• Garantia: cessao fiduciaria de recebiveis PPA","• Destinacao: 60% refi + 40% expansao Parnaiba VI"]),
        ("Descricao da Emissora", ["Eneva S.A.: maior geradora termica privada a gas do Brasil (5,9 GW).",
            "","## Modelo de Negocio","• Operacao integrada gas+energia na Bacia do Parnaiba (MA)","• 92% receita contratada via PPAs (15+ anos)",
            "","## Financials 2024","• Receita: R$ 8,2Bi | EBITDA: R$ 3,8Bi (46,3%) | Lucro: R$ 1,6Bi","• DL/EBITDA: 1,2x | EBITDA/Desp.Fin: 4,8x"]),
        ("Estrutura e Garantias", ["## Garantia Principal","Cessao fiduciaria de recebiveis PPA — cobertura 2,7x servico da divida.",
            "","## Covenants","• DL/EBITDA <= 3,5x | EBITDA/Desp.Fin >= 2,5x | DSCR >= 1,3x",
            "• Cross-default (>R$150M) | Change of control | Restricted Payments"]),
        ("Fatores de Risco", ["## Risco 1: Concentracao (Prob: Media | Impacto: Alto)","• 70% capacidade no MA — mitigante: diversificacao em andamento",
            "","## Risco 2: Reservas Gas (Prob: Baixa | Impacto: Alto)","• Mitigante: certificacao DeGolyer, reserva 15+ anos",
            "","## Risco 3: Regulatorio (Prob: Media | Impacto: Medio)","• PPAs firmados nao sao afetados retroativamente"]),
    ], "memorando_informacoes_debentures.pdf")

    gen_pdf("Research Report — SLC Agricola S.A.", [
        ("Executive Summary", ["Recomendacao: COMPRA | Preco-alvo: R$ 19,50 (+13,4% vs R$ 17,20)",
            "","## Tese","• Lider em escala (680k ha) com custo 12% abaixo da media","• 70% receita USD — hedge cambial natural",
            "• FCF yield 8,5% | Div. yield 4,2% | 6,8x EV/EBITDA (-9% vs peers)"]),
        ("Valuation", ["## DCF (WACC 13,2%, g 4,0%)","• EV: R$ 18,2Bi | Equity: R$ 14,0Bi | R$ 21,30/acao",
            "","## Comps","• EV/EBITDA mediana peers: 7,5x | SLC: 6,8x (desconto 9%)",
            "","## Football Field","• Range: R$ 17,50 — R$ 21,80 | Media ponderada: R$ 19,50"]),
        ("Riscos", ["## Commodities (Impacto: Alto)","• Hedge 60% producao em futuros (CBOT/B3)",
            "","## Red Flag: Contratos Arrendamento","• 3 contratos com rescisao unilateral — 22% area | Impacto: R$ 180M EBITDA",
            "• Probabilidade: baixa (2/3 renovados ultimos 10 anos)"]),
    ], "research_report_completo.pdf")

    gen_pdf("Term Sheet — CRA Movida Logistica", [
        ("Dados da Emissao", ["• Emissor: Eco Securitizadora (Movida Part.) | CRA Serie Unica | R$ 400M",
            "• Lastro: recebiveis locacao veiculos (frotas agro) | Rating: AA (S&P)","• CDI+2,10% semestral | 5 anos | Linear trimestral (12o mes+)",
            "• Registro RCVM 160 automatico | B3 CETIP21"]),
        ("Lastro e Garantias", ["## Composicao","• 4.200 contratos | Lastro R$ 580M (1,45x) | Inadimplencia: 1,8%",
            "","## Protecao","• Sobrecolateral 145% | Fundo reserva 3 meses | Trigger amort. se inadimpl. >4%"]),
        ("Covenants", ["• DL/EBITDA Movida <= 3,5x | Cobertura lastro >= 1,30x | Inadimpl. <= 4,0%",
            "","## Vencimento Antecipado","• Inadimplencia >3DU | Cobertura <1,15x (2 periodos) | Downgrade <A- | RJ/Falencia"]),
    ], "term_sheet_cra.pdf")

    gen_pdf("Executive Memo — Comite de Credito", [
        ("Proposta", ["Participacao como coordenador lider — Debentures Eneva R$ 800M.",
            "• Fee: 0,65% flat (R$ 5,2M) + 0,15% manutencao | Demanda: R$ 1,4Bi (1,75x)"]),
        ("Credito", ["## Metricas dez/2024","• DL/EBITDA: 1,18x (cov: 3,5x) | EBITDA/Fin: 4,8x (cov: 2,5x) | DSCR: 2,45x",
            "• Rating implicito: AA+ | Capacidade adicional: ~R$ 8,8Bi"]),
        ("Parecer", ["## RECOMENDACAO: GO","• Investment grade com track record solido","• Fluxo previsivel (92% PPA) + folga ampla em covenants",
            "","## Condicao: Parecer juridico garantias ate 25/03"]),
    ], "executive_memo_comite.pdf")

    gen_pdf("Investor Presentation — SLC Agricola Roadshow", [
        ("Investment Highlights", ["## Por que SLC?","• #1 escala (680k ha) | Custo -12% | 70% receita USD | FCF yield 8,5%",
            "","## Uso de Recursos (R$ 720M primario)","• 50% expansao area | 30% tecnologia | 20% deleverage"]),
        ("Financials", ["## Historico","• Receita: R$ 5,2Bi (2021) → R$ 7,8Bi (2024) — CAGR 14,5%",
            "• EBITDA: R$ 1,8Bi → R$ 2,9Bi — margem 37,4%","• FCF: R$ 1,46Bi (yield 8,5%)"]),
        ("ESG", ["• 32% area preservada | 0,42 tCO2e/ton (-30% media) | RTRS certificado",
            "• Novo Mercado B3 | 55,6% conselheiros independentes"]),
    ], "investor_presentation_roadshow.pdf")

    print(f"\nTotal: {len(list(OUT.glob('*')))} arquivos em {OUT.resolve()}")
